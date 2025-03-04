# cython: language_level=3
# cython: boundscheck=False
# cython: wraparound=True
# cython: cdivision=True
# cython: nonecheck=False
#! /usr/bin/env python3
# -*- coding: utf-8 -*-
# version: 2.1.0

import os
import shutil
import subprocess
import zlib
import hashlib
from typing import List, Union, Tuple, Optional, Callable, Any, Dict
import xml.etree.ElementTree as ET
import tempfile
import pkg_resources
import io
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
import sys
import time
from datetime import datetime
import uuid
import zipfile
from PIL import Image
import random
import urllib.parse
import re
import aiohttp
import asyncio
import configparser
from pathlib import Path
from fpdf import FPDF
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.platypus.frames import Frame
from reportlab.platypus.doctemplate import PageTemplate
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.utils import ImageReader
import cairosvg
from reportlab.platypus.tableofcontents import TableOfContents

def load_config() -> dict:
    """
    ~/.airun/airun.conf 파일에서 설정을 읽어옵니다.
    
    Returns:
        dict: 설정값들을 담은 딕셔너리
    """
    config = {}
    config_path = os.path.expanduser("~/.airun/airun.conf")
    
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if line and line.startswith('export '):
                    # 'export KEY="VALUE"' 형식 파싱
                    line = line.replace('export ', '')
                    key, value = line.split('=', 1)
                    key = key.strip()
                    value = value.strip().strip('"').strip("'")
                    config[key] = value
    except Exception as e:
        print(f"[WARNING] 설정 파일 로드 실패: {str(e)}")
        
    return config

# 설정 파일에서 SMTP 설정 로드
config = load_config()
SMTP_HOST = config.get("SMTP_HOST", "smtp.worksmobile.com")  # 기본값 설정
try:
    SMTP_PORT = int(config.get("SMTP_PORT", "587"))  # 기본값 설정
except (ValueError, TypeError):
    SMTP_PORT = 587  # 변환 실패시 기본값 사용
SMTP_USERNAME = config.get("SMTP_USERNAME", "")
SMTP_PASSWORD = config.get("SMTP_PASSWORD", "")

def install_if_missing(package: str, import_name: str = None) -> None:
    """Install a package if it's not already installed
    
    Args:
        package (str): Package name to install with pip
        import_name (str, optional): Module name to import. Defaults to package name.
    """
    import subprocess
    import sys
    import pkg_resources
    
    try:
        # 패키지가 이미 설치되어 있는지 확인
        pkg_resources.require(package)
        
        # pyhwp의 경우 특별 처리
        if package == "pyhwp":
            try:
                __import__("hwp5")
                return
            except ImportError:
                raise pkg_resources.DistributionNotFound(package)
        
        # 다른 패키지들의 경우 기존 방식대로 처리
        if import_name:
            __import__(import_name)
        else:
            __import__(package.replace('-', '_'))
        return  # 이미 설치되어 있고 import 가능하면 종료
    except (pkg_resources.DistributionNotFound, ImportError):
        try:
            # 가상환경의 pip 사용
            pip_path = os.path.join(os.path.dirname(sys.executable), 'pip')
            
            # 패키지 설치
            subprocess.run([
                pip_path, "install", "--quiet", package
            ], check=True)
            
            # print(f"[INFO] Successfully installed {package}")
            
            # 설치 후 import 확인
            if package == "pyhwp":
                try:
                    __import__("hwp5")
                except ImportError as e:
                    print(f"[WARNING] Package installed but import failed: {str(e)}")
                    raise
        except subprocess.CalledProcessError as e:
            print(f"[ERROR] Failed to install {package}: {str(e)}")
            raise

# Required packages
# 패키지 이름과 import 이름이 다른경우 처리
REQUIRED_PACKAGES = [
    ("requests", "requests"),
    ("pandas", "pandas"),
    ("numpy", "numpy"),
    ("matplotlib", "matplotlib"),
    ("openpyxl", "openpyxl"),
    ("fpdf", "fpdf"),
    ("pyhwp", "hwp5"),  # pyhwp의 실제 import 이름은 hwp5입니다
    ("olefile", "olefile"),  
    ("python-pptx", "pptx"),
    ("PyPDF2", "PyPDF2"),
    ("pdfminer.six", "pdfminer"),  # PDF 텍스트 추출을 위한 패키지
    ("pycryptodome", "Crypto"),    # PDF 암호화 처리를 위한 패키지
    ("Pillow", "PIL"),
    ("svglib", "svglib"),
    ("reportlab", "reportlab"),
    ("selenium", "selenium"),
    ("webdriver_manager", "webdriver_manager"),
    ("beautifulsoup4", "bs4"),
    ("lxml[html_clean]", "lxml"),  # html_clean feature 포함
    ("python-docx", "docx"),
    ("trafilatura", "trafilatura"),
    ("cairosvg", "cairosvg"),
    ("tabulate", "tabulate")  # 테이블 처리를 위한 패키지 추가
]

# 패키지 설치 상태 확인 및 설치
# print("\n[INFO] Checking required packages...")
for package, import_name in REQUIRED_PACKAGES:
    try:
        # importlib를 사용하여 모듈 import 시도
        if import_name:
            __import__(import_name)
        else:
            __import__(package.replace('-', '_'))
    except ImportError:
        print(f"[INFO] Installing missing package: {package}")
        install_if_missing(package, import_name)
        
# print("[INFO] All required packages are ready.")

# Then import all required modules
import smtplib
import webbrowser
import io
import zipfile
import olefile
import requests
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from fpdf import FPDF
from PIL import Image
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
from selenium.common.exceptions import TimeoutException, NoSuchElementException
from webdriver_manager.chrome import ChromeDriverManager
from urllib.parse import urlparse, urlunparse
from bs4 import BeautifulSoup
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication
from email.utils import formatdate
import cairosvg  # SVG 처리를 위한 라이브러리
# File extensions supported by pandas
PANDAS_EXTENSIONS = {
    '.xlsx': pd.read_excel,  # Excel files
    '.xls': pd.read_excel,
    '.csv': pd.read_csv,     # CSV files
    '.json': pd.read_json,   # JSON files
    '.html': pd.read_html,   # HTML files
    '.xml': pd.read_xml,     # XML files
    '.parquet': pd.read_parquet,  # Parquet files
    '.feather': pd.read_feather,  # Feather files
    '.pickle': pd.read_pickle,    # Pickle files
    '.sql': pd.read_sql,     # SQL files
    '.hdf': pd.read_hdf,     # HDF5 files
    '.sas': pd.read_sas,     # SAS files
    '.stata': pd.read_stata,  # Stata files
    '.spss': pd.read_spss    # SPSS files
}

# Pandas writers for different file types
PANDAS_WRITERS = {
    '.xlsx': lambda df, path: df.to_excel(path, index=False),
    '.xls': lambda df, path: df.to_excel(path, index=False),
    '.csv': lambda df, path: df.to_csv(path, index=False),
    '.json': lambda df, path: df.to_json(path),
    '.html': lambda df, path: df.to_html(path),
    '.xml': lambda df, path: df.to_xml(path),
    '.parquet': lambda df, path: df.to_parquet(path),
    '.feather': lambda df, path: df.to_feather(path),
    '.pickle': lambda df, path: df.to_pickle(path),
    '.sql': lambda df, path: df.to_sql(path),
    '.hdf': lambda df, path: df.to_hdf(path, 'data'),
}

# HWP file handlers
HWP_EXTENSIONS = {
    '.hwp': 'hwp',     # 한글 문서
    '.hwpx': 'hwpx'    # 한글 2018 이상 문서
}

# ChromeDriver 경로를 저장할 전역 변수
_chrome_driver_path = None

# ============================================================================
# 문서 처리 클래스 (Document Processing Classes)
# ============================================================================

def convert_image_to_rgb(img):
    """Convert image to RGB mode safely.
    
    Args:
        img: PIL Image object
    
    Returns:
        PIL Image object in RGB mode
    """
    if img.mode in ('RGBA', 'LA'):
        # RGBA나 LA 모드인 경우 알파 채널을 고려하여 변환
        background = Image.new('RGB', img.size, (255, 255, 255))
        if 'A' in img.mode:  # 알파 채널이 있는 경우
            background.paste(img, mask=img.split()[-1])
        else:
            background.paste(img)
        return background
    elif img.mode == 'P':  # 팔레트 모드
        return img.convert('RGB')
    elif img.mode != 'RGB':  # 그 외 모드
        return img.convert('RGB')
    return img

class HWPDocument:
    
    # 본문 스타일 정의 추가
    paragraph_styles = {
        # 0-2: 기본 크기 (왼쪽, 가운데, 오른쪽)
        'normal': {'char_pr_id': '0', 'align': 'left', 'font_size': 10, 'line_spacing': '120'},
        'center': {'char_pr_id': '0', 'align': 'center', 'line_spacing': '120'},
        'right': {'char_pr_id': '0', 'align': 'right', 'line_spacing': '120'},
        
        # 3-5: 중간 크기 (왼쪽, 가운데, 오른쪽)
        'medium': {'char_pr_id': '10', 'align': 'left', 'font_size': 12, 'line_spacing': '130', 'para_pr_id': '10'},
        'medium_center': {'char_pr_id': '10', 'align': 'center', 'font_size': 12, 'line_spacing': '130', 'para_pr_id': '11'},
        'medium_right': {'char_pr_id': '10', 'align': 'right', 'font_size': 12, 'line_spacing': '130', 'para_pr_id': '12'},
        
        # 6-8: 큰 크기 (왼쪽, 가운데, 오른쪽)
        'large': {'char_pr_id': '9', 'align': 'left', 'font_size': 14, 'bold': True, 'line_spacing': '150'},
        'large_center': {'char_pr_id': '9', 'align': 'center', 'font_size': 14, 'bold': True, 'line_spacing': '150'},
        'large_right': {'char_pr_id': '9', 'align': 'right', 'font_size': 14, 'bold': True, 'line_spacing': '150'},
        
        # 9-11: 중간 크기 굵은체 (왼쪽, 가운데, 오른쪽)
        'medium_bold': {'char_pr_id': '12', 'align': 'left', 'font_size': 12, 'bold': True, 'line_spacing': '130', 'para_pr_id': '13'},
        'medium_bold_center': {'char_pr_id': '12', 'align': 'center', 'font_size': 12, 'bold': True, 'line_spacing': '130', 'para_pr_id': '14'},
        'medium_bold_right': {'char_pr_id': '12', 'align': 'right', 'font_size': 12, 'bold': True, 'line_spacing': '130', 'para_pr_id': '15'},
        
        # 12-14: 큰 글씨 보통체 (왼쪽, 가운데, 오른쪽)
        'large_normal': {'char_pr_id': '11', 'align': 'left', 'font_size': 14, 'line_spacing': '150', 'para_pr_id': '16'},
        'large_normal_center': {'char_pr_id': '11', 'align': 'center', 'font_size': 14, 'line_spacing': '150', 'para_pr_id': '17'},
        'large_normal_right': {'char_pr_id': '11', 'align': 'right', 'font_size': 14, 'line_spacing': '150', 'para_pr_id': '18'},
        
        # 15-16: 기타 스타일
        'emphasis': {'char_pr_id': '0', 'align': 'left', 'font_size': 10, 'bold': True, 'line_spacing': '120'},
        'quote': {'char_pr_id': '0', 'align': 'left', 'indent': 20, 'line_spacing': '120'},
    }
        
    def __init__(self):
        self.elements = []  # (type, content, page_break, options) 튜플 리스트
        self.has_title = False  # 제목 존재 여부 추적
        # ~/.airun/templates 디렉토리에서 템플릿 파일 찾기
        self.template_path = os.path.expanduser('~/.airun/templates/blank.hwpx')
        if not os.path.exists(self.template_path):
            # 템플릿 디렉토리가 없으면 생성
            template_dir = os.path.dirname(self.template_path)
            os.makedirs(template_dir, exist_ok=True)
            raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {self.template_path}")
        self._temp_files = []  # Track temporary files for cleanup
        
    @staticmethod
    def _convert_image_to_rgb(img):
        """Convert image to RGB mode safely.
        
        Args:
            img: PIL Image object
        
        Returns:
            PIL Image object in RGB mode
        """
        if img.mode in ('RGBA', 'LA'):
            # RGBA나 LA 모드인 경우 알파 채널을 고려하여 변환
            background = Image.new('RGB', img.size, (255, 255, 255))
            if 'A' in img.mode:  # 알파 채널이 있는 경우
                background.paste(img, mask=img.split()[-1])
            else:
                background.paste(img)
            return background
        elif img.mode == 'P':  # 팔레트 모드
            return img.convert('RGB')
        elif img.mode != 'RGB':  # 그 외 모드
            return img.convert('RGB')
        return img        
        
    @staticmethod
    def _read_url(url: str) -> Union[str, bytes]:
        """
        Read and return the contents of a URL.
        URL의 내용을 읽어 반환합니다.
        
        Args:
            url (str): URL to read from
                읽을 URL
            
        Returns:
            Union[str, bytes]: Contents of the URL. Returns bytes for binary content (images, etc)
                            URL의 내용. 바이너리 콘텐츠(이미지 등)의 경우 bytes 반환
            
        Raises:
            requests.RequestException: If the URL request fails
                                    URL 요청 실패 시
        """
        response = requests.get(url, verify=False)
        response.raise_for_status()
        
        # Check content type
        content_type = response.headers.get('content-type', '').lower()
        if any(t in content_type for t in ['image/', 'video/', 'audio/', 'application/octet-stream']):
            return response.content
        return response.text
        
    def _preprocess_text(self, text):
        """특수문자와 글머리 기호를 처리하는 내부 메서드"""
        if not text:
            return text

        # ** 문자 제거 (가장 먼저 처리)
        text = text.replace('**', '')

        # 마크다운 형식이나 섹션 제목은 그대로 반환
        if re.match(r'^#{1,6}\s+', text.strip()) or \
           re.match(r'^\d+\.\s+', text.strip()) or \
           'KPI' in text:  # KPI가 포함된 경우 그대로 반환
            return text
                
        # 보존할 패턴 정의
        preserved_patterns = [
            # 숫자 관련
            (r'\d+(?:\.\d+)?%', lambda m: m.group()),  # 20%, 50.5% 등
            (r'\d+(?:\.\d+)?(?:천|만|억|조)?원', lambda m: m.group()),  # 5억 원, 10만 원 등
            (r'\d+(?:\.\d+)?(?:차)?년도', lambda m: m.group()),  # 1차년도, 2년도 등
            (r'\d+\s*(?:개|건|회|명)', lambda m: m.group()),  # 10개, 20건 등
            
            # 특수 형식
            (r'[A-Z]+(?:/[A-Z]+)*', lambda m: m.group()),  # KPI, ESG, R&D 등
            (r'\([^)]*\)', lambda m: m.group()),  # 괄호 안 내용
            
            # 날짜/시간
            (r'\d{4}[-/\.]\d{1,2}[-/\.]\d{1,2}', lambda m: m.group()),
        ]
        
        # 임시 토큰으로 보존할 패턴 치환
        preserved_tokens = {}
        for pattern, replacement in preserved_patterns:
            text = re.sub(pattern, lambda m: preserved_tokens.setdefault(f'__TOKEN_{len(preserved_tokens)}__', 
                         replacement(m) if callable(replacement) else replacement), text)
        
        # 특수문자 처리 규칙
        special_chars_map = {
            # 수학 기호
            '×': '×',  # 곱하기 기호 유지
            '÷': '÷',  # 나누기 기호 유지
            '±': '±',  # 플러스마이너스 유지
            '∓': '∓',  # 마이너스플러스 유지
            '∔': '+',  # 플러스로 변환
            '∸': '-',  # 마이너스로 변환
            '∹': '/',  # 나누기로 변환
            '⋅': '·',  # 가운뎃점 유지
            
            # 일반 특수문자
            '&': '&',      # & 유지
            '%': '%',      # % 유지
            '=': '=',      # = 유지
            '/': '/',      # / 유지
            
            # 공백 문자
            '\u3000': ' ',  # 전각 공백
            '\u200b': '',   # 제로 너비 공백
            '\ufeff': '',   # BOM
        }
        
        # 특수문자 치환
        for char, replacement in special_chars_map.items():
            text = text.replace(char, replacement)
        
        # 보존된 패턴 복원
        for token, original in preserved_tokens.items():
            text = text.replace(token, original)
        
        # XML 특수문자 이스케이프
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        text = text.replace('"', '&quot;')
        text = text.replace("'", '&apos;')
        
        # 연속된 공백을 하나로
        text = ' '.join(text.split())
        
        return text

    def _split_long_text(self, text, max_length=60):
        """긴 텍스트를 적절한 길이로 분리합니다."""
        sentences = []
        # 먼저 줄바꿈으로 분리
        for paragraph in text.split('\n'):
            paragraph = paragraph.strip()
            if not paragraph:
                sentences.append('')
                continue
                
            # 문장이 max_length보다 길면 추가로 분리
            while len(paragraph) > max_length:
                # 공백을 기준으로 단어 분리
                split_idx = paragraph[:max_length].rfind(' ')
                if split_idx == -1:  # 공백을 찾지 못한 경우
                    split_idx = max_length
                sentences.append(paragraph[:split_idx].strip())
                paragraph = paragraph[split_idx:].strip()
            
            if paragraph:  # 남은 문장 추가
                sentences.append(paragraph)
        
        return sentences

    def _join_broken_lines(self, lines):
        """잘린 문장을 하나로 합칩니다."""
        result = []
        current = []
        
        def should_start_new_line(line):
            """새로운 줄을 시작해야 하는지 확인"""
            stripped = line.strip()
            # 빈 줄
            if not stripped:
                return True
            # 글머리 기호로 시작하는 줄
            if stripped.startswith(('-', '•', '*', '+')):
                return True
            # 숫자 목록 (1., 1.1. 등)
            if re.match(r'^\d+\.(\d+\.)?\s', stripped):
                return True
            # 마크다운 헤더
            if stripped.startswith('#'):
                return True
            return False

        def should_end_line(line):
            """현재 줄을 종료해야 하는지 확인"""
            stripped = line.strip()
            # URL이나 특수 형식
            if any(marker in stripped for marker in ['http://', 'https://', '[Page']):
                return True
            # 문장 종결 부호로 끝나는 경우
            return stripped.endswith(('다.', '까?', '요.', '임.', '됨.', '함.', '.', '?', '!'))

        lines = [line.rstrip() for line in lines]  # 오른쪽 공백만 제거
        i = 0
        while i < len(lines):
            line = lines[i]
            
            # 빈 줄 처리
            if not line.strip():
                if current:
                    result.append(' '.join(current))
                    current = []
                result.append('')
                i += 1
                continue
            
            # 현재 줄이 하이픈으로 끝나고 다음 줄이 있는 경우
            if line.endswith('-') and i + 1 < len(lines):
                if current:
                    current.append(line[:-1])  # 하이픈 제거
                    current.append(lines[i + 1].strip())
                else:
                    current = [line[:-1], lines[i + 1].strip()]
                i += 2
                continue
            
            # 새로운 줄을 시작해야 하는 경우
            if should_start_new_line(line.strip()):
                if current:
                    result.append(' '.join(current))
                    current = []
                result.append(line.strip())
                i += 1
                continue
            
            # 현재 줄을 종료해야 하는 경우
            if should_end_line(line.strip()):
                if current:
                    current.append(line.strip())
                    result.append(' '.join(current))
                    current = []
                else:
                    result.append(line.strip())
                i += 1
                continue
            
            # 일반적인 문장 연결
            if current:
                current.append(line.strip())
            else:
                current = [line.strip()]
            i += 1
        
        # 남은 문장 처리
        if current:
            result.append(' '.join(current))
        
        # 연속된 빈 줄 정리 (최대 2개까지만 허용)
        final_result = []
        empty_count = 0
        for line in result:
            if not line:
                empty_count += 1
                if empty_count <= 2:
                    final_result.append(line)
            else:
                empty_count = 0
                final_result.append(line)
        
        return final_result

    def _normalize_content(self, content):
        """파일 내용을 정규화하면서 원본 포맷을 최대한 유지합니다."""
        if isinstance(content, (list, tuple)):
            content = '\n'.join(str(item) for item in content)
        
        if not isinstance(content, str):
            content = str(content)

        # ** 문자 제거
        content = re.sub(r'\*\*', '', content)
        
        lines = content.split('\n')
        lines = self._join_broken_lines(lines)
        
        result = []
        prev_empty = True
        
        for line in lines:
            if line.strip() in ['<그림>', '<표>']:
                continue
                
            if not line.strip():
                result.append('')
                prev_empty = True
                continue
            
            stripped = line.strip()
            
            # 마크다운 헤더 확인 (#, ##, ###)
            header_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
            if header_match:
                level = len(header_match.group(1))
                text = header_match.group(2).strip()
                # 레벨에 따른 스타일 매핑 (# = 3, ## = 4, ### = 5, ...)
                style_id = str(2 + level)  # level이 1(#)이면 3, 2(##)이면 4, 3(###)이면 5
                result.append(('heading', text, False, {'size': style_id}))
                prev_empty = False
                continue
            
            # 숫자로 시작하는 제목 패턴 (### 1. 형식)
            numbered_title_match = re.match(r'^(#{1,6})\s+(\d+\.)\s+(.+)$', stripped)
            if numbered_title_match:
                level = len(numbered_title_match.group(1))
                number = numbered_title_match.group(2)
                text = numbered_title_match.group(3).strip()
                style_id = str(2 + level)  # 레벨에 따른 스타일 매핑 사용
                result.append(('heading', f"{number} {text}", False, {'size': style_id}))
                prev_empty = False
                continue
            
            # 일반 텍스트
            result.append(line)
            prev_empty = False
        
        return result

    def add_heading(self, text: str, level: int = 1, page_break: bool = False, options: dict = None) -> None:
        """문서에 제목을 추가합니다.
        Args:
            text (str): 제목 텍스트
            level (int): 제목 레벨 (1: 문서 제목, 2: 제목 1, 3: 제목 2, 4: 제목 3, 5: 제목 4)
            page_break (bool): 페이지 나누기 여부
            options (dict): 추가 옵션
        """
        if not options:
            options = {}
    
        # 레벨에 따른 스타일 매핑
        level_styles = {
            1: {'style': '2', 'charPrIDRef': '2'},  # 문서 제목 (16pt, 굵게, 가운데)
            2: {'style': '3', 'charPrIDRef': '3'},  # 제목 1 (14pt, 굵게, 파란색)
            3: {'style': '4', 'charPrIDRef': '4'},  # 제목 2 (12pt, 굵게, 파란색)
            4: {'style': '5', 'charPrIDRef': '5'},  # 제목 3 (11pt, 굵게, 파란색)
            5: {'style': '6', 'charPrIDRef': '6'},  # 제목 4 (10pt, 굵게, 파란색)
        }
    
        # 기본 스타일에 사용자 옵션 병합
        style = level_styles.get(level, {'style': '2', 'charPrIDRef': '2'})
        style.update(options)
    
        text = self._preprocess_text(text)
        self.elements.append(('heading', text, page_break, style))
        self.has_title = True

    def add_paragraph(self, text: str, page_break: bool = False, options: Dict = None, style: Union[str, int] = None) -> None:
        """문단 추가
        
        Args:
            text: 추가할 텍스트 (str, list, DataFrame 등 다양한 타입 지원)
            page_break (bool): 페이지 나누기 여부
            options (Dict, optional): 문단 옵션
                - font_name (str): 폰트 이름
                - font_size (int): 폰트 크기
                - bold (bool): 굵게
                - align (str): 정렬 ('left', 'center', 'right')
                - indent (int): 들여쓰기
                - spacing_before (int): 문단 앞 간격
                - spacing_after (int): 문단 뒤 간격
                - tab_stops (List[Dict]): 탭 설정
                    - position (int): 탭 위치
                    - alignment (str): 탭 정렬
                    - leader (str): 리더 문자 (예: 'dot')
                - field_code (bool): 필드 코드 여부
                - bookmark (bool): 북마크 여부
            style (str or int, optional): 미리 정의된 스타일 이름 또는 인덱스 번호
                                  style이 제공되면 options는 추가 옵션으로 처리됩니다.
                                  
                                  인덱스 번호 설명:
                                  - 0-2: 기본 크기 (왼쪽, 가운데, 오른쪽)
                                  - 3-5: 중간 크기 (왼쪽, 가운데, 오른쪽)
                                  - 6-8: 큰 크기 굵은체 (왼쪽, 가운데, 오른쪽)
                                  - 9-11: 중간 크기 굵은체 (왼쪽, 가운데, 오른쪽)
                                  - 12-14: 큰 크기 보통체 (왼쪽, 가운데, 오른쪽)
                                  - 15: 강조 (굵게)
                                  - 16: 인용문 (들여쓰기)
        """
        # style이 제공된 경우 add_styled_paragraph 메서드 호출
        if style is not None:
            # style이 정수인 경우 스타일 목록에서 해당 인덱스의 스타일 이름을 가져옴
            if isinstance(style, int):
                style_keys = list(self.paragraph_styles.keys())
                if 0 <= style < len(style_keys):
                    style = style_keys[style]
                else:
                    raise ValueError(f"유효하지 않은 스타일 인덱스입니다: {style}. 유효한 범위: 0-{len(style_keys)-1}")
            
            return self.add_styled_paragraph(text, style=style, page_break=page_break, additional_options=options)
        # style이 제공되지 않은 경우 기본 문단 추가 로직 실행
        else:
            try:
                if not text:  # 빈 텍스트 처리
                    return
                    
                if options is None:
                    options = {}
                    
                # 필드 코드 처리
                if options.get('field_code'):
                    if text.startswith('{PAGE}'):
                        # 페이지 번호 필드
                        self.elements.append(('field', 'page_number', None, options))
                    elif text.startswith('{REF'):
                        # 페이지 참조 필드
                        self.elements.append(('field', 'page_ref', text[5:-4], options))
                    return
                    
                # 북마크 처리
                if options.get('bookmark'):
                    if text.startswith('{BM='):
                        # 북마크 추가
                        bookmark_name = text[4:-1]
                        self.elements.append(('bookmark', bookmark_name, None, options))
                    return
                    
                # DataFrame 처리
                if hasattr(text, 'to_string'):  # pandas DataFrame인 경우
                    # DataFrame을 표로 변환
                    header = text.columns.tolist()
                    data = text.values.tolist()
                    self.add_table(data, header=header)
                    return
                    
                # 리스트나 튜플을 문자열로 변환
                if isinstance(text, (list, tuple)):
                    text = '\n'.join(str(item) for item in text)
                # 그 외 타입은 str로 변환
                elif not isinstance(text, str):
                    text = str(text)
                    
                # 줄 단위로 처리
                lines = text.split('\n')
                for line in lines:
                    stripped = line.strip()
                    stripped = self._preprocess_text(stripped)
                    if not stripped:
                        continue
                        
                    # 마크다운 헤더 확인 (##, ###)
                    header_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
                    if header_match:
                        level = len(header_match.group(1))
                        title_text = header_match.group(2).strip()
                        style_id = str(2 + level)
                        self.add_heading(title_text, page_break=page_break, options={'size': style_id})
                        continue
                    
                    # 숫자로 시작하는 제목 패턴 (### 1. 형식)
                    numbered_title_match = re.match(r'^###\s+(\d+\.)\s+(.+)$', stripped)
                    if numbered_title_match:
                        number = numbered_title_match.group(1)
                        title_text = numbered_title_match.group(2).strip()
                        self.add_heading(f"{number} {title_text}", page_break=page_break, options={'size': '5'})
                        continue
                    
                    # 일반 텍스트는 add_text_content로 처리
                    self.add_text_content(line, options)
                    
            except Exception as e:
                print(f"문단 추가 실패: {str(e)}")
                raise

    def add_text_content(self, text, options=None):
        """일반 텍스트 내용을 문서에 추가합니다."""
        if not text:
            return
            
        if options is None:
            options = {}
            
        # ** 문자 제거
        text = text.replace('**', '')
        text = self._preprocess_text(text)

        # 글머리 기호 패턴 확인 함수
        def has_bullet_point(text: str) -> bool:
            # 숫자+괄호 패턴 (1), 2), 등)
            if re.match(r'^\s*\d+\)\s+', text):
                return True
            # 원문자 숫자 패턴 (①, ②, 등)
            if re.match(r'^\s*[①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮]\s+', text):
                return True
            # 알파벳+괄호 패턴 (a), b), 등)
            if re.match(r'^\s*[a-zA-Z]\)\s+', text):
                return True
            # 사각형, 다이아몬드 기호 패턴 (□, ■, ◆, ◇ 등)
            if re.match(r'^\s*[□■◆◇]\s+', text):
                return True
            return False
        
        # 특수 글머리 기호 패턴 확인 함수 (※, ✱, ❖, -, –, — 등)
        def has_special_bullet_point(text: str) -> bool:
            # ※, ✱, ❖ 기호 패턴
            if re.match(r'^\s*[※✱❖❍]\s+', text):
                return True
            # 대시(-) 패턴 (-, - , –, — 등)
            if re.match(r'^\s*[-–—]\s+', text):
                return True            
            return False
        
        # 글머리 기호가 있는 경우 공백 2개 추가
        if has_bullet_point(text):
            text = "  " + text
        
        # 특수 글머리 기호가 있는 경우 공백 4개 추가
        if has_special_bullet_point(text):
            text = "    " + text
            
        # 문단 옵션 설정
        paragraph_options = {
            'char_pr_id': options.get('char_pr_id', "0"), 
            'line_spacing': '120'
        }
        
        # 정렬 옵션 처리
        if 'align' in options:
            paragraph_options['align'] = options['align']
            
        # 폰트 옵션 처리
        if 'font_name' in options:
            paragraph_options['font_name'] = options['font_name']
        if 'font_size' in options:
            paragraph_options['font_size'] = options['font_size']
        if 'bold' in options:
            paragraph_options['bold'] = options['bold']
            
        # 들여쓰기 옵션 처리
        if 'indent' in options:
            paragraph_options['indent'] = options['indent']
            
        # 간격 옵션 처리
        if 'spacing_before' in options:
            paragraph_options['spacing_before'] = options['spacing_before']
        if 'spacing_after' in options:
            paragraph_options['spacing_after'] = options['spacing_after']
            
        # 긴 문단 처리
        if len(text) > 60:
            for sentence in self._split_long_text(text):
                if sentence.strip():
                    self.elements.append(('paragraph', sentence.strip(), False, paragraph_options))
        else:
            self.elements.append(('paragraph', text, False, paragraph_options))

    def add_tooltip(self, text):
        """툴팁 추가
        
        Args:
            text: 툴팁 텍스트
        """
        if not text:
            return
           
        data = [
            ["AI.RUN 2025. - Empowering Your AI Journey"],
            [text]
        ]
        self.add_table(data, style=4, header_style=3, text_align='left')

    def add_page_break(self):
        """
        빈 문단과 함께 페이지 넘김을 추가합니다.
        """
        self.elements.append(('paragraph', '', True, {'char_pr_id': "0"}))

    def add_image(self, image):
        """이미지 추가
        
        Args:
            image: 이미지 파일 경로, 바이너리 데이터, 또는 PIL Image 객체
        """
        try:
            if isinstance(image, bytes):
                # 바이너리 데이터를 PIL Image로 변환
                img_obj = Image.open(io.BytesIO(image))
            elif isinstance(image, str):
                # 파일 경로
                if not os.path.exists(image):
                    raise print(f"이미지 파일을 찾을 수 없습니다: {image}")
                img_obj = Image.open(image)
            elif isinstance(image, Image.Image):
                img_obj = image
            else:
                raise print("지원되지 않는 이미지 형식입니다")
            
            # 이미지 처리
            img_obj = self._convert_image_to_rgb(img_obj)
            img_obj = self._resize_image_to_page_width(img_obj)
            
            # 원본 이미지 형식 확인 및 유지
            original_format = img_obj.format
            if not original_format or original_format == 'JPEG':
                # 원본 형식이 없거나 JPEG인 경우 PNG로 저장 (무손실)
                img_format = 'PNG'
                file_ext = '.png'
            else:
                # 원본 형식 유지 (PNG, GIF, BMP 등)
                img_format = original_format
                file_ext = f'.{original_format.lower()}'
            
            # 임시 파일로 저장
            with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp_file:
                # 고품질 설정으로 저장
                if img_format == 'JPEG':
                    img_obj.save(temp_file.name, img_format, quality=95)  # 높은 품질 설정
                else:
                    img_obj.save(temp_file.name, img_format)
                    
                self._temp_files.append(temp_file.name)
                self.elements.append(('image', temp_file.name, False, {}))
            
        except Exception as e:
            raise print(f"이미지 처리 중 오류 발생: {str(e)}")

    def add_table(self, data, header=None, options: Dict = None, style: int = 2, align: str = 'center', text_align: str = 'center', header_style: int = None) -> None:
        """표를 추가합니다.
        
        Args:
            data: 표 데이터 (2차원 리스트)
            header: 헤더 데이터 (선택사항)
            options: 표 옵션 (선택사항)
            style: 테두리 스타일 (1~8, 기본값: 2)
                1: 회색 배경
                2: 회색 테두리
                3: 기본 테두리
                4: 굵은 테두리
                5: 이중 테두리
                6: 점선 테두리
                7: 굵은 실선
                8: 이중 실선
            align: 표 정렬 ('left', 'center', 'right', 기본값: 'center')
            text_align: 셀 내용 정렬 ('left', 'center', 'right', 기본값: 'center')
            header_style: 헤더 행 스타일 (1~8, 기본값: None - 지정하지 않으면 모든 행에 style 적용)
        """
        try:
            if not data:  # 빈 표 처리
                return
                
            if options is None:
                options = {}
                
            # align 값 검증
            align = align.lower()
            if align not in ['left', 'center', 'right']:
                raise ValueError("align은 'left', 'center', 'right' 중 하나여야 합니다.")
                
            # text_align 값 검증
            text_align = text_align.lower()
            if text_align not in ['left', 'center', 'right']:
                raise ValueError("text_align은 'left', 'center', 'right' 중 하나여야 합니다.")
                
            # HWPX 정렬 값으로 변환
            align_map = {
                'left': 'LEFT',
                'center': 'CENTER',
                'right': 'RIGHT'
            }
            hwpx_align = align_map[align]
            hwpx_text_align = align_map[text_align]
            
            # style 값 검증
            if not isinstance(style, int) or style < 1 or style > 8:
                raise ValueError("style은 1에서 8 사이의 정수여야 합니다.")
                
            # 스타일 매핑 (사용자 스타일 -> HWPX 스타일 ID)
            style_map = {
                1: "2",   # 보더 없음
                2: "3",   # 기본 테두리
                3: "12",  # 회색 배경
                4: "13",  # 양쪽 테두리
                5: "20",  # 상하단 굵은 테두리
                6: "22",  # 파란 테두리
                7: "23",  # 검은 테두리에 회색 배경
                8: "24"   # 빨간 테두리
            }
            
            # 스타일 설정 적용
            hwpx_style_id = style_map[style]
            
            # 헤더 스타일 설정
            hwpx_header_style_id = None
            if header_style is not None and header_style in style_map:
                hwpx_header_style_id = style_map[header_style]
            
            options["style"] = {
                "borderFillIDRef": hwpx_style_id,
                "cellBorderFillIDRef": hwpx_style_id,
                "headerFillIDRef": hwpx_header_style_id or hwpx_style_id
            }
            options["align"] = hwpx_align  # 정렬 옵션 추가
            options["text_align"] = hwpx_text_align  # 셀 내용 정렬 옵션 추가
            options["header_style"] = hwpx_header_style_id  # 헤더 스타일 옵션 추가
            
            # 빈 표 검사
            if not data or len(data) == 0:
                raise ValueError("빈 표는 추가할 수 없습니다")
            
            # 모든 행의 열 개수가 동일한지 확인
            col_count = len(data[0])
            if any(len(row) != col_count for row in data):
                raise ValueError("모든 행의 열 개수가 동일해야 합니다")
            
            def sanitize_cell_text(text):
                """표 셀의 텍스트를 안전하게 처리"""
                if not text:
                    return ""
                    
                # 문자열로 변환
                text = str(text).strip()
                
                # XML 특수문자 이스케이프
                text = text.replace('&', '&amp;')
                text = text.replace('<', '&lt;')
                text = text.replace('>', '&gt;')
                text = text.replace('"', '&quot;')
                text = text.replace("'", '&apos;')
                
                # 줄바꿈을 <hp:lineBreak/>로 변환
                if '\n' in text:
                    parts = text.split('\n')
                    text = '<hp:lineBreak/>'.join(p.strip() for p in parts if p.strip())
                    
                return text
            
            # 헤더와 데이터 처리
            processed_data = []
            if header:
                processed_data.append([sanitize_cell_text(cell) for cell in header])
            for row in data:
                processed_data.append([sanitize_cell_text(cell) for cell in row])
            
            # 표 데이터를 elements에 추가
            self.elements.append(('table', processed_data, False, options))
            
        except Exception as e:
            print(f"표 추가 실패: {str(e)}")
            raise

    def _resize_image_to_page_width(self, img_obj):
        """이미지를 페이지 너비에 맞게 크기 조정합니다."""
        # A4 페이지 크기 (mm)
        A4_WIDTH_MM = 210
        A4_HEIGHT_MM = 297
        
        # 여백을 제외한 실제 사용 가능한 너비 (페이지 너비의 80%)
        TARGET_WIDTH_MM = A4_WIDTH_MM * 0.8
        
        # 해상도를 96 DPI로 증가 (HWPX 표준에 더 가까움)
        MM_TO_PIXELS = 96 / 25.4  # 1mm = 96/25.4 pixels
        
        # 목표 너비 (픽셀)
        target_width = int(TARGET_WIDTH_MM * MM_TO_PIXELS)
        
        # 현재 이미지 크기
        current_width, current_height = img_obj.size
        
        # 이미지가 페이지 너비보다 작으면 원본 크기 유지
        if current_width <= target_width:
            return img_obj
        
        # 비율 계산
        ratio = target_width / current_width
        target_height = int(current_height * ratio)
        
        # 이미지 크기 조정 (고품질 리샘플링 사용)
        return img_obj.resize((target_width, target_height), Image.Resampling.BICUBIC)

    def save(self, output_path):
        try:
            # 저장 경로의 디렉토리가 없으면 생성
            os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

            # 디버깅을 위한 로그 추가
            # print("[DEBUG] Elements to save:")
            for element in self.elements:
                if element[0] == 'heading':
                    # print(f"[DEBUG] Heading: '{element[1]}' with style: {element[3].get('size', 'default')}")
                    pass

            # 제목이 없는 경우 자동으로 빈 제목 추가
            if not self.has_title:
                self.elements.insert(0, ('heading', "", False, {'size': '8'}))
                self.has_title = True

            with tempfile.TemporaryDirectory() as temp_dir:
                # 템플릿 파일 압축 해제
                with zipfile.ZipFile(self.template_path, 'r') as template_zip:
                    template_zip.extractall(temp_dir)

                # header.xml 파일 수정
                header_path = os.path.join(temp_dir, 'Contents', 'header.xml')
                with open(header_path, 'r', encoding='utf-8') as f:
                    header_content = f.read()
                    # print(f"[DEBUG] Original header.xml content: {header_content[:500]}...")  # 처음 500자만 출력

                # XML 선언과 네임스페이스 선언 확인
                if '<?xml' not in header_content:
                    header_content = '<?xml version="1.0" encoding="UTF-8" standalone="yes" ?>\n' + header_content

                # 스타일 정의 찾기
                style_start = header_content.find('<hh:style')
                style_end = header_content.find('</hh:style>')
                
                # 스타일 섹션 정의
                styles = '''<hh:style itemCnt="10">
                    <hh:stylePr id="0" paraPrIDRef="0" charPrIDRef="0" nextStyleIDRef="0" type="PARA" name="바탕글" engName="Normal" />
                    <hh:stylePr id="1" paraPrIDRef="1" charPrIDRef="1" nextStyleIDRef="1" type="PARA" name="본문" engName="Body" />
                    <hh:stylePr id="2" paraPrIDRef="2" charPrIDRef="2" nextStyleIDRef="2" type="PARA" name="제목" engName="Title" />
                    <hh:stylePr id="3" paraPrIDRef="3" charPrIDRef="3" nextStyleIDRef="3" type="PARA" name="제목 1" engName="Heading 1" />
                    <hh:stylePr id="4" paraPrIDRef="4" charPrIDRef="4" nextStyleIDRef="4" type="PARA" name="제목 2" engName="Heading 2" />
                    <hh:stylePr id="5" paraPrIDRef="5" charPrIDRef="5" nextStyleIDRef="5" type="PARA" name="제목 3" engName="Heading 3" />
                    <hh:stylePr id="6" paraPrIDRef="6" charPrIDRef="6" nextStyleIDRef="6" type="PARA" name="제목 4" engName="Heading 4" />
                    <hh:stylePr id="7" paraPrIDRef="7" charPrIDRef="7" nextStyleIDRef="7" type="PARA" name="제목 5" engName="Heading 5" />
                    <hh:stylePr id="8" paraPrIDRef="8" charPrIDRef="8" nextStyleIDRef="8" type="PARA" name="머리말" engName="Header" />
                    <hh:stylePr id="9" paraPrIDRef="9" charPrIDRef="9" nextStyleIDRef="9" type="PARA" name="표_셀" engName="Table Cell" />
                    <hh:stylePr id="10" paraPrIDRef="10" charPrIDRef="10" nextStyleIDRef="10" type="PARA" name="중간글씨" engName="Medium Text" />
                    <hh:stylePr id="11" paraPrIDRef="11" charPrIDRef="10" nextStyleIDRef="11" type="PARA" name="중간글씨가운데정렬" engName="Medium Center" />
                    <hh:stylePr id="12" paraPrIDRef="12" charPrIDRef="10" nextStyleIDRef="12" type="PARA" name="중간글씨오른쪽정렬" engName="Medium Right" />
                    <hh:stylePr id="13" paraPrIDRef="13" charPrIDRef="12" nextStyleIDRef="13" type="PARA" name="중간글씨굵게" engName="Medium Bold" />
                    <hh:stylePr id="14" paraPrIDRef="14" charPrIDRef="12" nextStyleIDRef="14" type="PARA" name="중간글씨굵게가운데정렬" engName="Medium Bold Center" />
                    <hh:stylePr id="15" paraPrIDRef="15" charPrIDRef="12" nextStyleIDRef="15" type="PARA" name="중간글씨굵게오른쪽정렬" engName="Medium Bold Right" />
                    <hh:stylePr id="16" paraPrIDRef="16" charPrIDRef="11" nextStyleIDRef="16" type="PARA" name="큰글씨보통" engName="Large Normal" />
                    <hh:stylePr id="17" paraPrIDRef="17" charPrIDRef="11" nextStyleIDRef="17" type="PARA" name="큰글씨보통가운데정렬" engName="Large Normal Center" />
                    <hh:stylePr id="18" paraPrIDRef="18" charPrIDRef="11" nextStyleIDRef="18" type="PARA" name="큰글씨보통오른쪽정렬" engName="Large Normal Right" />
                </hh:style>'''

                if style_start == -1 or style_end == -1:
                    # 스타일 섹션이 없으면 새로 추가
                    head_end = header_content.find('</hh:head>')
                    if head_end == -1:
                        raise ValueError("header.xml 파일의 구조가 올바르지 않습니다.")
                    
                    # 스타일 섹션을 </hh:head> 바로 앞에 추가
                    header_content = header_content[:head_end] + styles + header_content[head_end:]
                else:
                    # 기존 스타일 섹션 교체
                    header_content = header_content[:style_start] + styles + header_content[style_end + len('</hh:style>'):]

                # 수정된 header.xml 저장
                with open(header_path, 'w', encoding='utf-8') as f:
                    f.write(header_content)
                    # print("[DEBUG] Updated header.xml with styles")

                # 문단 모양 정의 추가
                para_pr_start = header_content.find('<hh:paraProperties')
                para_pr_end = header_content.find('</hh:paraProperties>')
                if para_pr_start == -1 or para_pr_end == -1:
                    raise ValueError("header.xml 파일의 구조가 올바르지 않습니다.")

                # 새로운 문단 모양 정의
                para_properties = '''<hh:paraProperties itemCnt="19">
                    <hh:paraPr id="0" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="0" next="0" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="1" tabPrIDRef="1" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="RIGHT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="2" tabPrIDRef="2" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="CENTER" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="3" tabPrIDRef="3" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="4" tabPrIDRef="4" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="5" tabPrIDRef="5" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="6" tabPrIDRef="6" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="7" tabPrIDRef="7" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="600" next="600" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="8" tabPrIDRef="8" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="0" next="0" />
                        <hh:lineSpacing type="PERCENT" value="160" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="9" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="CENTER" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="10" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="11" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="CENTER" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="12" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="RIGHT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="13" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="14" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="CENTER" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="15" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="130" />
                        <hh:align horizontal="RIGHT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="16" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="150" />
                        <hh:align horizontal="LEFT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="17" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="150" />
                        <hh:align horizontal="CENTER" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                    <hh:paraPr id="18" tabPrIDRef="0" condense="0" fontLineHeight="0" snapToGrid="0" suppressLineNumbers="0" checked="0">
                        <hh:margin left="0" right="0" prev="300" next="300" />
                        <hh:lineSpacing type="PERCENT" value="150" />
                        <hh:align horizontal="RIGHT" />
                        <hh:border borderFillIDRef="2" />
                    </hh:paraPr>
                </hh:paraProperties>'''
                # charProperties 섹션 전체 교체
                header_content = (
                    header_content[:para_pr_start] + 
                    para_properties +
                    header_content[para_pr_end + len('</hh:paraProperties>'):]
                )            
            
                # 글자 모양 정의 추가
                char_props_start = header_content.find('<hh:charProperties')
                char_props_end = header_content.find('</hh:charProperties>')
                
                if char_props_start == -1 or char_props_end == -1:
                    raise ValueError("header.xml 파일의 구조가 올바르지 않습니다.")
            

                # 글자 모양 정의
                char_props = '''<hh:charProperties itemCnt="13">
                    <hh:charPr id="0" height="1000" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                    </hh:charPr>
                    <hh:charPr id="1" height="1000" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                    </hh:charPr>
                    <hh:charPr id="2" height="1600" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="3" height="1400" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="4" height="1200" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="5" height="1100" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="6" height="1000" textColor="#2E74B5" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="7" height="900" textColor="#2E74B5" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="8" height="1000" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                    </hh:charPr>
                    <hh:charPr id="9" height="2400" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="9" latin="9" hanja="9" japanese="9" other="9" symbol="9" user="9"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                    <hh:charPr id="10" height="1200" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                    </hh:charPr>
                    <hh:charPr id="11" height="2400" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="9" latin="9" hanja="9" japanese="9" other="9" symbol="9" user="9"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                    </hh:charPr>
                    <hh:charPr id="12" height="1200" textColor="#000000" shadeColor="none" useFontSpace="0" useKerning="0" symMark="NONE" borderFillIDRef="2">
                        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
                        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
                        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
                        <hh:bold/>
                    </hh:charPr>
                </hh:charProperties>'''

                # charProperties 섹션 전체 교체
                header_content = (
                    header_content[:char_props_start] + 
                    char_props +
                    header_content[char_props_end + len('</hh:charProperties>'):]
                )

                # borderFills 섹션 찾기
                border_fills_start = header_content.find('<hh:borderFills')
                border_fills_end = header_content.find('</hh:borderFills>')
                
                if border_fills_start == -1 or border_fills_end == -1:
                    raise ValueError("header.xml 파일의 구조가 올바르지 않습니다.")

                # 새로운 테두리 스타일 정의
                border_fills = '''<hh:borderFills itemCnt="25">
                    <hh:borderFill id="1" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:rightBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:topBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:bottomBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                    </hh:borderFill>
                    <hh:borderFill id="2" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:rightBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:topBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:bottomBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="none" hatchColor="#999999" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="3" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                    </hh:borderFill>
                    <hh:borderFill id="4" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="5" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="6" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:rightBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:topBorder type="SOLID" width="0.4 mm" color="#5D83B0" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#CDD8E5" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#ADBFD5" hatchColor="#ADBFD5" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="7" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:rightBorder type="NONE" width="0.1 mm" color="#000000" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#CDD8E5" />
                        <hh:bottomBorder type="SOLID" width="0.4 mm" color="#5D83B0" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="8" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="9" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="10" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="11" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#5D5D5D" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#5D5D5D" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#5D5D5D" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="12" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="13" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="14" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="15" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="16" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="17" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="18" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.5 mm" color="#353535" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="19" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FFFFFF" />
                        <hh:bottomBorder type="SOLID" width="0.5 mm" color="#353535" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#DCDCDC" hatchColor="#DCDCDC" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="20" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:rightBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="21" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="NONE" width="0.12 mm" color="#353535" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#D6D6D6" />
                        <hh:topBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:bottomBorder type="SOLID" width="0.7 mm" color="#353535" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#FFFFFF" hatchColor="#FFFFFF" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="22" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#0000FF" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#0000FF" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#0000FF" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#0000FF" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                    </hh:borderFill>
                    <hh:borderFill id="23" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#000000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                        <hc:fillBrush>
                            <hc:winBrush faceColor="#D9D9D9" hatchColor="#000000" alpha="0" />
                        </hc:fillBrush>
                    </hh:borderFill>
                    <hh:borderFill id="24" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.12 mm" color="#FF0000" />
                        <hh:rightBorder type="SOLID" width="0.12 mm" color="#FF0000" />
                        <hh:topBorder type="SOLID" width="0.12 mm" color="#FF0000" />
                        <hh:bottomBorder type="SOLID" width="0.12 mm" color="#FF0000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                    </hh:borderFill>
                    <hh:borderFill id="25" threeD="0" shadow="0" centerLine="NONE" breakCellSeparateLine="0">
                        <hh:slash type="NONE" Crooked="0" isCounter="0" />
                        <hh:backSlash type="NONE" Crooked="0" isCounter="0" />
                        <hh:leftBorder type="SOLID" width="0.4 mm" color="#000000" />
                        <hh:rightBorder type="SOLID" width="0.4 mm" color="#000000" />
                        <hh:topBorder type="SOLID" width="0.4 mm" color="#000000" />
                        <hh:bottomBorder type="SOLID" width="0.4 mm" color="#000000" />
                        <hh:diagonal type="SOLID" width="0.1 mm" color="#000000" />
                    </hh:borderFill>
                </hh:borderFills>'''

                # borderFills 섹션 전체 교체
                header_content = (
                    header_content[:border_fills_start] + 
                    border_fills +
                    header_content[border_fills_end + len('</hh:borderFills>'):]
                )

                # 수정된 header.xml 저장
                with open(header_path, 'w', encoding='utf-8') as f:
                    f.write(header_content)

                # 이미지 처리를 위한 준비
                bindata_dir = os.path.join(temp_dir, 'BinData')
                contents_dir = os.path.join(temp_dir, 'Contents')
                preview_dir = os.path.join(temp_dir, 'Preview')
                
                # 필요한 모든 디렉토리 생성
                for directory in [bindata_dir, contents_dir, preview_dir]:
                    os.makedirs(directory, exist_ok=True)
                
                manifest_items = []
                image_count = 0

                # content.hpf 파일 수정 준비
                content_hpf_path = os.path.join(temp_dir, 'Contents', 'content.hpf')
                with open(content_hpf_path, 'r', encoding='utf-8') as f:
                    content_hpf = f.read()

                # 기존 이미지 항목 제거
                manifest_start = content_hpf.find('<opf:manifest>')
                manifest_end = content_hpf.find('</opf:manifest>')
                if manifest_start == -1 or manifest_end == -1:
                    raise ValueError("content.hpf 파일 구조가 올바르지 않습니다.")

                # 기본 manifest 항목
                manifest_items.append('''<opf:item id="header" href="Contents/header.xml" media-type="application/xml"/>
<opf:item id="section0" href="Contents/section0.xml" media-type="application/xml"/>
<opf:item id="settings" href="settings.xml" media-type="application/xml"/>''')

                # Section0.xml 파일 수정
                section_path = os.path.join(temp_dir, 'Contents', 'section0.xml')
                with open(section_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # 본문 시작 위치 찾기
                body_start = content.find('<hs:sec')
                body_end = content.find('</hs:sec>')
                if body_start == -1 or body_end == -1:
                    raise ValueError("문서 구조가 올바르지 않습니다.")
                
                # 기존 문서의 기본 설정 추출
                first_p_start = content.find('<hp:p', body_start)
                first_p_end = content.find('</hp:p>', first_p_start) + len('</hp:p>')
                if first_p_start == -1 or first_p_end == -1:
                    raise ValueError("문서 구조가 올바르지 않습니다.")

                # XML 네임스페이스와 기본 설정 보존
                header = content[body_start:first_p_end]
                
                # 문단 스타일 정의 추가
                header_xml_path = os.path.join(temp_dir, 'Contents', 'header.xml')
                with open(header_xml_path, 'r', encoding='utf-8') as f:
                    header_xml = f.read()
                
                # 문단 스타일 정의 위치 찾기
                para_pr_list_start = header_xml.find('<hh:paraStyleList>')
                para_pr_list_end = header_xml.find('</hh:paraStyleList>')
                
                if para_pr_list_start != -1 and para_pr_list_end != -1:
                    # 기존 문단 스타일 목록에 새 스타일 추가
                    para_style_definitions = '''
                    <hh:paraStyle id="1" name="오른쪽정렬" paraPrIDRef="1" charPrIDRef="0" nextStyleIDRef="1" />
                    <hh:paraStyle id="2" name="가운데정렬" paraPrIDRef="2" charPrIDRef="0" nextStyleIDRef="2" />
                    '''
                    
                    # 문단 스타일 목록에 추가
                    header_xml = header_xml[:para_pr_list_end] + para_style_definitions + header_xml[para_pr_list_end:]
                    
                    # 문단 속성 정의 위치 찾기
                    para_pr_def_start = header_xml.find('<hh:paraPrList>')
                    para_pr_def_end = header_xml.find('</hh:paraPrList>')
                    
                    if para_pr_def_start != -1 and para_pr_def_end != -1:
                        # 문단 속성 정의 추가
                        para_pr_definitions = '''
                        <hh:paraPr id="1">
                            <hh:margin left="0" right="0" prev="0" next="0" />
                            <hh:lineSpacing type="PERCENT" value="160" />
                            <hh:align horizontal="RIGHT" />
                        </hh:paraPr>
                        <hh:paraPr id="2">
                            <hh:margin left="0" right="0" prev="0" next="0" />
                            <hh:lineSpacing type="PERCENT" value="160" />
                            <hh:align horizontal="CENTER" />
                        </hh:paraPr>
                        '''
                        
                        # 문단 속성 목록에 추가
                        header_xml = header_xml[:para_pr_def_end] + para_pr_definitions + header_xml[para_pr_def_end:]
                        
                        # 수정된 header.xml 저장
                        with open(header_xml_path, 'w', encoding='utf-8') as f:
                            f.write(header_xml)
                
                # 새로운 본문 내용 생성
                new_body = header
                
                # 요소 순서대로 처리
                for element in self.elements:
                    element_type = element[0]
                    content_data = element[1]
                    page_break = element[2]
                    extra = element[3] if len(element) > 3 else {}
                    
                    if element_type == 'paragraph':
                        char_pr_id = extra.get('char_pr_id', "0")
                        
                        # 정렬 옵션 처리
                        align_option = "JUSTIFY"  # 기본값
                        para_pr_id = "0"  # 기본 문단 스타일
                        
                        # para_pr_id가 직접 지정된 경우 사용
                        if 'para_pr_id' in extra:
                            para_pr_id = extra['para_pr_id']
                            # para_pr_id에 따른 정렬 옵션 설정
                            if para_pr_id == "1":
                                align_option = "RIGHT"
                            elif para_pr_id == "2":
                                align_option = "CENTER"
                            elif para_pr_id == "10":
                                align_option = "LEFT"
                            elif para_pr_id == "11":
                                align_option = "CENTER"
                            elif para_pr_id == "12":
                                align_option = "RIGHT"
                        # 정렬 옵션으로 para_pr_id 설정
                        elif 'align' in extra:
                            align_value = extra['align'].upper() if isinstance(extra['align'], str) else extra['align']
                            if align_value in ['LEFT', 'CENTER', 'RIGHT', 'JUSTIFY']:
                                align_option = align_value
                                # 정렬에 따른 문단 스타일 ID 설정
                                if align_value == 'CENTER':
                                    para_pr_id = "2"  # 가운데 정렬 문단 스타일
                                elif align_value == 'RIGHT':
                                    para_pr_id = "1"  # 오른쪽 정렬 문단 스타일
                            elif isinstance(align_value, str) and align_value.lower() in ['left', 'center', 'right', 'justify']:
                                align_option = align_value.upper()
                                # 정렬에 따른 문단 스타일 ID 설정 (수정: 가운데와 오른쪽 정렬 매핑 변경)
                                if align_value.lower() == 'center':
                                    para_pr_id = "2"  # 가운데 정렬 문단 스타일
                                elif align_value.lower() == 'right':
                                    para_pr_id = "1"  # 오른쪽 정렬 문단 스타일
                        
                        # 각 줄을 별도의 문단으로 처리
                        if content_data:
                            for line in content_data.split('\n'):
                                # 문단 스타일 ID에 따라 정렬 옵션 설정 (수정: 가운데와 오른쪽 정렬 매핑 변경)
                                if para_pr_id == "2":
                                    align_option = "CENTER"
                                elif para_pr_id == "1":
                                    align_option = "RIGHT"
                                
                                paragraph_xml = f'''
                                <hp:p pageBreak="{1 if page_break else 0}" paraPrIDRef="{para_pr_id}" styleIDRef="{para_pr_id}">
                                    <hp:pPr>
                                        <hp:margin left="0" right="0" prev="200" next="200"/>
                                        <hp:lineSpacing type="PERCENT" value="120"/>
                                        <hp:lineWrap type="BREAK_WORD_BREAK_HANGUL"/>
                                        <hp:align horizontal="{align_option}"/>
                                    </hp:pPr>
                                    <hp:run charPrIDRef="{char_pr_id}">
                                        <hp:t>{line}</hp:t>
                                    </hp:run>
                                    <hp:linesegarray>
                                        <hp:lineseg textpos="0" vertpos="0" vertsize="1600" textheight="1600" 
                                                   baseline="1360" spacing="1600" horzpos="0" horzsize="42520" 
                                                   flags="1441792"/>
                                    </hp:linesegarray>
                                </hp:p>'''
                                new_body += paragraph_xml
                        else:
                            # 빈 문단 처리
                            # 문단 스타일 ID에 따라 정렬 옵션 설정 (수정: 가운데와 오른쪽 정렬 매핑 변경)
                            if para_pr_id == "2":
                                align_option = "CENTER"
                            elif para_pr_id == "1":
                                align_option = "RIGHT"
                                
                            paragraph_xml = f'''
                            <hp:p pageBreak="{1 if page_break else 0}" paraPrIDRef="{para_pr_id}" styleIDRef="{para_pr_id}">
                                <hp:pPr>
                                    <hp:margin left="0" right="0" prev="200" next="200"/>
                                    <hp:lineSpacing type="PERCENT" value="120"/>
                                    <hp:lineWrap type="BREAK_WORD_BREAK_HANGUL"/>
                                    <hp:align horizontal="{align_option}"/>
                                </hp:pPr>
                                <hp:run charPrIDRef="{char_pr_id}">
                                    <hp:t></hp:t>
                                </hp:run>
                                <hp:linesegarray>
                                    <hp:lineseg textpos="0" vertpos="0" vertsize="1600" textheight="1600" 
                                               baseline="1360" spacing="1600" horzpos="0" horzsize="42520" 
                                               flags="1441792"/>
                                </hp:linesegarray>
                            </hp:p>'''
                            new_body += paragraph_xml
                    
                    elif element_type == 'heading':
                        heading_style = extra.get('style', "2")
                        char_pr_id = extra.get('charPrIDRef', heading_style)  # 글자 모양 ID 가져오기
                        # 문서 제목(style 2)인 경우 가운데 정렬, 그 외는 왼쪽 정렬
                        align = 'CENTER' if heading_style == '2' else 'LEFT'
                        heading_xml = f'''
                        <hp:p pageBreak="{1 if page_break else 0}" paraPrIDRef="{heading_style}" styleIDRef="{heading_style}">
                            <hp:pPr>
                                <hp:margin left="0" right="0" prev="425" next="425"/>
                                <hp:lineSpacing type="PERCENT" value="160"/>
                                <hp:lineWrap type="BREAK_WORD_BREAK_HANGUL"/>
                                <hp:align horizontal="{align}"/>
                            </hp:pPr>
                            <hp:run charPrIDRef="{char_pr_id}">
                                <hp:t>{content_data}</hp:t>
                            </hp:run>
                            <hp:linesegarray>
                                <hp:lineseg textpos="0" vertpos="0" vertsize="2000" textheight="2000" 
                                           baseline="1700" spacing="2000" horzpos="0" horzsize="42520" 
                                           flags="393216"/>
                            </hp:linesegarray>
                        </hp:p>'''
                        new_body += heading_xml
                    
                    elif element_type == 'image':
                        # 이미지 처리
                        image_count += 1
                        image_path = content_data
                        
                        # 이미지 파일 복사 및 변환
                        img = Image.open(image_path)
                        img = HWPDocument._convert_image_to_rgb(img)  # RGB로 변환
                        
                        # 이미지 크기 계산 (A4 용지 기준 적절한 크기로 조정)
                        img_width, img_height = img.size
                        
                        # 이미지 크기 설정 (extra에서 가져오거나 기본값 사용)
                        img_options = extra if extra else {}
                        custom_width = img_options.get('width', 41550)  # 기본값: A4 용지 너비에 맞춤
                        # 높이는 원본 비율 유지
                        custom_height = img_options.get('height', int(custom_width * (img_height / img_width)))
                        
                        # 원본 이미지 확장자 유지
                        _, ext = os.path.splitext(image_path)
                        img_filename = f'image{image_count}{ext.lower()}'
                        img.save(os.path.join(bindata_dir, img_filename))

                        # content.hpf에 이미지 항목 추가
                        media_type = f"image/{ext[1:].lower()}"
                        manifest_items.append(
                            f'<opf:item id="image{image_count}" href="BinData/{img_filename}" '
                            f'media-type="{media_type}" isEmbeded="1"/>'
                        )

                        # 이미지 문단 추가
                        image_xml = f'''
                        <hp:p pageBreak="{1 if page_break else 0}" paraPrIDRef="0" styleIDRef="0">
                            <hp:pPr>
                                <hp:margin left="0" right="0" prev="200" next="200"/>
                            </hp:pPr>
                            <hp:run charPrIDRef="7">
                                <hp:pic id="{1000000 + image_count}" zOrder="{image_count}" 
                                       numberingType="PICTURE" textWrap="SQUARE" textFlow="BOTH_SIDES" 
                                       lock="0" dropcapstyle="None" href="" groupLevel="0" 
                                       instid="{682337706 + image_count}" reverse="0">
                                    <hp:offset x="0" y="0"/>
                                    <hp:orgSz width="{custom_width}" height="{custom_height}"/>
                                    <hp:curSz width="0" height="0"/>
                                    <hp:flip horizontal="0" vertical="0"/>
                                    <hp:rotationInfo angle="0" centerX="{custom_width//2}" centerY="{custom_height//2}" 
                                                    rotateimage="1"/>
                                    <hp:renderingInfo>
                                        <hc:transMatrix e1="1" e2="0" e3="0" e4="0" e5="1" e6="0"/>
                                        <hc:scaMatrix e1="1" e2="0" e3="0" e4="0" e5="1" e6="0"/>
                                        <hc:rotMatrix e1="1" e2="0" e3="0" e4="0" e5="1" e6="0"/>
                                    </hp:renderingInfo>
                                    <hp:imgRect>
                                        <hc:pt0 x="0" y="0"/>
                                        <hc:pt1 x="{custom_width}" y="0"/>
                                        <hc:pt2 x="{custom_width}" y="{custom_height}"/>
                                        <hc:pt3 x="0" y="{custom_height}"/>
                                    </hp:imgRect>
                                    <hp:imgClip left="0" right="96000" top="0" bottom="77400"/>
                                    <hp:inMargin left="0" right="0" top="0" bottom="0"/>
                                    <hc:img binaryItemIDRef="image{image_count}" bright="0" contrast="0" 
                                           effect="REAL_PIC" alpha="0"/>
                                    <hp:effects/>
                                    <hp:sz width="{custom_width}" widthRelTo="ABSOLUTE" height="{custom_height}" 
                                          heightRelTo="ABSOLUTE" protect="0"/>
                                    <hp:pos treatAsChar="0" affectLSpacing="0" flowWithText="1" 
                                           allowOverlap="1" holdAnchorAndSO="0" vertRelTo="PARA" 
                                           horzRelTo="PARA" vertAlign="TOP" horzAlign="CENTER" 
                                           vertOffset="0" horzOffset="0"/>
                                    <hp:outMargin left="0" right="0" top="0" bottom="0"/>
                                    <hp:shapeComment>그림입니다.</hp:shapeComment>
                                </hp:pic>
                            </hp:run>
                            <hp:linesegarray>
                                <hp:lineseg textpos="0" vertpos="0" vertsize="1000" textheight="1000" 
                                           baseline="850" spacing="600" horzpos="0" horzsize="42520" 
                                           flags="393216"/>
                            </hp:linesegarray>
                        </hp:p>'''
                        new_body += image_xml

                    elif element_type == 'table':
                        table_data, page_break, options = content_data, page_break, extra
                        row_count = len(table_data)
                        col_count = len(table_data[0])
                        style = options.get("style", {
                            "borderFillIDRef": "3",
                            "cellBorderFillIDRef": "3",
                            "headerFillIDRef": "3"
                        })
                        
                        # 표 XML 시작
                        table_xml = f'''
                        <hp:p pageBreak="{1 if page_break else 0}" paraPrIDRef="0" styleIDRef="0">
                            <hp:run charPrIDRef="7">
                                <hp:tbl id="{1000000 + len(new_body)}" zOrder="0" numberingType="TABLE" 
                                       textWrap="TOP_AND_BOTTOM" textFlow="BOTH_SIDES" lock="0" 
                                       dropcapstyle="None" pageBreak="CELL" repeatHeader="1" 
                                       rowCnt="{row_count}" colCnt="{col_count}" cellSpacing="0" 
                                       borderFillIDRef="{style['borderFillIDRef']}" noAdjust="0">
                                    <hp:sz width="42520" widthRelTo="ABSOLUTE" height="5000" 
                                          heightRelTo="ABSOLUTE" protect="0"/>
                                    <hp:pos treatAsChar="0" affectLSpacing="0" flowWithText="1" 
                                           allowOverlap="1" holdAnchorAndSO="0" vertRelTo="PARA" 
                                           horzRelTo="PARA" vertAlign="TOP" horzAlign="{options.get('align', 'CENTER')}" 
                                           vertOffset="0" horzOffset="0"/>
                                    <hp:outMargin left="283" right="283" top="283" bottom="283"/>
                                    <hp:inMargin left="510" right="510" top="141" bottom="141"/>'''

                        # 각 행 추가
                        for row_idx, row in enumerate(table_data):
                            table_xml += '<hp:tr>'
                            for col_idx, cell in enumerate(row):
                                table_xml += f'''
                                    <hp:tc name="" header="0" hasMargin="0" protect="0" editable="0" 
                                          dirty="0" borderFillIDRef="{
                                            options.get('header_style') if row_idx == 0 and options.get('header_style') else style['cellBorderFillIDRef']
                                          }">
                                        <hp:subList id="" textDirection="HORIZONTAL" lineWrap="BREAK" 
                                                   vertAlign="CENTER" linkListIDRef="0" 
                                                   linkListNextIDRef="0" textWidth="0" textHeight="0" 
                                                   hasTextRef="0" hasNumRef="0">
                                            <hp:p paraPrIDRef="{
                                                '0' if options.get('text_align', 'CENTER') == 'LEFT' else 
                                                '1' if options.get('text_align', 'CENTER') == 'RIGHT' else 
                                                '2'
                                            }" styleIDRef="{
                                                '0' if options.get('text_align', 'CENTER') == 'LEFT' else 
                                                '1' if options.get('text_align', 'CENTER') == 'RIGHT' else 
                                                '2'
                                            }" pageBreak="0" columnBreak="0" merged="0">
                                                <hp:pPr>
                                                    <hp:align horizontal="{options.get('text_align', 'CENTER')}"/>
                                                    <hp:margin left="0" right="0" prev="300" next="300"/>
                                                </hp:pPr>
                                                <hp:run charPrIDRef="1">
                                                    <hp:t>{str(cell)}</hp:t>
                                                </hp:run>
                                            </hp:p>
                                        </hp:subList>
                                        <hp:cellAddr colAddr="{col_idx}" rowAddr="{row_idx}"/>
                                        <hp:cellSpan colSpan="1" rowSpan="1"/>
                                        <hp:cellSz width="{42520 // col_count}" height="2000"/>
                                        <hp:cellMargin left="510" right="510" top="300" bottom="300"/>
                                    </hp:tc>'''
                            table_xml += '</hp:tr>'

                        # 표 XML 종료
                        table_xml += '''
                                </hp:tbl>
                            </hp:run>
                        </hp:p>'''
                        
                        new_body += table_xml

                    elif element_type == 'field':
                        if content == 'page_number':
                            # 페이지 번호 필드 추가
                            # HWPX 필드 코드 형식으로 변환
                            pass
                        elif content == 'page_ref':
                            # 페이지 참조 필드 추가
                            # HWPX 필드 코드 형식으로 변환
                            pass
                    elif element_type == 'bookmark':
                        # 북마크 추가
                        # HWPX 북마크 형식으로 변환
                        pass

                new_body += '</hs:sec>'

                # 전체 내용 교체
                content = content[:body_start] + new_body

                # 수정된 내용 저장
                with open(section_path, 'w', encoding='utf-8') as f:
                    f.write(content)

                # content.hpf 파일 업데이트
                new_manifest = '<opf:manifest>\n' + '\n'.join(manifest_items) + '\n</opf:manifest>'
                content_hpf = (
                    content_hpf[:manifest_start] + 
                    new_manifest + 
                    content_hpf[manifest_end + len('</opf:manifest>'):]
                )
                with open(content_hpf_path, 'w', encoding='utf-8') as f:
                    f.write(content_hpf)

                # Preview/PrvText.txt 수정
                preview_text = '\n'.join(
                    content_data for type_, content_data, _, _ in self.elements 
                    if type_ == 'paragraph' and content_data
                )
                preview_path = os.path.join(temp_dir, 'Preview', 'PrvText.txt')
                with open(preview_path, 'w', encoding='utf-8') as f:
                    f.write(preview_text)

                # 새로운 HWPX 파일 생성
                with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as output_zip:
                    for root, _, files in os.walk(temp_dir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arc_name = os.path.relpath(file_path, temp_dir)
                            output_zip.write(file_path, arc_name)

            # 임시 파일 정리
            for temp_file in self._temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.unlink(temp_file)
                except Exception as e:
                    print(f"Warning: Failed to delete temp file {temp_file}: {e}")
            self._temp_files.clear()  # 리스트 비우기

            return True
        except Exception as e:
            raise print(f"문서 저장 중 오류 발생: {str(e)}")

    def add_field(self, field_type: str, bookmark: str = None) -> None:
        """문서에 필드 코드 추가
        
        Args:
            field_type (str): 필드 타입 (page_ref, page_number 등)
            bookmark (str, optional): 참조할 북마크 이름
        """
        try:
            # HWPX 필드 코드 형식에 맞게 구성
            if field_type == "page_ref" and bookmark:
                field_code = f"{{REF {bookmark} \\p}}"  # 페이지 참조 필드
            elif field_type == "page_number":
                field_code = "{PAGE}"  # 현재 페이지 번호
            else:
                raise ValueError(f"지원하지 않는 필드 타입: {field_type}")
            
            # 필드 코드를 포함하는 문단 추가
            self.add_paragraph(field_code, options={
                'font_name': '맑은 고딕',
                'font_size': 10,
                'field_code': True  # 필드 코드임을 표시
            })
            
        except Exception as e:
            print(f"필드 추가 실패: {str(e)}")
            raise
    
    def add_bookmark(self, bookmark_name: str) -> None:
        """문서에 북마크 추가
        
        Args:
            bookmark_name (str): 북마크 이름
        """
        try:
            # HWPX 북마크 형식에 맞게 구성
            bookmark_code = f"{{BM={bookmark_name}}}"
            
            # 북마크 코드 추가
            self.add_paragraph(bookmark_code, options={
                'bookmark': True  # 북마크임을 표시
            })
            
        except Exception as e:
            print(f"북마크 추가 실패: {str(e)}")
            raise

    def add_styled_paragraph(self, text: str, style: str = 'normal', page_break: bool = False, additional_options: Dict = None) -> None:
        """미리 정의된 스타일을 사용하여 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            style (str): 사용할 스타일 이름 
                - 'normal': 기본 스타일
                - 'large': 큰 글씨
                - 'medium': 중간 글씨
                - 'center': 가운데 정렬
                - 'right': 오른쪽 정렬
                - 'medium_center': 중간 글씨 가운데 정렬
                - 'medium_right': 중간 글씨 오른쪽 정렬
                - 'large_center': 큰 글씨 가운데 정렬
                - 'large_right': 큰 글씨 오른쪽 정렬
                - 'emphasis': 강조 (굵게)
                - 'quote': 인용문 (들여쓰기)
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션 (기본 스타일에 추가로 적용)
        """
        if style not in self.paragraph_styles:
            raise ValueError(f"지원하지 않는 스타일입니다: {style}. 사용 가능한 스타일: {', '.join(self.paragraph_styles.keys())}")
            
        # 기본 스타일 옵션 가져오기
        options = self.paragraph_styles[style].copy()
        
        # 추가 옵션 적용
        if additional_options:
            options.update(additional_options)
            
        # 문단 추가
        self.add_paragraph(text, page_break=page_break, options=options)
        
    def add_large_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """큰 글씨 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="large", page_break=page_break, additional_options=additional_options)
        
    def add_centered_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """가운데 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="center", page_break=page_break, additional_options=additional_options)
        
    def add_large_centered_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """큰 글씨 가운데 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="large_center", page_break=page_break, additional_options=additional_options)
        
    def add_emphasized_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """강조(굵게) 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="emphasis", page_break=page_break, additional_options=additional_options)
        
    def add_quote(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """인용문 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="quote", page_break=page_break, additional_options=additional_options)
        
    def add_medium_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 글씨 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium", page_break=page_break, additional_options=additional_options)
        
    def add_medium_centered_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 글씨 가운데 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium_center", page_break=page_break, additional_options=additional_options)
        
    def add_medium_right_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 글씨 오른쪽 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium_right", page_break=page_break, additional_options=additional_options)
        
    def add_medium_bold_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 크기 굵은체 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium_bold", page_break=page_break, additional_options=additional_options)
        
    def add_medium_bold_centered_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 크기 굵은체 가운데 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium_bold_center", page_break=page_break, additional_options=additional_options)
        
    def add_medium_bold_right_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """중간 크기 굵은체 오른쪽 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="medium_bold_right", page_break=page_break, additional_options=additional_options)
        
    def add_large_normal_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """큰 글씨 보통체 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="large_normal", page_break=page_break, additional_options=additional_options)
        
    def add_large_normal_centered_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """큰 글씨 보통체 가운데 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="large_normal_center", page_break=page_break, additional_options=additional_options)
        
    def add_large_normal_right_text(self, text: str, page_break: bool = False, additional_options: Dict = None) -> None:
        """큰 글씨 보통체 오른쪽 정렬 스타일로 문단 추가
        
        Args:
            text (str): 추가할 텍스트
            page_break (bool): 페이지 나누기 여부
            additional_options (Dict, optional): 추가 옵션
        """
        self.add_styled_paragraph(text, style="large_normal_right", page_break=page_break, additional_options=additional_options)

class PDFDocument:
    def __init__(self):
        try:
            self.font_path = safe_path_join(os.path.expanduser("~"), ".airun", "Pretendard-Regular.ttf")
            self.font_path_bold = safe_path_join(os.path.expanduser("~"), ".airun", "Pretendard-Bold.ttf")
            if not os.path.exists(self.font_path):
                raise FileNotFoundError(f"Font file not found: {self.font_path}")
            if not os.path.exists(self.font_path_bold):
                raise FileNotFoundError(f"Font file not found: {self.font_path_bold}")
            
            # font registration
            pdfmetrics.registerFont(TTFont('Pretendard', self.font_path))
            pdfmetrics.registerFont(TTFont('Pretendard-Bold', self.font_path_bold))
            
            # 폰트 패밀리 등록
            pdfmetrics.registerFontFamily(
                'Pretendard',
                normal='Pretendard',
                bold='Pretendard-Bold'
            )
            
            self.elements = []
            
            # 스타일 설정
            self.styles = getSampleStyleSheet()
            
            # 기본 스타일 복사 및 한글 폰트 적용
            self.styles.add(ParagraphStyle(
                name='Korean',
                parent=self.styles['Normal'],
                fontName='Pretendard',
                fontSize=10,
                leading=15,
                encoding='utf-8',  # 인코딩 명시적 지정
                wordWrap='CJK'     # CJK 워드랩 사용
            ))
            
            # 제목 레벨별 크기 설정
            heading_sizes = {
                1: {'fontSize': 16, 'leading': 24, 'spaceBefore': 20, 'spaceAfter': 10},
                2: {'fontSize': 14, 'leading': 21, 'spaceBefore': 15, 'spaceAfter': 8},
                3: {'fontSize': 12, 'leading': 18, 'spaceBefore': 12, 'spaceAfter': 6},
                4: {'fontSize': 10, 'leading': 15, 'spaceBefore': 10, 'spaceAfter': 5},
                5: {'fontSize': 10, 'leading': 12, 'spaceBefore': 8, 'spaceAfter': 4}
            }
            
            # 기존 Heading 스타일을 유지하면서 한글 제목 스타일 추가
            for level, sizes in heading_sizes.items():
                # 기존 Heading 스타일 수정
                heading_style = self.styles[f'Heading{level}']
                heading_style.fontSize = sizes['fontSize']
                heading_style.leading = sizes['leading']
                heading_style.spaceBefore = sizes['spaceBefore']
                heading_style.spaceAfter = sizes['spaceAfter']
                
                # 한글 제목 스타일 추가
                self.styles.add(ParagraphStyle(
                    name=f'KoreanHeading{level}',
                    parent=heading_style,  # 기존 Heading 스타일 상속
                    fontName='Pretendard-Bold',
                    fontSize=sizes['fontSize'],  # 레벨별 폰트 크기
                    leading=sizes['leading'],    # 레벨별 줄간격
                    spaceBefore=sizes['spaceBefore'],  # 레벨별 상단 여백
                    spaceAfter=sizes['spaceAfter'],    # 레벨별 하단 여백
                    alignment=heading_style.alignment  # 기존 정렬 유지
                ))
            
            # 목차 초기화
            self.toc = TableOfContents()
            self.toc.dotsMinLevel = 0  # 모든 레벨에 점선 표시
            
            # 목차 스타일 설정
            for i in range(1, 4):
                toc_style = ParagraphStyle(
                    name=f'TOCHeading{i}',
                    parent=self.styles[f'KoreanHeading{i}'],
                    fontName='Pretendard',
                    fontSize=12 - (i-1),  # 목차 항목 크기: 12pt, 11pt, 10pt
                    leading=16,           # 목차 줄간격
                    leftIndent=20*(i-1),  # 들여쓰기
                    firstLineIndent=0,
                    spaceBefore=3,
                    spaceAfter=3
                )
                self.styles.add(toc_style)
            
            self.toc.levelStyles = [
                self.styles[f'TOCHeading{i}'] for i in range(1, 4)
            ]
            
            # 머릿말/꼬릿말 설정
            self.header_text = ""
            self.footer_text = ""
            self.header_align = "left"
            self.footer_align = "left"
            
        except Exception as e:
            print(f"[ERROR] Failed to initialize PDF document: {str(e)}")
            raise

    def _header_footer(self, canvas, doc):
        canvas.saveState()
        
        # 머릿말 추가
        if self.header_text:
            canvas.setFont('Pretendard', 9)
            if self.header_align == 'center':
                canvas.drawCentredString(A4[0]/2, A4[1] - 20*mm, self.header_text)
            elif self.header_align == 'right':
                canvas.drawRightString(A4[0] - 20*mm, A4[1] - 20*mm, self.header_text)
            else:
                canvas.drawString(20*mm, A4[1] - 20*mm, self.header_text)
            
            # 구분선
            canvas.line(20*mm, A4[1] - 25*mm, A4[0] - 20*mm, A4[1] - 25*mm)

        # 꼬릿말 추가
        if self.footer_text:
            canvas.setFont('Pretendard', 9)
            if self.footer_align == 'center':
                canvas.drawCentredString(A4[0]/2, 15*mm, self.footer_text)
            elif self.footer_align == 'right':
                canvas.drawRightString(A4[0] - 20*mm, 15*mm, self.footer_text)
            else:
                canvas.drawString(20*mm, 15*mm, self.footer_text)
            
            # 구분선
            canvas.line(20*mm, 20*mm, A4[0] - 20*mm, 20*mm)

        # 페이지 번호
        canvas.setFont('Pretendard', 9)
        canvas.drawCentredString(A4[0]/2, 10*mm, f"- {doc.page} -")
        
        canvas.restoreState()

    def set_header(self, text: str, align: str = 'left'):
        self.header_text = text
        self.header_align = align

    def set_footer(self, text: str, align: str = 'left'):
        self.footer_text = text
        self.footer_align = align

    def add_heading(self, text: str, level: int = 1, align: str = 'left', use_korean: bool = True):
        """Add a heading to the document and register it in the TOC."""
        if not text:
            return
            
        # 스타일 선택 (한글 또는 기본)
        style_name = f'KoreanHeading{level}' if use_korean else f'Heading{level}'
        base_style = self.styles[style_name]
        
        # 정렬이 다른 경우 새로운 스타일 생성
        if align != 'left':
            style = ParagraphStyle(
                f'{style_name}_{align}',
                parent=base_style,
                alignment=self._get_alignment(align)
            )
        else:
            style = base_style
        
        # 제목 추가
        heading = Paragraph(text, style)
        self.elements.append(heading)
        self.elements.append(Spacer(1, 6))
        
        # 목차 항목 추가 - 임시로 페이지 번호 1 사용
        # multiBuild 과정에서 실제 페이지 번호로 업데이트됨
        self.toc.addEntry(level-1, text, 1)

    def _get_alignment(self, align):
        if align == 'left':
            return 0
        elif align == 'center':
            return 1
        elif align == 'right':
            return 2
        else:
            raise ValueError("Invalid alignment format")

    def add_paragraph(self, text: str, font_size: int = 10, align: str = 'left'):
        """Add content text to the document."""
        if not text:
            return
            
    # 마크다운 스타일 강조 문자 제거
        def clean_markdown(text: str) -> str:
            # 볼드체 (**text**) 제거
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
            # 이탤릭체 (*text* 또는 _text_) 제거
            text = re.sub(r'\*(.*?)\*', r'\1', text)
            text = re.sub(r'_(.*?)_', r'\1', text)
            # 인라인 코드 (`text`) 제거
            text = re.sub(r'`(.*?)`', r'\1', text)
            return text
        
        style = ParagraphStyle(
            'Content',
            parent=self.styles['Normal'],
            fontName='Pretendard',
            fontSize=font_size,
            leading=font_size * 2.0,
            alignment={'left': 0, 'center': 1, 'right': 2}[align],
            spaceAfter=10,
            leftIndent=10*mm,  # 좌측 여백 추가
            rightIndent=10*mm  # 우측 여백 추가
        )
        
        # 문단 나누기
        paragraphs = text.split('\n')
        for paragraph in paragraphs:
            if paragraph.strip():
                # 마크다운 강조 문자 제거 후 Paragraph 객체 생성
                cleaned_text = clean_markdown(paragraph)
                p = Paragraph(cleaned_text, style)
                self.elements.append(p)

    def add_tooltip(self, text: str, header_text: str = "AI.RUN 2025. - Empowering Your AI Journey", font_size: int = 9, background_color=colors.white):
        """Add a tooltip-style text box with blue border and header.
        
        
        Args:
            text (str): Text to display in the tooltip
            font_size (int, optional): Font size for the tooltip text. Defaults to 9.
            background_color (Color, optional): Background color of the tooltip. Defaults to white.
        """
        try:
            if not text:
                return
                
            # 헤더 스타일 설정
            header_style = ParagraphStyle(
                'TooltipHeader',
                parent=self.styles['Normal'],
                fontName='Pretendard-Bold',  # 폰트 변경
                fontSize=font_size,
                leading=font_size * 1.5,
                alignment=0,  # 왼쪽 정렬
                textColor=colors.white,
                encoding='utf-8',
                wordWrap='CJK',
                allowWidows=0,
                allowOrphans=0
            )
            
            # 본문 스타일 설정
            body_style = ParagraphStyle(
                'TooltipBody',
                parent=self.styles['Normal'],
                fontName='Pretendard',  # 폰트 변경
                fontSize=font_size,
                leading=font_size * 1.5,
                alignment=0,  # 왼쪽 정렬
                spaceBefore=3,
                spaceAfter=3,
                leftIndent=5,
                rightIndent=5,
                encoding='utf-8',
                wordWrap='CJK',
                allowWidows=0,
                allowOrphans=0
            )
            
            # 특수 문자 처리를 위한 텍스트 전처리
            def preprocess_text(text):
                # 특수 문자 매핑 정의 (수정된 버전)
                char_map = {
                    '–': '-',    # en dash를 일반 하이픈으로
                    '—': '-',    # em dash를 일반 하이픈으로
                    '−': '-',    # 유니코드 마이너스를 일반 하이픈으로
                    '"': '"',    # 직선형 따옴표를 곡선형으로
                    '"': '"',
                    "'": "'",    # 직선형 작은따옴표를 곡선형으로
                    "'": "'",
                    '...': '…',  # 마침표 3개를 말줄임표로
                    '©': '(c)',  # 저작권 기호
                    '®': '(R)',  # 등록 상표
                    '™': '(TM)', # 상표
                    '•': '-',    # 글머리 기호를 하이픈으로
                    '·': '-',    # 가운뎃점을 하이픈으로
                    '×': 'x',    # 곱하기 기호를 x로
                    '÷': '/',    # 나누기 기호를 /로
                    '±': '+-',   # 플러스마이너스
                    '≠': '!=',   # 같지 않음
                    '≤': '<=',   # 작거나 같음
                    '≥': '>=',   # 크거나 같음
                    '∞': 'inf',  # 무한대
                    '°': 'deg',  # 도
                    '′': "'",    # 프라임
                    '″': '"',    # 더블 프라임
                    '→': '->',   # 화살표
                    '←': '<-',
                    '↑': '^',
                    '↓': 'v',
                    '⇒': '=>',   # 이중 화살표
                    '⇐': '<=',
                    '⇔': '<=>',
                }
                
                # 특수 문자 변환
                for old, new in char_map.items():
                    text = text.replace(old, new)
                
                # 연속된 공백 정리
                text = ' '.join(text.split())
                
                # 줄바꿈 문자 정리
                text = text.replace('\r\n', '\n').replace('\r', '\n')
                
                # HTML 엔티티 디코딩
                text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
                
                return text
                
            # 헤더와 본문 Paragraph 객체 생성
            header = Paragraph(preprocess_text(header_text), header_style)
            
            # 텍스트 줄바꿈 처리
            paragraphs = []
            for line in text.split('\n'):
                if line.strip():  # 빈 줄 제외
                    paragraphs.append(Paragraph(preprocess_text(line), body_style))
            
            # 테이블 데이터 준비 (2행 1열)
            table_data = [
                [header],  # 헤더 행
                [paragraphs[0] if paragraphs else '']  # 첫 번째 문단
            ]
            
            # 나머지 문단들을 추가
            for p in paragraphs[1:]:
                table_data.append([p])
            
            # 테이블 너비 계산 (문서 너비에서 좌우 여백 고려)
            available_width = A4[0] - 60*mm  # 좌우 각각 10mm 추가 여백
            
            # 테이블 생성
            table = Table(table_data, colWidths=[available_width])
            table.hAlign = 'CENTER'  # 테이블 중앙 정렬
            
            # 테이블 스타일 설정
            table_style = TableStyle([
                # 전체 테두리
                ('BOX', (0, 0), (-1, -1), 0.5, colors.gray),
                
                # 헤더 행 스타일
                ('BACKGROUND', (0, 0), (-1, 0), colors.black),  # 헤더 배경색
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),  # 헤더 텍스트 색상
                
                # 본문 행 스타일
                ('BACKGROUND', (0, 1), (-1, -1), background_color),  # 본문 배경색
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),  # 본문 텍스트 색상
                
                # 정렬 설정
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),  # 전체 왼쪽 정렬
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),  # 수직 중앙 정렬
                
                # 패딩 설정
                ('LEFTPADDING', (0, 0), (-1, -1), 10),  # 왼쪽 패딩
                ('RIGHTPADDING', (0, 0), (-1, -1), 10),  # 오른쪽 패딩
                ('TOPPADDING', (0, 0), (-1, 0), 5),  # 헤더 상단 패딩
                ('BOTTOMPADDING', (0, 0), (-1, 0), 5),  # 헤더 하단 패딩
                ('TOPPADDING', (0, 1), (-1, -1), 3),  # 본문 상단 패딩
                ('BOTTOMPADDING', (0, 1), (-1, -1), 3),  # 본문 하단 패딩
            ])
            
            table.setStyle(table_style)
            
            # 테이블 추가
            self.elements.append(table)
            self.elements.append(Spacer(1, 8*mm))  # 툴팁 아래 여백
            
        except Exception as e:
            print(f"[ERROR] Failed to add tooltip: {str(e)}")
            raise

    def add_table(self, data, header=None):
        """Add a table to the document."""
        try:
            import pandas as pd
            import re
            
            # 특수문자 처리 함수
            def clean_text(text):
                if pd.isna(text):
                    return ''
                text = str(text).strip()
                text = re.sub(r'[^\uAC00-\uD7A3a-zA-Z0-9\s\.,\-\(\)/%]', '', text)
                return text
            
            # DataFrame 변환 및 데이터 전처리
            if isinstance(data, pd.DataFrame):
                df = data.applymap(clean_text)
            else:
                if header:
                    df = pd.DataFrame(data, columns=header)
                else:
                    df = pd.DataFrame(data)
                df = df.applymap(clean_text)
            
            # 컬럼 너비 계산 수정
            available_width = A4[0] - 60*mm  # 좌우 각각 10mm 추가 여백
            col_width = available_width / len(df.columns)
            
            # 데이터를 Paragraph 객체로 변환
            def to_paragraph(text):
                style = ParagraphStyle(
                    'TableCell',
                    fontName='Pretendard',
                    fontSize=10,
                    leading=12,
                    alignment=1,
                    wordWrap='CJK'
                )
                return Paragraph(str(text), style)
            
            # 테이블 데이터 준비
            table_data = []
            if header is not None:
                table_data.append([to_paragraph(col) for col in df.columns])
            for _, row in df.iterrows():
                table_data.append([to_paragraph(cell) for cell in row])
            
            # 테이블 생성
            table = Table(table_data, colWidths=[col_width] * len(df.columns))
            table.hAlign = 'CENTER'  # 테이블 중앙 정렬
            
            # 테이블 스타일 설정
            style = [
                ('FONT', (0, 0), (-1, -1), 'Pretendard'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                ('PADDING', (0, 0), (-1, -1), 6),
                ('WORDWRAP', (0, 0), (-1, -1), True),
            ]
            
            if header is not None:
                style.extend([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                ])
            
            table.setStyle(TableStyle(style))
            self.elements.append(table)
            self.elements.append(Spacer(1, 10*mm))
            
        except Exception as e:
            print(f"Error in add_table: {str(e)}")
            raise

    def add_image(self, image_path: str, width: float = None, height: float = None, scale_small_images: bool = False):
        """
        Add an image to the PDF document.
        """
        try:
            temp_file = None
            
            try:
                # 페이지 최대 너비와 높이 (여백 고려)
                max_width = A4[0] - 60*mm  # 좌우 각각 30mm 여백
                max_height = A4[1] - 90*mm  # 상하 여백을 더 넉넉하게 설정
                
                # 이미지 크기 계산
                from PIL import Image as PILImage
                img = PILImage.open(image_path)
                img_w, img_h = img.size
                
                # 이미지 비율 계산
                aspect_ratio = img_w / img_h
                
                # 이미지 크기 조정
                if width is None and height is None:
                    # 먼저 높이를 기준으로 조정
                    if img_h > max_height:
                        height = max_height
                        width = height * aspect_ratio
                    else:
                        height = img_h
                        width = img_w
                    
                    # 너비가 최대 너비를 초과하는 경우 다시 조정
                    if width > max_width:
                        width = max_width
                        height = width / aspect_ratio
                    
                    # 최종 안전 검사
                    if height > max_height:
                        height = max_height
                        width = height * aspect_ratio
                    
                    # 최소 크기 설정 (너무 작은 이미지 방지)
                    if scale_small_images:
                        min_width = max_width * 0.3
                        if width < min_width:
                            width = min_width
                            height = width / aspect_ratio
                
                # reportlab Image 객체 생성
                from reportlab.platypus import Image as RLImage
                img = RLImage(image_path, width=width, height=height)
                img.hAlign = 'CENTER'
                
                # 이미지 추가
                self.elements.append(img)
                self.elements.append(Spacer(1, 10*mm))
                
            finally:
                if temp_file and os.path.exists(temp_file):
                    os.unlink(temp_file)
                    
        except Exception as e:
            print(f"[ERROR] Failed to add image: {str(e)}")

    def add_spacing(self, points: int):
        """Add vertical spacing."""
        self.elements.append(Spacer(1, points))

    def add_page_break(self):
        """Add a page break to the document."""
        try:
            self.elements.append(PageBreak())
        except Exception as e:
            print(f"[WARNING] Failed to add page break: {str(e)}")

    def save(self, filename: str, include_toc: bool = False):
        try:
            doc = SimpleDocTemplate(
                filename,
                pagesize=A4,
                leftMargin=20*mm,
                rightMargin=20*mm,
                topMargin=30*mm,
                bottomMargin=30*mm
            )
            
            story = []
            
            if include_toc:
                # 목차 제목 추가
                toc_title = Paragraph("목 차", self.styles['KoreanHeading1'])
                story.append(toc_title)
                story.append(Spacer(1, 10*mm))
                story.append(self.toc)
                story.append(PageBreak())
            
            # 본문 내용 추가
            story.extend(self.elements)
            
            # 문서 빌드
            doc.multiBuild(
                story,
                onFirstPage=self._header_footer,
                onLaterPages=self._header_footer,
                canvasmaker=NumberedCanvas
            )
                
        except Exception as e:
            print(f"[ERROR] Failed to save PDF: {str(e)}")
            raise
        

class NumberedCanvas(canvas.Canvas):
    """페이지 번호를 위한 캔버스"""
    def __init__(self, *args, **kwargs):
        canvas.Canvas.__init__(self, *args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        """add page info to each page (page x of y)"""
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            canvas.Canvas.showPage(self)
        canvas.Canvas.save(self)

    def draw_page_number(self, page_count):
        """페이지 번호 그리기"""
        self.setFont("Pretendard", 9)
        self.drawCentredString(
            A4[0]/2,
            10*mm,
            f"- {self._pageNumber} -"
        )

# ============================================================================
# 파일 시스템 기본 유틸리티 (File System Core Utilities)
# ============================================================================

def normalize_path(path: str) -> str:
    """
    Normalize file path by handling spaces, special characters, and user paths.
    파일 경로의 공백, 특수문자, 사용자 경로를 처리합니다.
    """
    try:
        # Expand user path (~/...)
        expanded_path = os.path.expanduser(path)
        
        # Convert to absolute path
        abs_path = os.path.abspath(expanded_path)
        
        # Windows 경로 특수 처리
        if os.name == 'nt':
            # UNC 경로 처리 (네트워크 경로)
            if abs_path.startswith('\\\\'):
                return abs_path
            # 긴 경로 처리 (260자 제한 우회)
            if not abs_path.startswith('\\\\?\\'):
                if len(abs_path) >= 260:
                    abs_path = '\\\\?\\' + abs_path
        
        return abs_path
        
    except Exception as e:
        print("[ERROR] Path normalization failed: %s" % str(e))
        raise

def safe_path_join(*paths: str) -> str:
    """
    Safely join path components.
    안전하게 경로를 결합합니다.
        
    Args:
        *paths: Path components to join
                결합할 경로들
        
    Returns:
        str: Normalized joined path
             정규화된 결합 경로
    """
    try:
        processed_paths = []
        for path in paths:
            path_str = str(path)
            
            # 홈 디렉토리 처리
            if path_str.startswith('~'):
                path_str = os.path.expanduser(path_str)
            
            # Windows 경로 구분자 정규화
            if os.name == 'nt':
                path_str = path_str.replace('/', '\\')
                
            processed_paths.append(path_str)
        
        # 경로 결합 및 정규화
        joined_path = os.path.join(*processed_paths)
        normalized_path = os.path.normpath(joined_path)
        
        # 절대 경로로 변환
        if not os.path.isabs(normalized_path):
            normalized_path = os.path.abspath(normalized_path)
            
        # Windows 긴 경로 처리
        if os.name == 'nt' and len(normalized_path) >= 260:
            if not normalized_path.startswith('\\\\?\\'):
                normalized_path = '\\\\?\\' + normalized_path
            
        return normalized_path
        
    except Exception as e:
        print("[ERROR] Path join failed: %s" % str(e))
        raise

def list_directory(path: str) -> List[str]:
    """
    Return a list of files and folders in the directory.
    디렉토리의 파일/폴더 목록을 반환합니다.
    
    Args:
        path (str): Path to the directory to scan
               탐색할 디렉토리 경로
        
    Returns:
        List[str]: List of file and folder names
                  파일과 폴더 이름 목록
        
    Raises:
        FileNotFoundError: If the directory does not exist
                          디렉토리가 존재하지 않는 경우
        PermissionError: If there is no access permission
                        디렉토리 접근 권한이 없는 경우
    """
    return os.listdir(path)

def read_file(path: str, sheet_name: str = None) -> Union[str, pd.DataFrame, bytes]:
    """
    Read and return the contents of a file based on its extension.
    파일 확장자에 따라 내용을 읽어 반환합니다.

    Args:
        path (str): Path to the file to read
        sheet_name (str, optional): Sheet name for Excel files. Defaults to None.
    """
    file_ext = os.path.splitext(path)[1].lower()
    
    try:
        # Convert to raw path
        raw_path = os.path.expanduser(path)
        
        if not os.path.exists(raw_path):
            print(f"[ERROR] File not found: {raw_path}")
            raise FileNotFoundError(f"File not found: {raw_path}")

        # Office 문서, PDF, HWP 처리
        try:
            if file_ext in ['.doc', '.docx']:
                return extract_from_doc(raw_path)
            elif file_ext in ['.ppt', '.pptx']:
                return extract_from_ppt(raw_path)
            elif file_ext == '.pdf':
                return extract_from_pdf(raw_path)
            elif file_ext in ['.hwp', '.hwpx']:
                return extract_from_hwp(raw_path)
        except Exception as e:
            if "hwp5txt is not installed" in str(e):
                print("[ERROR] hwp5txt is not installed. Please install it using 'pip install --user pyhwp'")
                raise
            elif "Not a valid HWP file or file is corrupted" in str(e):
                print("[ERROR] Invalid or corrupted HWP file: %s" % raw_path)
                raise
            elif "Failed to convert file" in str(e):
                # print("[ERROR] Failed to convert document: %s" % raw_path)
                return None
            else:
                raise

        # Pandas-supported files
        PANDAS_EXTENSIONS = {
            '.xlsx': lambda p: pd.read_excel(p, sheet_name=sheet_name) if sheet_name else pd.read_excel(p),  # Excel files
            '.xls': lambda p: pd.read_excel(p, sheet_name=sheet_name) if sheet_name else pd.read_excel(p),
            '.csv': pd.read_csv,     # CSV files
            '.json': pd.read_json,   # JSON files
            '.html': pd.read_html,   # HTML files
            '.xml': pd.read_xml,     # XML files
            '.parquet': pd.read_parquet,  # Parquet files
            '.feather': pd.read_feather,  # Feather files
            '.pickle': pd.read_pickle,    # Pickle files
            '.sql': pd.read_sql,     # SQL files
            '.hdf': pd.read_hdf,     # HDF5 files
            '.sas': pd.read_sas,     # SAS files
            '.stata': pd.read_stata,  # Stata files
            '.spss': pd.read_spss    # SPSS files
        }
        
        if file_ext in PANDAS_EXTENSIONS:
            try:
                return PANDAS_EXTENSIONS[file_ext](raw_path)
            except Exception as e:
                print(f"[ERROR] Failed to read {file_ext} file: {str(e)}")
                raise
                    
        # Text files
        TEXT_EXTENSIONS = ['.txt', '.log', '.yaml', '.yml', '.md', '.cfg', '.conf']
        if file_ext in TEXT_EXTENSIONS:
            try:
                with open(raw_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception as e:
                print(f"[ERROR] Failed to read text file: {str(e)}")
                raise
                
        # Binary files
        BINARY_EXTENSIONS = [
            # Images
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
            # Audio
            '.mp3', '.wav', '.ogg', '.flac',
            # Video
            '.mp4', '.avi', '.mkv', '.mov',
            # Archives
            '.zip', '.rar', '.7z', '.tar', '.gz'
        ]
        
        if file_ext in BINARY_EXTENSIONS:
            with open(raw_path, 'rb') as f:
                return f.read()
                
        # Unknown files
        try:
            with open(raw_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(raw_path, 'rb') as f:
                return f.read()
                
    except Exception as e:
        print(f"[ERROR] Failed to read file: {str(e)}")
        raise

def write_file(path: str, content: Union[str, pd.DataFrame, bytes], mode: str = 'w', encoding: str = 'utf-8') -> None:
    """
    Write content to a file based on its type and extension.
    내용을 파일 형식에 맞게 저장합니다.
    
    Args:
        path (str): Path to write the file to
        content (Union[str, pd.DataFrame, bytes]): Content to write
        mode (str, optional): File open mode ('w', 'a', 'wb', 'ab'). Defaults to 'w'
        encoding (str, optional): Text encoding. Defaults to 'utf-8'
    """
    try:
        raw_path = os.path.expanduser(path)
        
        directory = os.path.dirname(raw_path)
        if directory:
            os.makedirs(directory, exist_ok=True)
            
        # DataFrame to text for .txt files
        if isinstance(content, pd.DataFrame):
            content = content.to_string()
            
        # Write content based on type
        if isinstance(content, bytes):
            with open(raw_path, 'wb' if 'b' not in mode else mode) as f:
                f.write(content)
        else:
            with open(raw_path, mode, encoding=encoding) as f:
                f.write(str(content))
                
    except Exception as e:
        print(f"[ERROR] Failed to write file: {str(e)}")
        raise

def rename_file_or_directory(old_path: str, new_name: str) -> None:
    """
    Rename a file or directory.
    파일 또는 디렉토리의 이름을 변경합니다.
    
    Args:
        old_path (str): Current path of the file/directory
                       변경할 파일/디렉토리의 현재 경로
        new_name (str): New name for the file/directory
                       새로운 이름
        
    Raises:
        FileNotFoundError: If the file/directory does not exist
                          파일/디렉토리가 존재하지 않는 경우
        PermissionError: If there is no permission to rename
                        변경 권한이 없는 경우
    """
    new_path = os.path.join(os.path.dirname(old_path), new_name)
    os.rename(old_path, new_path)

def remove_file(path: str) -> None:
    """
    Delete a file.
    파일을 삭제합니다.
    
    Args:
        path (str): Path to the file to delete
               삭제할 파일의 경로
        
    Raises:
        FileNotFoundError: If the file does not exist
                          파일이 존재하지 않는 경우
        PermissionError: If there is no permission to delete
                        삭제 권한이 없는 경우
    """
    os.remove(path)

def remove_directory_recursively(path: str) -> None:
    """
    Delete a directory and all its contents recursively.
    디렉토리를 재귀적으로 삭제합니다.
    
    Args:
        path (str): Path to the directory to delete
               삭제할 디렉토리의 경로
        
    Raises:
        FileNotFoundError: If the directory does not exist
                          디렉토리가 존재하지 않는 경우
        PermissionError: If there is no permission to delete
                        삭제 권한이 없는 경우
    """
    shutil.rmtree(path)

# ============================================================================
# 시스템 유틸리티 (System Utilities)
# ============================================================================

def run_command(command: str) -> Tuple[str, str, int]:
    """
    Execute a shell command and return its output.
    셸 명령어를 실행하고 결과를 반환합니다.
    
    Args:
        command (str): Command to execute
                      실행할 명령어
        
    Returns:
        Tuple[str, str, int]: (stdout, stderr, return_code)
                             (표준출력, 표준에러, 반환코드)
        
    Raises:
        subprocess.SubprocessError: If the command execution fails
                                  명령어 실행 실패 시
    """
    process = subprocess.Popen(
        command,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True
    )
    stdout, stderr = process.communicate()
    return stdout, stderr, process.returncode

def which_command(command: str) -> Optional[str]:
    """
    Check if a command exists in system PATH.
    시스템 PATH에 명령어가 존재하는지 확인합니다.
    
    Args:
        command (str): Command to check
                      확인할 명령어
        
    Returns:
        Optional[str]: Full path to the command if found, None otherwise
                      명령어가 존재하면 전체 경로, 없으면 None
    """
    try:
        result = subprocess.run(
            ['which', command],
            capture_output=True,
            text=True
        )
        if result.returncode == 0:
            return result.stdout.strip()
        return None
    except subprocess.SubprocessError:
        return None

def apt_install(package_name: str) -> Tuple[bool, str]:
    """
    Install a package using apt-get.
    apt-get을 사용하여 패키지를 설치합니다.
    
    Args:
        package_name (str): Name of the package to install
                          설치할 패키지 이름
        
    Returns:
        Tuple[bool, str]: (success, message)
                         (성공 여부, 메시지)
        
    Note:
        Requires sudo privileges
        sudo 권한이 필요합니다
    """
    try:
        # Check if running as root
        if os.geteuid() != 0:
            return False, "This function requires root privileges"
            
        cmd = f"apt-get install -y {package_name}"
        result = subprocess.run(
            cmd,
            shell=True,
            capture_output=True,
            text=True
        )
        
        if result.returncode == 0:
            return True, f"Successfully installed {package_name}"
        else:
            return False, f"Failed to install {package_name}: {result.stderr}"
            
    except subprocess.SubprocessError as e:
        return False, f"Installation error: {str(e)}"

def is_package_installed(package_name: str) -> bool:
    """
    Check if a package is installed via apt.
    apt로 패키지가 설치되어 있는지 확인합니다.
    
    Args:
        package_name (str): Name of the package to check
                          확인할 패키지 이름
        
    Returns:
        bool: True if installed, False otherwise
              설치되어 있으면 True, 아니면 False
    """
    try:
        result = subprocess.run(
            ['dpkg', '-l', package_name],
            capture_output=True,
            text=True
        )
        return result.returncode == 0
    except subprocess.SubprocessError:
        return False

# ============================================================================
# 문서 변환 유틸리티 (Document Conversion Utilities)
# ============================================================================

def extract_from_hwp_hwp2html(hwp_path: str) -> str:
    """
    HWP/HWPX 파일에서 텍스트를 추출합니다 (hwp5html 사용).
    """
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        import platform
        import subprocess
        import tempfile
        import os
        from bs4 import BeautifulSoup
        import glob
        import olefile
        
        # 파일 존재 여부 확인
        if not os.path.exists(hwp_path):
            print("[WARNING] File not found: %s" % hwp_path)
            return ""
            
        print("\nExtracting text from: %s" % hwp_path)
        
        # 파일 형식 확인
        is_hwpx = False
        is_hwp = False
        
        try:
            # HWPX 확인 (ZIP 파일 형식)
            try:
                with zipfile.ZipFile(hwp_path) as zf:
                    if 'Contents/section0.xml' in zf.namelist():
                        is_hwpx = True
            except zipfile.BadZipFile:
                pass
                
            # HWP 확인 (OLE2 파일 형식)
            if not is_hwpx:
                try:
                    with olefile.OleFileIO(hwp_path) as ole:
                        if ole.exists('FileHeader'):
                            is_hwp = True
                except:
                    pass
                    
        except Exception as e:
            print("[WARNING] Failed to check file format: %s" % str(e))
            return ""
            
        if not (is_hwpx or is_hwp):
            print("[WARNING] Unsupported file format: %s" % hwp_path)
            return ""
            
        if is_hwpx:
            try:
                with zipfile.ZipFile(hwp_path) as zf:
                    # section0.xml 파싱
                    with zf.open('Contents/section0.xml') as f:
                        content = f.read().decode('utf-8')
                        
                        # XML 파싱
                        root = ET.fromstring(content)
                        
                        # 네임스페이스 정의
                        ns = {
                            'hp': 'http://www.hancom.co.kr/hwpml/2011/paragraph',
                            'hc': 'http://www.hancom.co.kr/hwpml/2011/core',
                            'ha': 'http://www.hancom.co.kr/hwpml/2011/app',
                            'hs': 'http://www.hancom.co.kr/hwpml/2011/section'
                        }
                        
                        text_parts = []
                        
                        # 1. 문단(p) 처리
                        for para in root.findall('.//hp:p', ns):
                            para_text = []
                            
                            # 1.1 일반 텍스트(t)
                            for run in para.findall('.//hp:run', ns):
                                for t in run.findall('.//hp:t', ns):
                                    if t.text:
                                        para_text.append(t.text)
                            
                            # 1.2 표 처리
                            for tbl in para.findall('.//hp:tbl', ns):
                                for tc in tbl.findall('.//hp:tc', ns):
                                    cell_text = []
                                    for t in tc.findall('.//hp:t', ns):
                                        if t.text:
                                            cell_text.append(t.text.strip())
                                    if cell_text:
                                        para_text.append(' '.join(cell_text))
                            
                            if para_text:
                                text_parts.append(' '.join(para_text))
                        
                        text = '\n'.join(text_parts)
                        return text if text.strip() else ""
            except Exception as e:
                print("[WARNING] Failed to parse HWPX file: %s" % str(e))
                return ""
                
        else:  # HWP
            try:
                # 임시 디렉토리 생성
                with tempfile.TemporaryDirectory() as temp_dir:
                    # hwp5html 명령어로 HTML 파일 생성
                    output_dir = os.path.join(temp_dir, 'output')
                    os.makedirs(output_dir, exist_ok=True)
                    cmd = ['hwp5html', '--output', output_dir, hwp_path]
                    
                    try:
                        # UTF-8 인코딩으로 출력 설정
                        env = os.environ.copy()
                        env['PYTHONIOENCODING'] = 'utf-8'
                        
                        # 명령어 실행
                        result = subprocess.run(cmd, 
                                             env=env, 
                                             capture_output=True, 
                                             text=True, 
                                             encoding='utf-8')
                        
                        if result.returncode == 0:
                            # HTML 파일들을 순서대로 읽기
                            text_parts = []
                            html_files = sorted(glob.glob(os.path.join(output_dir, '*.html')))
                            
                            if not html_files:
                                # PrvText 스트림에서 시도
                                try:
                                    with olefile.OleFileIO(hwp_path) as ole:
                                        if ole.exists('PrvText'):
                                            prvtext = ole.openstream('PrvText')
                                            text = prvtext.read().decode('utf-16-le').strip()
                                            text = text.replace('\r\n', '\n')  # 개행문자 통일
                                            text = text.replace('\0', '')      # null 문자 제거
                                            return text if text.strip() else ""
                                except:
                                    pass
                                    
                                print("[WARNING] No text content found in HWP file: %s" % hwp_path)
                                return ""
                            
                            for html_file in html_files:
                                # HTML 파일 읽기
                                with open(html_file, 'r', encoding='utf-8') as f:
                                    html_content = f.read()
                                
                                # BeautifulSoup으로 HTML 파싱
                                soup = BeautifulSoup(html_content, 'html.parser')
                                
                                # 텍스트 추출 및 정제
                                text = soup.get_text(separator='\n', strip=True)
                                if text.strip():
                                    text_parts.append(text)
                            
                            # 모든 텍스트 합치기
                            text = '\n\n'.join(text_parts)
                            text = text.replace('\r\n', '\n')  # 개행문자 통일
                            
                            return text if text.strip() else ""
                        else:
                            error_msg = result.stderr.strip()
                            print("[WARNING] HWP 파일 변환 실패: %s" % error_msg)
                            return ""
                            
                    except subprocess.SubprocessError as e:
                        print("[WARNING] Failed to execute hwp5html: %s" % str(e))
                        return ""
                        
            except Exception as e:
                print("[WARNING] Failed to extract text from HWP file: %s" % str(e))
                return ""
            
    except Exception as e:
        print("[WARNING] Error processing document: %s" % str(e))
        return ""

def extract_from_hwp(hwp_path: str) -> str:
    """
    HWP/HWPX 파일에서 텍스트를 추출합니다.
    
    Args:
        hwp_path: HWP 파일 경로
        
    Returns:
        str: 추출된 텍스트
    """
    try:
        import os
        
        # 1. hwp2txt 방식으로 시도
        # print("\nAttempting to extract text using hwp2txt method...")
        text_content = extract_from_hwp_hwp2txt(hwp_path)
        if text_content:
            text_content = clean_hwp_text(text_content)
            print(f"Extracted text length (hwp2txt): {len(text_content)}")
            sections = extract_structure(text_content)
            return '\n\n'.join(sections)
            
        print("hwp2txt method failed")
        print("Falling back to PDF conversion method")
        
        # 2. PDF로 변환 시도
        # print(f"Attempting to convert HWP to PDF: {hwp_path}")
        pdf_path = convert_hwp_to_pdf(hwp_path)
        
        if pdf_path and os.path.exists(pdf_path):
            print(f"Successfully converted to PDF: {pdf_path}")
            print("Extracting text from PDF...")
            
            # PDF에서 텍스트 추출
            text_content = extract_from_pdf(pdf_path)
            if text_content:
                print(f"Successfully extracted text from PDF, length: {len(text_content)}")
                text_content = clean_hwp_text(text_content)
                print(f"Cleaned text length: {len(text_content)}")
                
                # 문서 구조화
                sections = extract_structure(text_content)
                text_content = '\n\n'.join(sections)
                print(f"Final text length after structuring: {len(text_content)}")
                return text_content
        
        print("Both extraction methods failed")
        return ""
        
    except Exception as e:
        print(f"[ERROR] Failed to process HWP file: {str(e)}")
        return ""

def extract_from_hwp_hwp2txt(hwp_path: str) -> str:
    """
    HWP/HWPX 파일을 텍스트로 변환합니다.
    """
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        import platform
        import subprocess
        import tempfile
        import os
        
        # 파일 확장자 확인
        ext = os.path.splitext(hwp_path)[1].lower()
        
        if ext == '.hwpx':
            try:
                with zipfile.ZipFile(hwp_path) as zf:
                    # section0.xml 파싱
                    with zf.open('Contents/section0.xml') as f:
                        content = f.read().decode('utf-8')
                        
                        # XML 파싱
                        root = ET.fromstring(content)
                        
                        # 네임스페이스 정의
                        ns = {
                            'hp': 'http://www.hancom.co.kr/hwpml/2011/paragraph',
                            'hc': 'http://www.hancom.co.kr/hwpml/2011/core',
                            'ha': 'http://www.hancom.co.kr/hwpml/2011/app',
                            'hs': 'http://www.hancom.co.kr/hwpml/2011/section'
                        }
                        
                        text_parts = []
                        
                        # 1. 문단(p) 처리
                        for para in root.findall('.//hp:p', ns):
                            para_text = []
                            
                            # 1.1 일반 텍스트(t)
                            for run in para.findall('.//hp:run', ns):
                                for t in run.findall('.//hp:t', ns):
                                    if t.text:
                                        para_text.append(t.text)
                            
                            # 1.2 표 처리
                            for tbl in para.findall('.//hp:tbl', ns):
                                for tc in tbl.findall('.//hp:tc', ns):
                                    cell_text = []
                                    for t in tc.findall('.//hp:t', ns):
                                        if t.text:
                                            cell_text.append(t.text.strip())
                                    if cell_text:
                                        para_text.append(' '.join(cell_text))
                            
                            if para_text:
                                text_parts.append(' '.join(para_text))
                        
                        text = '\n'.join(text_parts)
                        text = clean_hwp_text(text)
                        sections = extract_structure(text)
                        return '\n\n'.join(sections) if sections else "No text content found in the HWPX file."
                    
            except Exception as e:
                raise Exception(f"Failed to process HWPX file: {str(e)}")
                
        elif ext == '.hwp':
            # Windows 환경에서는 임시 파일을 사용하여 처리
            if platform.system() == 'Windows':
                # 임시 디렉토리 생성
                with tempfile.TemporaryDirectory() as temp_dir:
                    # hwp5txt 명령어로 텍스트 파일 생성
                    temp_txt = os.path.join(temp_dir, 'output.txt')
                    cmd = ['hwp5txt', '--output', temp_txt, hwp_path]
                    
                    try:
                        # UTF-8 인코딩으로 출력 설정
                        env = os.environ.copy()
                        env['PYTHONIOENCODING'] = 'utf-8'
                        
                        # 명령어 실행
                        result = subprocess.run(cmd, env=env, capture_output=True, text=True, encoding='utf-8')
                        if result.returncode != 0:
                            raise Exception(result.stderr)
                        
                        # 생성된 텍스트 파일 읽기
                        if os.path.exists(temp_txt):
                            with open(temp_txt, 'r', encoding='utf-8') as f:
                                text = f.read()
                                text = clean_hwp_text(text)
                                sections = extract_structure(text)
                                return '\n\n'.join(sections)
                        else:
                            text = result.stdout
                            text = clean_hwp_text(text)
                            sections = extract_structure(text)
                            return '\n\n'.join(sections)
                    except subprocess.SubprocessError as e:
                        raise Exception(f"Failed to convert HWP file: {str(e)}")
            else:
                # Linux/Mac 환경에서는 기존 방식 유지
                result = subprocess.run(
                    ['hwp5txt', hwp_path],
                    capture_output=True,
                    text=True,
                    encoding='utf-8'
                )
                if result.returncode == 0:
                    text = result.stdout
                    text = clean_hwp_text(text)
                    sections = extract_structure(text)
                    return '\n\n'.join(sections)
                else:
                    raise Exception(result.stderr)
        else:
            raise ValueError("Unsupported file format. Only .hwp and .hwpx files are supported.")
            
    except subprocess.SubprocessError as e:
        raise Exception(f"Failed to convert HWP file: {str(e)}")
    except Exception as e:
        # print(f"[ERROR] Failed to convert file: {str(e)}")
        return None

def extract_from_doc(doc_path: str) -> str:
    """
    DOC/DOCX 파일에서 텍스트를 추출합니다.
    
    Args:
        doc_path (str): DOC/DOCX 파일 경로
        
    Returns:
        str: 추출된 텍스트
        
    Raises:
        ImportError: 필요한 패키지가 설치되지 않은 경우
        Exception: 파일 처리 중 오류가 발생한 경우
    """
    try:
        # 파일 존재 여부 확인
        if not os.path.exists(doc_path):
            raise FileNotFoundError(f"File not found: {doc_path}")
            
        # 파일 확장자 확인
        ext = os.path.splitext(doc_path)[1].lower()
        if ext not in ['.doc', '.docx']:
            raise ValueError("File must have .doc or .docx extension")
            
        print(f"\nExtracting text from: {doc_path}")
        
        if ext == '.docx':
            # DOCX 파일 처리
            from docx import Document
            doc = Document(doc_path)
            
            # 텍스트 추출
            paragraphs = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)
                    
            # 표 처리
            for table in doc.tables:
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            row_texts.append(text)
                    if row_texts:
                        paragraphs.append(" | ".join(row_texts))
            
            return "\n".join(paragraphs)
        else:
            # DOC 파일 처리 (antiword 사용)
            import subprocess
            
            # antiword 설치 여부 확인
            result = subprocess.run(['which', 'antiword'], capture_output=True, text=True)
            if result.returncode != 0:
                raise ImportError("antiword is not installed. Please install it using 'sudo apt-get install antiword'")
            
            # antiword로 텍스트 추출
            result = subprocess.run(['antiword', doc_path], capture_output=True, text=True)
            if result.returncode != 0:
                raise Exception(f"Failed to extract text using antiword: {result.stderr}")
            
            return result.stdout
        
    except ImportError as e:
        if 'antiword' in str(e):
            raise ImportError(str(e))
        raise ImportError("Required packages not found. Please install 'python-docx' for DOCX files and 'antiword' for DOC files")
    except Exception as e:
        raise Exception(f"Failed to extract text from document: {str(e)}")

def extract_from_ppt(ppt_path: str) -> str:
    """
    PPT/PPTX 파일에서 텍스트를 추출합니다.
    
    Args:
        ppt_path (str): PPT/PPTX 파일 경로
        
    Returns:
        str: 추출된 텍스트
        
    Raises:
        ImportError: 필요한 패키지가 설치되지 않은 경우
        Exception: 파일 처리 중 오류가 발생한 경우
    """
    try:
        # 파일 존재 여부 확인
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"File not found: {ppt_path}")
            
        # 파일 확장자 확인
        ext = os.path.splitext(ppt_path)[1].lower()
        if ext not in ['.ppt', '.pptx']:
            raise ValueError("File must have .ppt or .pptx extension")
            
        print(f"\nExtracting text from: {ppt_path}")
        
        if ext == '.pptx':
            # PPTX 파일 처리
            from pptx import Presentation
            prs = Presentation(ppt_path)
            
            # 슬라이드별 텍스트 추출
            all_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                print(f"Processing slide {i}/{len(prs.slides)}...")
                
                # 도형에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                            
                    # 표 처리
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                text = cell.text.strip()
                                if text:
                                    row_text.append(text)
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.extend(table_text)
                
                if slide_text:
                    all_text.append(f"[Slide {i}]\n" + "\n".join(slide_text))
            
            return "\n\n".join(all_text)
        else:
            # PPT 파일 처리 (catppt 사용)
            import subprocess
            
            # catppt 설치 여부 확인
            result = subprocess.run(['which', 'catppt'], capture_output=True, text=True)
            if result.returncode != 0:
                raise ImportError("catppt is not installed. Please install it using 'sudo apt-get install catdoc'")
            
            # catppt로 텍스트 추출
            result = subprocess.run(['catppt', ppt_path], capture_output=True, text=True)
            if result.returncode != 0:
                raise Exception(f"Failed to extract text using catppt: {result.stderr}")
            
            return result.stdout
        
    except ImportError as e:
        if 'catppt' in str(e):
            raise ImportError(str(e))
        raise ImportError("Required packages not found. Please install 'python-pptx' for PPTX files and 'catdoc' for PPT files")
    except Exception as e:
        raise Exception(f"Failed to extract text from presentation: {str(e)}")

def extract_from_pdf(pdf_path: str) -> str:
    """
    PDF 파일에서 텍스트를 추출합니다.
    
    Args:
        pdf_path (str): PDF 파일 경로
        
    Returns:
        str: 추출된 텍스트
        
    Raises:
        ImportError: PyPDF2 패키지가 설치되지 않은 경우
        Exception: 파일 처리 중 오류가 발생한 경우
    """
    try:
        from PyPDF2 import PdfReader
        import re
        
        # 파일 존재 여부 확인
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"File not found: {pdf_path}")
            
        # 파일 확장자 확인
        ext = os.path.splitext(pdf_path)[1].lower()
        if ext != '.pdf':
            raise ValueError("File must have .pdf extension")
            
        print(f"\nExtracting text from PDF: {pdf_path}")
        
        # PDF 파일 열기
        reader = PdfReader(pdf_path)
        
        # 페이지별 텍스트 추출
        all_text = []
        total_pages = len(reader.pages)
        
        # 텍스트 정제를 위한 패턴들
        cleanup_patterns = [
            # 목차 관련 패턴
            (r'\*{2,}', ''),  # 연속된 별표(*) 제거
            (r'\.{2,}', ''),  # 연속된 점(...) 제거
            (r'_{2,}', ''),   # 연속된 밑줄(_) 제거
            (r'={2,}', ''),   # 연속된 등호(=) 제거
            (r'-{2,}', '-'),  # 연속된 하이픈(-)을 하나로
            (r'~{2,}', ''),   # 연속된 물결표(~) 제거
            
            # 제어 문자 및 특수 문자 패턴
            (r'[\u0000-\u0008\u000B\u000C\u000E-\u001F]', ''),  # 제어 문자 제거
            (r'[\u2000-\u200F\u2028-\u202F]', ' '),  # 특수 공백 및 제어 문자
            (r'[\uFFF0-\uFFFF]', ''),  # 특수 유니코드 영역
            (r'[\u0080-\u00FF]', ''),  # 확장 ASCII 영역
            (r'[\u0100-\u017F]', ''),  # 라틴 확장 문자
            (r'[\u0180-\u024F]', ''),  # 라틴 확장 추가
            (r'[\u0250-\u02AF]', ''),  # IPA 확장
            (r'[\u0300-\u036F]', ''),  # 결합 발음 구별 부호
            (r'[\u2500-\u257F]', ''),  # 박스 드로잉
            (r'[\u2580-\u259F]', ''),  # 블록 요소
            (r'[\u3040-\u309F]', ''),  # 히라가나
            (r'[\u30A0-\u30FF]', ''),  # 가타카나
            (r'[\u31F0-\u31FF]', ''),  # 가타카나 음성 확장
            (r'[\uFF00-\uFFEF]', ''),  # 전각 문자
            (r'[\u3000-\u303F]', ' '),  # CJK 기호 및 문장 부호
            
            # 한글과 기본 문자 처리
            (r'[\uAC00-\uD7AF가-힣]', lambda m: m.group()),  # 한글은 유지
            (r'[^\uAC00-\uD7AF가-힣\w\s.,!?():\-\[\]\/]', ' '),  # 나머지 문자는 공백으로
            
            # 공백 정리
            (r'\s+', ' '),  # 연속된 공백을 하나로
        ]
        
        for i, page in enumerate(reader.pages, 1):
            print(f"Processing page {i}/{total_pages}...")
            text = page.extract_text()
            
            # 텍스트 정제
            if text.strip():
                # 패턴 적용
                for pattern, replacement in cleanup_patterns:
                    if callable(replacement):
                        text = re.sub(pattern, replacement, text)
                    else:
                        text = re.sub(pattern, replacement, text)
                
                # 줄바꿈 처리
                text = re.sub(r'\n\s*\n', '\n\n', text)  # 빈 줄 정리
                text = re.sub(r'([^.])\n([^\n])', r'\1 \2', text)  # 문장 중간 줄바꿈 처리
                
                # 목차 번호와 내용 사이의 공백 정리
                text = re.sub(r'(\d+[\.\)］\]])\s*', r'\1 ', text)  # 목차 번호 뒤 공백 정리
                
                # 최종 정제
                text = text.strip()
                if text:
                    all_text.append(f"[Page {i}]\n{text}")
        
        if not all_text:
            print("Warning: No text content found in PDF")
            return ""
            
        return "\n\n".join(all_text)
        
    except ImportError:
        raise ImportError("Required package 'PyPDF2' not found")
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")

def extract_tables_from_hwp(hwp_path: str, excel_path: str) -> bool:
    """
    HWP/HWPX 파일에서 표를 추출하여 Excel 파일로 저장합니다.
    """
    try:
        import pandas as pd
        import subprocess
        import os
        from bs4 import BeautifulSoup
        import tempfile        
        
        # 파일 존재 여부 확인
        if not os.path.exists(hwp_path):
            raise FileNotFoundError(f"File not found: {hwp_path}")
            
        # 파일 확장자 확인
        ext = os.path.splitext(hwp_path)[1].lower()
        if ext not in ['.hwp', '.hwpx']:
            raise ValueError("File must have .hwp or .hwpx extension")
        
        print(f"\nExtracting tables from: {hwp_path}")
        tables = []
        
        if ext == '.hwpx':
            # HWPX 파일 처리
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(hwp_path) as zf:
                # section0.xml 파싱
                with zf.open('Contents/section0.xml') as f:
                    content = f.read().decode('utf-8')
                    root = ET.fromstring(content)
                    
                    # 네임스페이스 정의
                    ns = {
                        'hp': 'http://www.hancom.co.kr/hwpml/2011/paragraph',
                        'hc': 'http://www.hancom.co.kr/hwpml/2011/core',
                        'ha': 'http://www.hancom.co.kr/hwpml/2011/app',
                        'hs': 'http://www.hancom.co.kr/hwpml/2011/section'
                    }
                    
                    # 표 추출
                    table_count = 0
                    for tbl in root.findall('.//hp:tbl', ns):
                        table_count += 1
                        print(f"\nProcessing table {table_count}...")
                        
                        table_data = []
                        max_cols = 0
                        
                        # 모든 행을 순회하며 최대 열 수 찾기
                        for tr in tbl.findall('.//hp:tr', ns):
                            cols = len(tr.findall('.//hp:tc', ns))
                            max_cols = max(max_cols, cols)
                        
                        # 행 처리
                        row_count = 0
                        for tr in tbl.findall('.//hp:tr', ns):
                            row_count += 1
                            row_data = [''] * max_cols  # 빈 셀로 초기화
                            
                            # 셀 처리
                            for i, tc in enumerate(tr.findall('.//hp:tc', ns)):
                                cell_text = []
                                
                                # 병합 셀 정보
                                rowspan = int(tc.get('rowspan', '1'))
                                colspan = int(tc.get('colspan', '1'))
                                
                                # 셀 내용 추출
                                for t in tc.findall('.//hp:t', ns):
                                    if t.text:
                                        cell_text.append(t.text.strip())
                                
                                # 셀 내용 저장
                                cell_content = ' '.join(cell_text) if cell_text else ''
                                
                                # 병합 셀 처리
                                for col in range(i, min(i + colspan, max_cols)):
                                    row_data[col] = cell_content
                            
                            table_data.append(row_data)
                        
                        if table_data:
                            # 빈 행/열 제거
                            df = pd.DataFrame(table_data)
                            # 모든 값이 빈 문자열인 행/열 제거
                            df = df.loc[:, (df != '').any()]
                            df = df.loc[(df != '').any(axis=1)]
                            
                            if not df.empty:
                                tables.append(df)
                                print(f"Found table with {len(df)} rows and {len(df.columns)} columns")
                            else:
                                print("Table is empty after cleaning")
                        else:
                            print("No data found in table")
                            
        else:
            # HWP 파일 처리
            
            # 파일 존재 여부 확인
            if not os.path.exists(hwp_path):
                raise FileNotFoundError(f"File not found: {hwp_path}")
                
            # 파일 확장자 확인
            ext = os.path.splitext(hwp_path)[1].lower()
            if ext not in ['.hwp', '.hwpx']:
                raise ValueError("File must have .hwp or .hwpx extension")
            
            print(f"\nExtracting tables from: {hwp_path}")
            tables = []
            
            # 임시 디렉토리 생성
            temp_dir = tempfile.mkdtemp(prefix='hwp_')
            temp_html_path = os.path.join(temp_dir, 'output.html')
            
            try:
                # HWP 파일을 HTML로 변환
                result = subprocess.run(['hwp5html', '--output', temp_dir, hwp_path],
                                    capture_output=True, text=True)
                if result.returncode != 0:
                    raise Exception(f"Failed to convert HWP to HTML: {result.stderr}")
                
                # index.xhtml 파일 찾기
                index_path = os.path.join(temp_dir, 'index.xhtml')
                if not os.path.exists(index_path):
                    raise Exception("HTML conversion failed: index.xhtml not found")
                
                # HTML 파일 읽기
                with open(index_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # BeautifulSoup으로 HTML 파싱
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # 표 추출
                table_elements = soup.find_all('table')
                print(f"Found {len(table_elements)} tables")
                
                for i, table in enumerate(table_elements, 1):
                    print(f"\nProcessing table {i}...")
                    
                    # 표 데이터 추출
                    table_data = []
                    for row in table.find_all('tr'):
                        row_data = []
                        for cell in row.find_all(['td', 'th']):
                            # 셀 병합 정보 처리
                            rowspan = int(cell.get('rowspan', 1))
                            colspan = int(cell.get('colspan', 1))
                            
                            # 셀 내용 가져오기
                            cell_text = cell.get_text(strip=True)
                            
                            # 병합된 셀 처리
                            for _ in range(colspan):
                                row_data.append(cell_text)
                        
                        if row_data:  # 빈 행 제외
                            table_data.append(row_data)
                    
                    if table_data:
                        # DataFrame으로 변환
                        df = pd.DataFrame(table_data)
                        
                        # 빈 행/열 제거
                        df = df.loc[:, (df != '').any()]
                        df = df.loc[(df != '').any(axis=1)]
                        
                        if not df.empty:
                            tables.append(df)
                            print(f"Added table with {len(df)} rows and {len(df.columns)} columns")
                        else:
                            print("Table is empty after cleaning")
                    else:
                        print("No data found in table")
                
            finally:
                # 임시 디렉토리 삭제
                import shutil
                shutil.rmtree(temp_dir, ignore_errors=True)
        
        # 추출된 표가 있는 경우
        if tables:
            # Excel 파일 저장 디렉토리 생성
            os.makedirs(os.path.dirname(excel_path), exist_ok=True)
            
            print(f"\nSaving {len(tables)} tables to Excel file: {excel_path}")
            
            # 여러 시트로 Excel 파일 저장
            with pd.ExcelWriter(excel_path) as writer:
                for i, df in enumerate(tables):
                    sheet_name = f'Table_{i+1}'
                    df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                    print(f"Saved table {i+1} to sheet '{sheet_name}'")
            
            return True
        
        print("No valid tables found in the document")
        return False
        
    except ImportError:
        raise ImportError("Required packages not found. Please install 'pyhwp' and 'olefile'")
    except (FileNotFoundError, ValueError) as e:
        raise type(e)(str(e))
    except Exception as e:
        raise Exception(f"Failed to extract tables: {str(e)}")

def extract_images_from_hwp(hwp_path: str, output_dir: str) -> List[str]:
    """
    HWP/HWPX 파일에서 이미지를 추출하여 지정된 디렉토리에 저장합니다.
    """
    try:
        import os
        import hashlib
        import zlib
        import zipfile
        from hwp5.filestructure import Hwp5File
        from hwp5.storage.ole import OleStorage
        import olefile
        
        # 파일 존재 여부 확인
        if not os.path.exists(hwp_path):
            raise FileNotFoundError(f"File not found: {hwp_path}")
            
        # 파일 확장자 확인
        ext = os.path.splitext(hwp_path)[1].lower()
        if ext not in ['.hwp', '.hwpx']:
            raise ValueError("File must have .hwp or .hwpx extension")
        
        # 출력 디렉토리 생성
        os.makedirs(output_dir, exist_ok=True)
        
        print(f"\nExtracting images from: {hwp_path}")
        print(f"Output directory: {output_dir}")
        
        # 저장된 이미지 파일 경로 리스트
        saved_images = []
        
        if ext == '.hwpx':
            # HWPX 파일 처리
            with zipfile.ZipFile(hwp_path) as zf:
                # BinData 디렉토리의 파일 목록 가져오기
                bindata_files = [f for f in zf.namelist() if f.startswith('BinData/')]
                print(f"\nFound {len(bindata_files)} files in BinData directory")
                
                for bin_file in bindata_files:
                    print(f"\nProcessing: {bin_file}")
                    
                    # 파일 데이터 읽기
                    data = zf.read(bin_file)
                    
                    # 이미지 파일 확장자 확인
                    ext = None
                    if data.startswith(b'\xFF\xD8'):  # JPEG
                        ext = '.jpg'
                    elif data.startswith(b'\x89PNG'):  # PNG
                        ext = '.png'
                    elif data.startswith(b'GIF8'):  # GIF
                        ext = '.gif'
                    elif data.startswith(b'BM'):  # BMP
                        ext = '.bmp'
                    elif data.startswith(b'\x00\x00\x01\x00'):  # ICO
                        ext = '.ico'
                    elif data.startswith(b'\x00\x00\x02\x00'):  # CUR
                        ext = '.cur'
                    elif data.startswith(b'%PDF'):  # PDF
                        ext = '.pdf'
                    elif data[0:4] in [b'RIFF', b'WEBP']:  # WEBP
                        ext = '.webp'
                    
                    if ext:
                        print(f"Detected image type: {ext}")
                        # 파일명 생성 (중복 방지를 위해 해시 사용)
                        hash_name = hashlib.md5(data).hexdigest()[:8]
                        image_path = os.path.join(output_dir, f'image_{hash_name}{ext}')
                        
                        # 이미지 파일 저장
                        with open(image_path, 'wb') as f:
                            f.write(data)
                        saved_images.append(image_path)
                        print(f"Saved image to: {image_path}")
                    else:
                        print(f"Not an image file (magic bytes: {data[:8].hex()})")

        else:
            # HWP 파일 처리
            ole = olefile.OleFileIO(hwp_path)
            
            try:
                print("\nScanning for BinData...")
                
                # BinData 스토리지에서 파일 목록 가져오기
                bindata_list = []
                for entry in ole.listdir():
                    if 'BinData' in entry:
                        bindata_list.append(entry)
                
                print(f"Found {len(bindata_list)} BinData items")
                
                # 각 BinData 처리
                for entry in bindata_list:
                    print(f"\nProcessing: {'/'.join(entry)}")
                    
                    # 스트림 데이터 읽기
                    stream = ole.openstream(entry)
                    data = stream.read()
                    stream.close()
                    
                    # 압축 해제 시도
                    try:
                        decompressed = zlib.decompress(data, -15)  # 압축 해제 시도
                        print("Successfully decompressed data")
                        data = decompressed
                    except zlib.error:
                        print("Data is not compressed with zlib")
                    
                    # 이미지 파일 확장자 확인
                    ext = None
                    if data.startswith(b'\xFF\xD8'):  # JPEG
                        ext = '.jpg'
                    elif data.startswith(b'\x89PNG'):  # PNG
                        ext = '.png'
                    elif data.startswith(b'GIF8'):  # GIF
                        ext = '.gif'
                    elif data.startswith(b'BM'):  # BMP
                        ext = '.bmp'
                    elif data.startswith(b'\x00\x00\x01\x00'):  # ICO
                        ext = '.ico'
                    elif data.startswith(b'\x00\x00\x02\x00'):  # CUR
                        ext = '.cur'
                    elif data.startswith(b'%PDF'):  # PDF
                        ext = '.pdf'
                    elif data[0:4] in [b'RIFF', b'WEBP']:  # WEBP
                        ext = '.webp'
                    
                    # 압축 해제된 데이터의 매직 바이트 출력
                    print(f"Magic bytes after processing: {data[:8].hex()}")
                    
                    if ext:
                        print(f"Detected image type: {ext}")
                        # 파일명 생성 (중복 방지를 위해 해시 사용)
                        hash_name = hashlib.md5(data).hexdigest()[:8]
                        image_path = os.path.join(output_dir, f'image_{hash_name}{ext}')
                        
                        # 이미지 파일 저장
                        with open(image_path, 'wb') as f:
                            f.write(data)
                        saved_images.append(image_path)
                        print(f"Saved image to: {image_path}")
                    else:
                        print(f"Not an image file (magic bytes: {data[:8].hex()})")
                
                if not saved_images:
                    print("\nNo valid images found in BinData storage")
                    
            finally:
                ole.close()
        
        if saved_images:
            print(f"\nSuccessfully extracted {len(saved_images)} images")
            return saved_images
        else:
            print("\nNo images found in the document")
            return []
        
    except ImportError:
        raise ImportError("Required packages not found. Please install 'pyhwp' and 'olefile'")
    except (FileNotFoundError, ValueError) as e:
        raise type(e)(str(e))
    except Exception as e:
        raise Exception(f"Failed to extract images: {str(e)}")

def create_pdf_document() -> PDFDocument:
    """
    Create a new PDF document with Korean support.
    한글을 지원하는 새로운 PDF 문서를 생성합니다.
    
    Returns:
        PDFDocument: A new PDF document instance
                    새로운 PDF 문서 인스턴스
                    
    Raises:
        FileNotFoundError: If Korean font file is not found
                          한글 폰트 파일이 없는 경우
    """
    return PDFDocument()

# ============================================================================
# 시각화 유틸리티 (Visualization Utilities)
# ============================================================================

def create_matplotlib(figsize: tuple = (6, 4)) -> tuple[plt.Figure, plt.Axes, fm.FontProperties]:
    """
    한글 폰트가 설정된 matplotlib 그래프를 생성합니다.
    
    Args:
        figsize (tuple): 그래프 크기 (width, height)
        
    Returns:
        tuple[plt.Figure, plt.Axes, fm.FontProperties]: (figure, axis, font_properties)
        
    Raises:
        FileNotFoundError: 한글 폰트 파일이 없는 경우
    """
    try:
        # 한글 폰트 설정
        font_path = safe_path_join(os.path.expanduser("~"), ".airun", "Pretendard-Regular.ttf")
        if not os.path.exists(font_path):
            raise FileNotFoundError("Korean font file not found. Please ensure airun is properly installed.")
        
        # 폰트 속성 설정
        font_prop = fm.FontProperties(fname=font_path)
        plt.rcParams["font.family"] = font_prop.get_name()
        fm.fontManager.addfont(font_path)
        plt.rcParams["axes.unicode_minus"] = False
        plt.rcParams["font.size"] = 12
        
        # 그래프 생성
        fig, ax = plt.subplots(figsize=figsize)
        
        # 안전한 텍스트 처리를 위한 메서드 오버라이드
        original_set_title = ax.set_title
        original_set_xlabel = ax.set_xlabel
        original_set_ylabel = ax.set_ylabel
        original_legend = ax.legend
        
        def safe_set_title(title, **kwargs):
            if 'fontproperties' not in kwargs:
                kwargs['fontproperties'] = font_prop
            return original_set_title(str(title), **kwargs)
            
        def safe_set_xlabel(label, **kwargs):
            if 'fontproperties' not in kwargs:
                kwargs['fontproperties'] = font_prop
            return original_set_xlabel(str(label), **kwargs)
            
        def safe_set_ylabel(label, **kwargs):
            if 'fontproperties' not in kwargs:
                kwargs['fontproperties'] = font_prop
            return original_set_ylabel(str(label), **kwargs)
            
        def safe_legend(*args, **kwargs):
            if 'prop' not in kwargs:
                kwargs['prop'] = font_prop
            if args and isinstance(args[0], (list, tuple)):
                args = list(args)
                args[0] = [str(label) for label in args[0]]
            return original_legend(*args, **kwargs)
        
        # 안전한 메서드로 교체
        ax.set_title = safe_set_title
        ax.set_xlabel = safe_set_xlabel
        ax.set_ylabel = safe_set_ylabel
        ax.legend = safe_legend
        
        # 숫자 포맷팅 헬퍼 함수 추가
        def add_formatter(formatter_func):
            """숫자 포맷터를 y축에 추가"""
            import matplotlib.ticker as ticker
            ax.yaxis.set_major_formatter(ticker.FuncFormatter(formatter_func))
        
        # 축에 헬퍼 함수 추가
        ax.add_formatter = add_formatter
        
        return fig, ax, font_prop
        
    except Exception as e:
        print(f"[ERROR] Failed to create matplotlib figure: {str(e)}")
        raise

def save_plot(fig: plt.Figure, filename: str, dpi: int = 300) -> None:
    """
    matplotlib 그래프를 파일로 저장합니다.
    
    Args:
        fig (plt.Figure): matplotlib figure 객체
        filename (str): 저장할 파일 경로
        dpi (int): 이미지 해상도
    """
    fig.savefig(filename, dpi=dpi, bbox_inches="tight")
    plt.close(fig)

def convert_dot_to_svg(dot_path: str, output_path: str = None) -> str:
    """
    DOT 파일을 SVG로 변환합니다.
    
    Args:
        dot_path (str): DOT 파일 경로
        output_path (str, optional): 출력할 SVG 파일 경로. 지정하지 않으면 DOT 파일과 같은 위치에 생성
        
    Returns:
        str: 생성된 SVG 파일 경로
        
    Raises:
        ImportError: graphviz 패키지가 설치되지 않은 경우
        FileNotFoundError: DOT 파일이 존재하지 않는 경우
        Exception: 변환 중 오류가 발생한 경우
    """
    try:
        # graphviz 패키지 설치 확인 및 설치
        install_if_missing('graphviz')
        
        import os
        import graphviz
        
        # DOT 파일 존재 여부 확인
        if not os.path.exists(dot_path):
            raise FileNotFoundError(f"DOT file not found: {dot_path}")
            
        # 출력 경로가 지정되지 않은 경우 기본값 설정
        if output_path is None:
            output_path = os.path.splitext(dot_path)[0] + '.svg'
            
        # DOT 파일 읽기
        with open(dot_path, 'r', encoding='utf-8') as f:
            dot_content = f.read()
            
        # DOT 내용을 Source 객체로 변환
        src = graphviz.Source(dot_content)
        
        # SVG로 렌더링
        # render()는 확장자를 자동으로 추가하므로, .svg 확장자 제거
        output_path_without_ext = os.path.splitext(output_path)[0]
        rendered_path = src.render(filename=output_path_without_ext, format='svg', cleanup=True)
        
        print(f"Successfully converted {dot_path} to {rendered_path}")
        return rendered_path
        
    except ImportError:
        raise ImportError("Required package 'graphviz' not found")
    except Exception as e:
        raise Exception(f"Failed to convert DOT to SVG: {str(e)}")

# ============================================================================
# 웹 관련 유틸리티 (Web Utilities)
# ============================================================================

def read_url(url: str) -> Union[str, bytes]:
    """
    Read and return the contents of a URL.
    URL의 내용을 읽어 반환합니다.
    
    Args:
        url (str): URL to read from
              읽을 URL
        
    Returns:
        Union[str, bytes]: Contents of the URL. Returns bytes for binary content (images, etc)
                          URL의 내용. 바이너리 콘텐츠(이미지 등)의 경우 bytes 반환
        
    Raises:
        requests.RequestException: If the URL request fails
                                 URL 요청 실패 시
    """
    response = requests.get(url, verify=False)
    response.raise_for_status()
    
    # Check content type
    content_type = response.headers.get('content-type', '').lower()
    if any(t in content_type for t in ['image/', 'video/', 'audio/', 'application/octet-stream']):
        return response.content
    return response.text

def get_selenium_driver():
    """셀레니움 웹드라이버를 생성합니다."""
    try:
        import platform
        import os
        import sys
        
        # 운영체제 확인
        is_windows = platform.system().lower() == 'windows'
        
        if is_windows:
            # Windows용 설정
            chrome_options = Options()
            chrome_options.add_argument('--headless=new')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument('--disable-gpu')
            chrome_options.add_argument('--window-size=1920,1080')
            chrome_options.add_argument('--disable-extensions')
            chrome_options.add_argument('--disable-software-rasterizer')
            chrome_options.add_argument('--ignore-certificate-errors')
            chrome_options.add_argument('--log-level=3')
            chrome_options.add_argument('--silent')
            chrome_options.add_argument('--disable-logging')
            chrome_options.add_argument('--disable-background-networking')
            chrome_options.add_argument('--disable-background-timer-throttling')
            chrome_options.add_argument('--disable-backgrounding-occluded-windows')
            chrome_options.add_argument('--disable-breakpad')
            chrome_options.add_argument('--disable-client-side-phishing-detection')
            chrome_options.add_argument('--disable-default-apps')
            chrome_options.add_argument('--disable-features=site-per-process')
            chrome_options.add_argument('--disable-hang-monitor')
            chrome_options.add_argument('--disable-popup-blocking')
            chrome_options.add_argument('--disable-prompt-on-repost')
            chrome_options.add_argument('--disable-sync')
            chrome_options.add_argument('--metrics-recording-only')
            chrome_options.add_argument('--no-first-run')
            chrome_options.add_argument('--safebrowsing-disable-auto-update')
            chrome_options.add_argument('--password-store=basic')
            
            # Chrome 실행 파일 경로 직접 지정
            chrome_path = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
            if os.path.exists(chrome_path):
                chrome_options.binary_location = chrome_path
            
            driver = webdriver.Chrome(options=chrome_options)
            
        else:
            # Linux용 설정
            global _chrome_driver_path
            
            chrome_options = Options()
            chrome_options.add_argument('--headless=new')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument('--disable-gpu')
            chrome_options.add_argument('--window-size=1920x1080')
            chrome_options.add_argument('--disable-blink-features=AutomationControlled')
            chrome_options.add_argument('--lang=ko_KR')
            chrome_options.add_argument('--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36')
            chrome_options.add_experimental_option('excludeSwitches', ['enable-automation'])
            chrome_options.add_experimental_option('useAutomationExtension', False)
            
        try:
            # 먼저 시스템에 설치된 ChromeDriver 사용 시도
            driver = webdriver.Chrome(options=chrome_options)
        except:
            try:
                # 실패하면 webdriver_manager 사용
                _chrome_driver_path = ChromeDriverManager().install()
                service = Service(_chrome_driver_path)
                driver = webdriver.Chrome(service=service, options=chrome_options)
            except Exception as e:
                print(f"[ERROR] ChromeDriver 초기화 실패: {str(e)}", file=sys.stderr)
                return None
        
        driver.set_page_load_timeout(30)
        return driver
        
    except Exception as e:
        print(f"[ERROR] Selenium WebDriver 초기화 실패: {str(e)}", file=sys.stderr)
        return None

def extract_web_content(url: str, extract_type: str = 'all', max_items: int = 10, session: requests.Session = None) -> dict:
    """
    웹사이트의 주요 콘텐츠를 추출합니다.
    
    Args:
        url (str): 콘텐츠를 추출할 웹사이트의 URL
        extract_type (str): 추출할 콘텐츠 타입 ('all', 'text', 'links', 'media')
        max_items (int): 추출할 최대 아이템 수
        session (requests.Session): 기존 세션 사용 (선택사항)
        
    Returns:
        dict: {
            'title': 페이지 제목,
            'content': 주요 콘텐츠 텍스트,
            'links': [{'url': 링크URL, 'text': 링크텍스트, 'type': 링크타입}, ...],
            'media': [{'url': 미디어URL, 'type': 미디어타입, 'title': 미디어제목}, ...],
            'metadata': {'description': 설명, 'keywords': [키워드들], 'author': 작성자, 'published_date': 발행일}
        }
    """
    try:
        # 필요한 패키지 설치
        install_if_missing('trafilatura')
        install_if_missing('requests')
        
        import trafilatura
        import requests
        from urllib.parse import urlparse
        from bs4 import BeautifulSoup
        import chardet
        import urllib3
        
        # print(f"\n[DEBUG] extract_web_content 시작")
        # print(f"[DEBUG] URL: {url}")
        # print(f"[DEBUG] extract_type: {extract_type}")
        # print(f"[DEBUG] max_items: {max_items}")
        
        # SSL 경고 메시지 비활성화
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        
        # URL 유효성 검사
        parsed_url = urlparse(url)
        if not all([parsed_url.scheme, parsed_url.netloc]):
            raise ValueError("Invalid URL format")
            
        # 세션이 없으면 새로 생성
        if session is None:
            session = requests.Session()
            session.headers.update({
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7'
            })
        
        # print("[DEBUG] 웹페이지 다운로드 시작")
        # 웹페이지 다운로드 (SSL 검증 비활성화)
        response = session.get(url, timeout=10, verify=False)
        
        # 인코딩 자동 감지
        if response.encoding.lower() == 'iso-8859-1':
            encoding_detect = chardet.detect(response.content)
            detected_encoding = encoding_detect['encoding']
            if detected_encoding:
                response.encoding = detected_encoding
                # print(f"[DEBUG] 감지된 인코딩: {detected_encoding}")
        
        downloaded = response.text
        # print("[DEBUG] 웹페이지 다운로드 완료")
        
        if not downloaded:
            raise Exception("Failed to download webpage")
            
        # 결과 저장할 딕셔너리
        result = {
            'title': '',
            'content': '',
            'links': [],
            'media': [],
            'metadata': {}
        }
        
        # BeautifulSoup으로 파싱
        # print("[DEBUG] BeautifulSoup 파싱 시작")
        soup = BeautifulSoup(downloaded, 'lxml')
        
        # 제목 추출
        result['title'] = soup.title.string.strip() if soup.title else ''
        # print(f"[DEBUG] 추출된 제목: {result['title']}")
        
        # 메타데이터 추출
        meta_tags = soup.find_all('meta')
        metadata = {}
        for tag in meta_tags:
            if 'name' in tag.attrs and 'content' in tag.attrs:
                metadata[tag['name']] = tag['content']
        
        result['metadata'] = {
            'description': metadata.get('description', ''),
            'keywords': metadata.get('keywords', '').split(',') if metadata.get('keywords') else [],
            'author': metadata.get('author', ''),
            'published_date': metadata.get('published_date', '')
        }
        # print(f"[DEBUG] 추출된 메타데이터: {result['metadata']}")
        
        # 주요 콘텐츠 추출
        if extract_type in ['all', 'text']:
            # print("[DEBUG] 주요 콘텐츠 추출 시작")
            content = trafilatura.extract(downloaded, include_links=False, include_images=False, output_format='txt')
            result['content'] = content if content else "콘텐츠를 추출할 수 없습니다."
            # print(f"[DEBUG] 콘텐츠 추출 길이: {len(result['content'])}")
            
        # 링크 추출
        if extract_type in ['all', 'links']:
            # print("[DEBUG] 링크 추출 시작")
            links = []
            for link in soup.find_all('a', href=True)[:max_items]:
                href = link['href']
                # 상대 URL을 절대 URL로 변환
                if not href.startswith(('http://', 'https://')):
                    href = requests.compat.urljoin(url, href)
                links.append({
                    'url': href,
                    'text': link.get_text(strip=True),
                    'type': 'video' if 'youtube.com/watch' in href else 'link'
                })
            result['links'] = links
            # print(f"[DEBUG] 추출된 링크 수: {len(result['links'])}")
            
        # 미디어 추출
        if extract_type in ['all', 'media']:
            # print("[DEBUG] 미디어 추출 시작")
            media = []
            # 이미지 추출
            for img in soup.find_all('img', src=True)[:max_items]:
                src = img['src']
                if not src.startswith(('http://', 'https://')):
                    src = requests.compat.urljoin(url, src)
                media.append({
                    'url': src,
                    'type': 'image',
                    'title': img.get('alt', '')
                })
                # print(f"[DEBUG] 추출된 이미지 URL: {src}")
            # 비디오 추출
            for video in soup.find_all('video', src=True)[:max_items]:
                src = video['src']
                if not src.startswith(('http://', 'https://')):
                    src = requests.compat.urljoin(url, src)
                media.append({
                    'url': src,
                    'type': 'video',
                    'title': video.get('title', '')
                })
                # print(f"[DEBUG] 추출된 비디오 URL: {src}")
            result['media'] = media
            # print(f"[DEBUG] 추출된 미디어 수: {len(result['media'])}")
            
        # print("[DEBUG] extract_web_content 완료")
        return result
        
    except ImportError as e:
        raise ImportError(f"Required package not found: {str(e)}")
    except requests.RequestException as e:
        raise Exception(f"Failed to download webpage: {str(e)}")
    except ValueError as e:
        raise ValueError(f"Invalid parameter: {str(e)}")
    except Exception as e:
        raise Exception(f"Failed to extract web content: {str(e)}")

def open_in_browser(url: str) -> bool:
    """
    URL을 기본 브라우저에서 엽니다.
    
    Args:
        url (str): 열고자 하는 URL
        
    Returns:
        bool: 성공 여부
    """
    try:
        import webbrowser
        
        # URL 형식 검증
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
            
        # print(f"[INFO] 브라우저에서 URL 열기: {url}")
        webbrowser.open(url)
        return True
        
    except Exception as e:
        print(f"[ERROR] URL을 열 수 없습니다: {str(e)}")
        return False

# ============================================================================
# 검색 관련 함수 (Search Functions)
# ============================================================================

def search_content(path: str, query: str, file_types: List[str] = None) -> List[Tuple[str, List[str], int]]:
    """
    지정된 경로(파일 또는 디렉토리)에서 특정 파일들의 내용을 검색합니다.
    
    Args:
        path (str): 검색할 파일 또는 디렉토리 경로
        query (str): 검색할 텍스트
        file_types (List[str], optional): 검색할 파일 확장자 목록 (기본값: ['.txt', '.doc', '.docx', '.ppt', '.pptx', '.pdf', '.hwp', '.hwpx'])
    
    Returns:
        List[Tuple[str, List[str], int]]: [(파일경로, [검색된 라인들], 검색된 총 개수), ...]
        
    Raises:
        FileNotFoundError: 파일 또는 디렉토리가 존재하지 않는 경우
    """
    try:
        import os
        import docx
        from pptx import Presentation
        from PyPDF2 import PdfReader
        
        if file_types is None:
            file_types = ['.txt', '.doc', '.docx', '.ppt', '.pptx', '.pdf', '.hwp', '.hwpx']
            
        results = []
        query = query.lower()  # 대소문자 구분 없이 검색
        
        def process_file(file_path: str) -> Optional[Tuple[str, List[str], int]]:
            """단일 파일을 처리하는 내부 함수"""
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext not in file_types:
                    return None
                    
                content_lines = []
                match_count = 0
                
                # 텍스트 파일
                if ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8') as f:
                        lines = f.readlines()
                        for line in lines:
                            if query in line.lower():
                                content_lines.append(line.strip())
                                match_count += 1
                
                # Word 문서
                elif ext in ['.docx']:
                    doc = docx.Document(file_path)
                    for para in doc.paragraphs:
                        text = para.text
                        if query in text.lower():
                            content_lines.append(text.strip())
                            match_count += 1
                
                # PowerPoint 문서
                elif ext in ['.pptx']:
                    prs = Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text = shape.text
                                if query in text.lower():
                                    content_lines.append(text.strip())
                                    match_count += 1
                
                # PDF 문서
                elif ext == '.pdf':
                    reader = PdfReader(file_path)
                    for page in reader.pages:
                        text = page.extract_text()
                        lines = text.split('\n')
                        for line in lines:
                            if query in line.lower():
                                content_lines.append(line.strip())
                                match_count += 1
                
                # 한글 문서
                elif ext in ['.hwp', '.hwpx']:
                    text = extract_from_hwp(file_path)
                    lines = text.split('\n')
                    for line in lines:
                        if query in line.lower():
                            content_lines.append(line.strip())
                            match_count += 1
                
                if match_count > 0:
                    return (file_path, content_lines, match_count)
                return None
                
            except Exception as e:
                print("Warning: Failed to process %s: %s" % (file_path, str(e)))
                return None
        
        # 파일 또는 디렉토리 존재 여부 확인
        if not os.path.exists(path):
            raise FileNotFoundError("Path not found: %s" % path)
        
        # 단일 파일인 경우
        if os.path.isfile(path):
            result = process_file(path)
            if result:
                results.append(result)
        
        # 디렉토리인 경우
        else:
            for root, _, files in os.walk(path):
                for file in files:
                    file_path = os.path.join(root, file)
                    result = process_file(file_path)
                    if result:
                        results.append(result)
        
        return results
        
    except ImportError as e:
        raise ImportError("Required package not found: %s" % str(e))
    except Exception as e:
        raise Exception("Search failed: %s" % str(e))

def get_search_url(query: str) -> Tuple[str, str, str]:
    """
    자연어 검색 요청을 분석하여 적절한 검색 URL과 키워드를 반환합니다.
    
    Args:
        query (str): 자연어 검색 요청 (예: "유튜브에서 최신음악 찾아줘", "네이버 뉴스에서 속보 검색")
        
    Returns:
        Tuple[str, str, str]: (검색 URL 패턴, 실제 검색 키워드, 검색 타입)
    """
    # 검색 사이트 매핑
    SEARCH_SITES = {
        '유튜브': {
            'base_url': 'https://www.youtube.com',
            'keywords': ['유튜브', '유튭', 'youtube'],
            'type': 'video',
            'search_patterns': {
                'default': '/results?search_query={keyword}',
                'music': '/results?search_query={keyword}&sp=EgIQAQ%253D%253D',  # 음악 필터
                'live': '/results?search_query={keyword}&sp=EgJAAQ%253D%253D',  # 실시간 필터
                'playlist': '/results?search_query={keyword}&sp=EgIQAw%253D%253D'  # 재생목록 필터
            }
        },
        '네이버': {
            'base_url': 'https://search.naver.com',
            'keywords': ['네이버', 'naver'],
            'type': 'all',
            'search_patterns': {
                'default': '/search.naver?where=nexearch&query={keyword}',
                'news': '/search.naver?where=news&query={keyword}',
                'blog': '/search.naver?where=blog&query={keyword}',
                'cafe': '/search.naver?where=article&query={keyword}',
                'shopping': 'https://search.shopping.naver.com/search/all?query={keyword}'
            }
        },
        '구글': {
            'base_url': 'https://www.google.com',
            'keywords': ['구글', 'google'],
            'type': 'all',
            'search_patterns': {
                'default': '/search?q={keyword}',
                'news': '/search?tbm=nws&q={keyword}',
                'image': '/search?tbm=isch&q={keyword}',
                'video': '/search?tbm=vid&q={keyword}'
            }
        },
        '다음': {
            'base_url': 'https://search.daum.net',
            'keywords': ['다음', 'daum'],
            'type': 'all',
            'search_patterns': {
                'default': '/search?w=tot&q={keyword}',
                'news': '/search?w=news&q={keyword}',
                'image': '/search?w=img&q={keyword}',
                'video': '/search?w=vclip&q={keyword}',
                'blog': '/search?w=blog&q={keyword}'
            }
        }
    }
    
    # 검색어 전처리
    query = query.lower().strip()
    
    # 기본값 설정
    default_site = SEARCH_SITES['네이버']
    site_info = default_site
    search_type = 'default'
    
    # 검색 사이트 감지
    detected_site = None
    for site, info in SEARCH_SITES.items():
        for keyword in info['keywords']:
            if keyword in query.lower():
                detected_site = site
                site_info = info
                # 검색어에서 사이트 키워드와 관련된 부분 제거
                for k in info['keywords']:
                    query = query.replace(k, '').strip()
                break
        if detected_site:
            break
    
    # 검색 타입 감지
    search_keywords = {
        'news': ['뉴스', '속보', '신문'],
        'blog': ['블로그', '후기', '리뷰'],
        'image': ['이미지', '사진', '그림'],
        'video': ['동영상', '영상'],
        'shopping': ['쇼핑', '상품', '물건'],
        'music': ['음악', '노래', '뮤직'],
        'live': ['라이브', '실시간', '생방송'],
        'playlist': ['플레이리스트', '재생목록']
    }
    
    for stype, keywords in search_keywords.items():
        for keyword in keywords:
            if keyword in query:
                search_type = stype
                query = query.replace(keyword, '').strip()
                break
    
    # 검색어에서 불필요한 단어 제거
    remove_words = ['에서', '검색', '찾아', '찾아줘', '검색해줘', '보여줘', '재생', '재생해줘']
    for word in remove_words:
        query = query.replace(word, '').strip()
    
    # 검색 패턴 선택
    if search_type not in site_info['search_patterns']:
        search_type = 'default'
    
    search_pattern = site_info['search_patterns'][search_type]
    
    # 최종 URL 생성을 위한 기본 URL과 패턴 반환
    return site_info['base_url'], search_pattern, query

def create_search_result(title: str, url: str, description: str = "", image_url: str = None, result_type: str = "web", date: str = None) -> dict:
    """
    검색 결과를 표준화된 형식으로 생성하는 함수입니다.
    
    Args:
        title (str): 검색 결과 제목
        url (str): 검색 결과 URL
        description (str, optional): 검색 결과 설명
        image_url (str, optional): 이미지 URL (있는 경우)
        result_type (str, optional): 결과 타입 (web, video, news 등)
        date (str, optional): 검색 결과의 날짜 (YYYY-MM-DD 형식)
        
    Returns:
        dict: 표준화된 검색 결과
    """
    result = {
        'title': title,
        'url': url,
        'description': description,
        'type': result_type,
        'thumbnail': None,
        'image_url': None,
        'date': date  # 날짜 정보 추가
    }
    
    if image_url:
        result['thumbnail'] = image_url
        result['image_url'] = image_url
        
    return result

def search_youtube(query: str, max_results: int = None) -> List[dict]:
    """YouTube 검색을 수행합니다."""
    driver = None
    try:
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        
        driver = get_selenium_driver()
        wait = WebDriverWait(driver, 10)
        
        results = []
        search_url = f"https://www.youtube.com/results?search_query={query}"
        
        driver.get(search_url)
        video_elements = wait.until(
            EC.presence_of_all_elements_located((By.CSS_SELECTOR, "ytd-video-renderer"))
        )
        
        limit = max_results if max_results is not None else 10
        for video in video_elements[:limit]:
            try:
                title_element = video.find_element(By.CSS_SELECTOR, "#video-title")
                title = title_element.get_attribute('title')
                url = title_element.get_attribute('href')
                thumbnail = video.find_element(By.CSS_SELECTOR, "#thumbnail img").get_attribute('src')
                description = video.find_element(By.CSS_SELECTOR, "#description-text").text
                
                results.append(create_search_result(
                    title=title,
                    url=url,
                    description=description,
                    image_url=thumbnail,
                    result_type='video'
                ))
            except Exception as e:
                continue
        
        return results
        
    except Exception as e:
        return []
        
    finally:
        if driver:
            try:
                driver.quit()
            except:
                pass

def search_naver(query: str, max_results: int = None) -> List[dict]:
    """네이버 검색을 수행합니다."""
    driver = None
    try:
        import sys
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        import time
        
        driver = get_selenium_driver()
        if not driver:
            return []
            
        results = []
        search_url = f"https://search.naver.com/search.naver?where=news&query={query}"
        
        # debug_print("[DEBUG] 네이버 검색 시작")
        driver.get(search_url)
        time.sleep(2)  # 페이지 로딩 대기
        
        # 검색 결과 요소 찾기
        search_results = driver.find_elements(By.CSS_SELECTOR, "div.news_wrap")

        if not search_results:
            debug_print("[DEBUG] 검색 결과가 없습니다.")
            return []        
        
        limit = max_results if max_results is not None else 10
        for result in search_results[:limit]:
            try:
                # 제목과 링크 추출
                title_element = result.find_element(By.CSS_SELECTOR, "a.news_tit")
                title = title_element.text
                url = title_element.get_attribute("href")
                
                # 설명 추출
                description = ""
                try:
                    description = result.find_element(By.CSS_SELECTOR, "div.news_dsc").text
                except:
                    pass
                
                if title and url:
                    results.append(create_search_result(
                        title=title,
                        url=url,
                        description=description,
                        result_type='news'
                    ))
            except Exception as e:
                debug_print(f"[DEBUG] 네이버 결과 처리 중 오류: {str(e)}")
                continue
        
        return results
        
    except Exception as e:
        debug_print(f"[ERROR] 네이버 검색 실패: {str(e)}")
        return []
        
    finally:
        if driver:
            try:
                driver.quit()
            except:
                pass

def search_google(query: str, max_results: int = None) -> List[dict]:
    """구글 검색을 수행합니다."""
    driver = None
    try:
        import sys
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        import time
        
        driver = get_selenium_driver()
        if not driver:
            debug_print("[ERROR] Google 드라이버 초기화 실패")
            return []
            
        results = []
        search_url = f"https://www.google.co.kr/search?q={query}&hl=ko&lr=lang_ko"
        
        # debug_print("[DEBUG] 구글 검색 시작")
        driver.get(search_url)
        time.sleep(3)  # 페이지 로딩 대기
        
        # reCAPTCHA 체크
        if "recaptcha" in driver.page_source.lower() or "비정상적인 트래픽" in driver.page_source:
            debug_print("[DEBUG] reCAPTCHA 감지됨, 모바일 버전으로 재시도")
            # 모바일 버전으로 재시도
            driver.get(f"https://www.google.co.kr/search?q={query}&hl=ko&lr=lang_ko&source=mobile")
            time.sleep(3)
        
        # 검색 결과 요소 찾기
        search_results = driver.find_elements(By.CSS_SELECTOR, "div.g")
        if not search_results:
            search_results = driver.find_elements(By.CSS_SELECTOR, "div.xpd")  # 모바일 버전 선택자
        
        # debug_print(f"[DEBUG] Google 검색 결과 {len(search_results)}개 발견")
        if not search_results:  # 두 선택자 모두에서 결과가 없는 경우
            debug_print("[DEBUG] 검색 결과가 없습니다.")
            return []        
        
        limit = max_results if max_results is not None else 10
        for result in search_results[:limit]:
            try:
                # 제목과 링크 추출
                title_element = result.find_element(By.CSS_SELECTOR, "h3, div[role='heading']")
                title = title_element.text
                url = result.find_element(By.CSS_SELECTOR, "a").get_attribute("href")
                
                # 설명 추출
                description = ""
                try:
                    description = result.find_element(By.CSS_SELECTOR, "div.VwiC3b, div.BNeawe").text
                except:
                    pass
                
                if title and url:
                    results.append(create_search_result(
                        title=title,
                        url=url,
                        description=description,
                        result_type='web'
                    ))
            except Exception as e:
                debug_print(f"[DEBUG] 결과 처리 중 오류: {str(e)}")
                continue
        
        return results
        
    except Exception as e:
        debug_print(f"[ERROR] Google 검색 실패: {str(e)}")
        return []
        
    finally:
        if driver:
            try:
                driver.quit()
            except:
                pass

def search_daum(query: str, max_results: int = 5) -> List[dict]:
    """
    다음 검색을 수행하고 결과를 반환합니다.
    """
    driver = None
    try:
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.common.exceptions import TimeoutException, NoSuchElementException
        import time

        # debug_print("[DEBUG] 다음 검색 시작")
        driver = get_selenium_driver()
        if not driver:
            debug_print("[ERROR] Daum 드라이버 초기화 실패")
            return []

        wait = WebDriverWait(driver, 10)
        
        search_url = f"https://search.daum.net/search?w=fusion&nil_search=btn&DA=NTB&q={query}"
        driver.get(search_url)
        time.sleep(2)  # 페이지 로딩 대기
                   
        results = []
        try:
            # 검색 결과 요소 찾기 - 새로운 HTML 구조
            docs = wait.until(EC.presence_of_all_elements_located((By.CSS_SELECTOR, "c-card")))
            
            limit = max_results if max_results is not None else 10
            for doc in docs[:limit]:
                try:
                    # 제목과 링크 추출
                    title_element = doc.find_element(By.CSS_SELECTOR, "div.item-title strong.tit-g a")
                    title = title_element.text.strip()
                    url = title_element.get_attribute("href")
                    
                    # 작성자/출처 정보 추출
                    source_info = ""
                    try:
                        source_info = doc.find_element(By.CSS_SELECTOR, "div.area_tit a.item-writer").text.strip()
                    except NoSuchElementException:
                        pass
                    
                    # 설명 추출
                    description = ""
                    try:
                        description = doc.find_element(By.CSS_SELECTOR, "p.conts-desc").text.strip()
                        if source_info:
                            description = f"[{source_info}] {description}"
                    except NoSuchElementException:
                        pass
                    
                    # 날짜 추출
                    date = ""
                    try:
                        date = doc.find_element(By.CSS_SELECTOR, "span.txt_desc").text.strip()
                    except NoSuchElementException:
                        pass
                    
                    if title and url:
                        results.append(create_search_result(
                            title=title,
                            url=url,
                            description=description,
                            result_type='web',
                            date=date
                        ))
                except Exception as e:
                    debug_print(f"[DEBUG] 다음 검색 결과 항목 처리 중 오류: {str(e)}")
                    continue
                    
        # except TimeoutException:
        #     debug_print("[DEBUG] 다음 검색 결과 로딩 시간 초과")
        except Exception as e:
            debug_print(f"[DEBUG] 다음 검색 결과 처리 중 오류: {str(e)}")
        
        # debug_print(f"[DEBUG] 다음 검색 완료: {len(results)}개 결과")
        return results
        
    except Exception as e:
        debug_print(f"[DEBUG] 다음 검색 중 오류 발생: {str(e)}")
        return []
        
    finally:
        if driver:
            try:
                driver.quit()
            except:
                pass

def search_web(query: str, site: str = None, max_results: int = None) -> List[dict]:
    """
    웹 검색을 수행하는 함수
    
    Args:
        query (str): 검색어
        site (str, optional): 검색할 사이트 ('youtube', 'naver', 'google', 'daum')
        max_results (int, optional): 최대 검색 결과 수
    
    Returns:
        List[dict]: 검색 결과 목록
    """
    try:
        # 검색 사이트 지정이 있는 경우
        if site:
            site_lower = site.lower()
            if site_lower == 'youtube':
                return search_youtube(query, max_results)
            elif site_lower == 'naver':
                return search_naver(query, max_results)
            elif site_lower == 'google':
                return search_google(query, max_results)
            elif site_lower == 'daum':
                return search_daum(query, max_results)
            else:
                return []
        
        # 검색 사이트를 지정하지 않은 경우 구글에서만 검색
        return search_google(query, max_results)
            
    except Exception as e:
        return []

def open_search_result(results: List[dict], index: int = 0) -> bool:
    """
    검색 결과 중 지정된 인덱스의 URL을 브라우저에서 엽니다.
    
    Args:
        results (List[dict]): 검색 결과 리스트
        index (int): 열고자 하는 결과의 인덱스 (기본값: 0, 첫 번째 결과)
        
    Returns:
        bool: 성공 여부
    """
    try:
        if not results or len(results) <= index:
            print(f"[ERROR] 인덱스 {index}의 검색 결과가 없습니다.")
            return False
            
        result = results[index]
        if 'url' not in result:
            print("[ERROR] 검색 결과에 URL이 없습니다.")
            return False
            
        return open_in_browser(result['url'])
        
    except Exception as e:
        print(f"[ERROR] 검색 결과를 열 수 없습니다: {str(e)}")
        return False

# ============================================================================
# 이메일 관련 (Email Functions)
# ============================================================================

def send_email(to_email: str, subject: str, body: str, attachments: List[str] = None, html_body: str = None) -> bool:
    """
    이메일을 발송합니다.
    
    Args:
        to_email (str): 받는 사람 이메일 주소
        subject (str): 이메일 제목
        body (str): 이메일 본문 (텍스트)
        attachments (List[str], optional): 첨부 파일 경로 리스트
        html_body (str, optional): HTML 형식의 이메일 본문
        
    Returns:
        bool: 발송 성공 여부
    """
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        from email.mime.application import MIMEApplication
        from email.utils import formatdate
        import os
        
        # SMTP_SECURE 설정 가져오기
        config = load_config()
        secure = config.get("SMTP_SECURE", "YES").upper() not in ["NO", "N"]
        
        print(f"\n[INFO] SMTP Settings:")
        print(f"- Server: {SMTP_HOST}")
        print(f"- Port: {SMTP_PORT}")
        print(f"- Account: {SMTP_USERNAME}")
        print(f"- Secure: {secure}")
        
        # 메시지 생성
        msg = MIMEMultipart('alternative')
        msg['From'] = SMTP_USERNAME
        msg['To'] = to_email
        msg['Subject'] = subject
        msg['Date'] = formatdate(localtime=True)
        
        # 텍스트 본문 추가
        msg.attach(MIMEText(body, 'plain', 'utf-8'))
        
        # HTML 본문이 있는 경우 추가
        if html_body:
            msg.attach(MIMEText(html_body, 'html', 'utf-8'))
        
        # 첨부 파일 처리
        if attachments:
            for file_path in attachments:
                try:
                    if os.path.exists(file_path):
                        with open(file_path, 'rb') as f:
                            part = MIMEApplication(f.read(), Name=os.path.basename(file_path))
                            part['Content-Disposition'] = f'attachment; filename="{os.path.basename(file_path)}"'
                            msg.attach(part)
                    else:
                        print(f"[WARNING] 첨부 파일을 찾을 수 없습니다: {file_path}")
                except Exception as e:
                    print(f"[WARNING] 첨부 파일 처리 중 오류 발생: {str(e)}")
                    continue
        
        # SMTP 서버 연결 및 이메일 발송
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            if secure:
                server.starttls()  # TLS 보안 연결 (SMTP_SECURE가 true인 경우에만)
            server.login(SMTP_USERNAME, SMTP_PASSWORD)
            server.send_message(msg)
            
        print(f"[INFO] 이메일이 성공적으로 발송되었습니다: {to_email}")
        return True
        
    except Exception as e:
        print(f"[ERROR] 이메일 발송 실패: {str(e)}")
        return False

# ============================================================================
# DeepL 관련 클래스와 함수 (DeepL Related Classes and Functions)
# ============================================================================

class DeepLTranslator:
    def __init__(self):
        self._install_required_packages()
        
        # 설정 파일에서 API 키 읽기
        config = load_config()
        self.api_key = config.get('DEEPL_API_KEY') or os.getenv('DEEPL_API_KEY')
        
        if not self.api_key:
            raise ValueError("DeepL API key is required")
        
        # Free API 엔드포인트 사용
        self.base_url = "https://api-free.deepl.com/v2"
        self.max_length = 4500  # Free API 안전 제한

    def _install_required_packages(self):
        packages = ['aiohttp']
        for package in packages:
            install_if_missing(package)
            
    def _split_text(self, text: str) -> list:
        """텍스트를 문장 단위로 분할합니다."""
        chunks = []
        current_chunk = []
        current_length = 0
        
        # 문단 단위로 먼저 분리
        paragraphs = text.split('\n\n')
        
        for paragraph in paragraphs:
            # 문장 단위로 분리
            sentences = paragraph.replace('. ', '.\n').replace('? ', '?\n').replace('! ', '!\n').split('\n')
            
            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue
                
                # 문장이 단독으로 제한을 초과하는 경우
                if len(sentence) > self.max_length:
                    # 현재 청크가 있으면 먼저 추가
                    if current_chunk:
                        chunks.append('\n'.join(current_chunk))
                        current_chunk = []
                        current_length = 0
                    
                    # 긴 문장을 강제로 분할
                    while sentence:
                        chunks.append(sentence[:self.max_length])
                        sentence = sentence[self.max_length:]
                    continue
                
                # 현재 청크에 문장을 추가했을 때 제한을 초과하는 경우
                if current_length + len(sentence) > self.max_length:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
                
                current_chunk.append(sentence)
                current_length += len(sentence)
        
        # 마지막 청크 처리
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
        
        return chunks

    async def translate(self, text: str, target_lang: str = 'EN', formality: str = 'default') -> dict:
        try:
            chunks = self._split_text(text)
            translated_chunks = []
            
            headers = {
                'Authorization': f'DeepL-Auth-Key {self.api_key}',
                'Content-Type': 'application/json'
            }
            
            async with aiohttp.ClientSession() as session:
                for chunk in chunks:
                    if not chunk.strip():
                        continue
                        
                    params = {
                        'text': [chunk],  # text 파라미터를 배열로 변경
                        'target_lang': target_lang.upper(),
                        'formality': formality
                    }
                    
                    async with session.post(f"{self.base_url}/translate", 
                                          headers=headers,
                                          json=params) as response:
                        if response.status == 200:
                            result = await response.json()
                            translated_chunks.append(result['translations'][0]['text'])
                        else:
                            error_text = await response.text()
                            return {
                                'success': False,
                                'translated_text': '',
                                'error': f'Translation failed: {error_text}'
                            }
            
            return {
                'success': True,
                'translated_text': '\n'.join(translated_chunks),
                'error': ''
            }
        except Exception as e:
            return {
                'success': False,
                'translated_text': '',
                'error': str(e)
            }

    async def get_supported_languages(self) -> list:
        """지원되는 언어 목록을 반환합니다."""
        headers = {
            'Authorization': f'DeepL-Auth-Key {self.api_key}'
        }
        
        async with aiohttp.ClientSession() as session:
            async with session.get(f"{self.base_url}/languages", 
                                 headers=headers) as response:
                if response.status == 200:
                    languages = await response.json()
                    return languages
                return []
            
# ============================================================================
# AI 관련 클래스와 함수 (AI Related Classes and Functions)
# ============================================================================

class AIProvider:
    def __init__(self):
        # Install and import required packages
        self._install_required_packages()
        self._import_providers()
        
        # Load configuration
        self.config = config
        self.provider = config.get('USE_LLM', 'openai')
        self.language = config.get('LANGUAGE', 'en')
        
        # Default max tokens for each provider and model
        self.default_max_tokens = {
            'openai': {
                'gpt-3.5-turbo': 16385,
                'gpt-4': 8192,
                'gpt-4-32k': 32768,
                'gpt-4-turbo': 128000,
                'gpt-4-vision': 128000,
                'gpt-4-all': 128000
            },
            'anthropic': {
                'claude-3-opus': 200000,
                'claude-3-sonnet': 200000,
                'claude-3-haiku': 200000
            },
            'gemini': {
                'gemini-pro': 1000000,
                'gemini-ultra': 1000000
            },
            'groq': {
                'mixtral-8x7b-32768': 32768,
                'llama2-70b-4096': 4096
            },
            'ollama': {
                'llama3': 8192,
                'mistral': 8192,
                'mixtral': 32768
            }
        }
        
        # Default prompts
        self.default_prompts = {
            'summarize': "Please provide a clear and concise summary of the following text:",
            'translate': "Please translate the following text accurately:",
            'review': "Please review the following code for quality, bugs, and performance issues:"
        }
        
        # Load custom prompts if available
        self.custom_prompts = {
            'summarize': self._load_prompt('summarize.prompt'),
            'translate': self._load_prompt('translate.prompt'),
            'review': self._load_prompt('review.prompt')
        }
    
    def _install_required_packages(self):
        packages = ['openai', 'anthropic', 'google-generativeai', 'groq', 'requests']
        for package in packages:
            install_if_missing(package)
    
    def _import_providers(self):
        """Import required packages for each provider"""
        try:
            import openai
            self.openai = openai
        except ImportError:
            self.openai = None
            
        try:
            import anthropic
            self.anthropic = anthropic
        except ImportError:
            self.anthropic = None
            
        try:
            import google.generativeai as genai
            self.genai = genai
        except ImportError:
            self.genai = None
            
        try:
            from groq import Groq
            self.groq = Groq
        except ImportError:
            self.groq = None

    def _load_prompt(self, filename):
        try:
            # ~/.airun/ 디렉토리에서 프롬프트 파일 찾기
            prompt_path = os.path.join(os.path.expanduser("~"), ".airun", filename)
            # print(f"Loading prompt from: {prompt_path}")
            if os.path.exists(prompt_path):
                with open(prompt_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                # print(f"Successfully loaded prompt: {filename}")
                return content
            # print(f"Prompt file not found: {prompt_path}, using default prompt")
            default_prompt = self.default_prompts.get(filename.split('.')[0])
            if default_prompt:
                # print(f"Using default prompt for {filename}")
                return default_prompt
            raise ValueError(f"No default prompt available for {filename}")
        except Exception as e:
            print(f"Error loading prompt file {filename}: {str(e)}")
            default_prompt = self.default_prompts.get(filename.split('.')[0])
            if default_prompt:
                print(f"Using default prompt for {filename}")
                return default_prompt
            raise ValueError(f"Failed to load prompt and no default available for {filename}")

    def _prepare_content(self, content):
        if isinstance(content, str) and os.path.exists(content):
            return read_file(content)
        return content

    def _call_provider(self, messages, max_tokens=None):
        # If max_tokens is not specified, use provider's default
        if max_tokens is None:
            max_tokens = self.default_max_tokens.get(self.provider, 2048)

        try:
            if self.provider == 'openai':
                if not config.get('OPENAI_API_KEY'):
                    raise ValueError("OpenAI API key is not configured.")
                if not self.openai:
                    raise ImportError("Failed to import openai package")
                    
                model = config.get('OPENAI_MODEL', 'gpt-4-turbo')
                client = self.openai.OpenAI(api_key=config.get('OPENAI_API_KEY'))
                response = client.chat.completions.create(
                    model=model,
                    messages=messages,
                    max_tokens=max_tokens
                )
                return response.choices[0].message.content

            elif self.provider == 'anthropic':
                if not config.get('ANTHROPIC_API_KEY'):
                    raise ValueError("Anthropic API key is not configured.")
                    
                model = config.get('ANTHROPIC_MODEL', 'claude-3-opus-20240229')
                response = self.anthropic.messages.create(
                    model=model,
                    messages=messages,
                    max_tokens=max_tokens
                )
                return response.content

            elif self.provider == 'gemini':
                if not config.get('GOOGLE_API_KEY'):
                    raise ValueError("Google API key is not configured.")
                
                import google.generativeai as genai
                genai.configure(api_key=config.get('GOOGLE_API_KEY'))
                model = config.get('GEMINI_MODEL', 'gemini-1.5-pro')
                
                # Convert messages to Gemini format
                prompt = "\n".join([m['content'] for m in messages])
                response = self.genai.GenerativeModel(model).generate_content(
                    prompt,
                    generation_config=self.genai.types.GenerationConfig(
                        max_output_tokens=max_tokens
                    )
                )
                return response.text

            elif self.provider == 'groq':
                if not config.get('GROQ_API_KEY'):
                    raise ValueError("Groq API key is not configured.")
                
                from groq import Groq
                client = Groq(api_key=config.get('GROQ_API_KEY'))
                model = config.get('GROQ_MODEL', 'mixtral-8x7b-32768')
                
                response = client.chat.completions.create(
                    model=model,
                    messages=messages,
                    max_tokens=max_tokens
                )
                return response.choices[0].message.content

            elif self.provider == 'ollama':
                import requests
                
                ollama_url = config.get('OLLAMA_PROXY_SERVER', 'http://localhost:11434')
                model = config.get('OLLAMA_MODEL', 'llama3:latest')
                
                if config.get('OLLAMA_PROXY_SERVER'):
                    response = requests.post(
                        ollama_url,
                        json={
                            'proxybody': {
                                'model': model,
                                'messages': messages,
                                'stream': False
                            }
                        }
                    )
                else:
                    response = requests.post(
                        'http://localhost:11434/api/chat',
                        json={
                            'model': model,
                            'messages': messages,
                            'stream': False
                        }
                    )
                
                response.raise_for_status()
                return response.json()['message']['content']

            else:
                raise ValueError(f"Unsupported AI provider: {self.provider}")

        except Exception as e:
            print("Error occurred while calling AI provider: %s" % str(e))
            raise

    def _get_model_max_tokens(self):
        """현재 사용 중인 모델의 최대 토큰 수를 반환합니다."""
        if self.provider == 'openai':
            model = self.config.get('OPENAI_MODEL', 'gpt-4-turbo')
            return self.default_max_tokens['openai'].get(model, 8192)  # 기본값 8192
        elif self.provider == 'anthropic':
            model = self.config.get('ANTHROPIC_MODEL', 'claude-3-opus')
            return self.default_max_tokens['anthropic'].get(model, 200000)
        elif self.provider == 'gemini':
            model = self.config.get('GEMINI_MODEL', 'gemini-pro')
            return self.default_max_tokens['gemini'].get(model, 1000000)
        elif self.provider == 'groq':
            model = self.config.get('GROQ_MODEL', 'mixtral-8x7b-32768')
            return self.default_max_tokens['groq'].get(model, 32768)
        elif self.provider == 'ollama':
            model = self.config.get('OLLAMA_MODEL', 'llama3')
            return self.default_max_tokens['ollama'].get(model, 8192)
        return 4000  # 기본값
    
class TextProcessor(AIProvider):
    # 모델별 비율 설정을 클래스 변수로 정의
    MODEL_RATIOS = {
        'openai': {
            'gpt-3.5-turbo': {'chunk': 0.4, 'summary': 0.3},
            'gpt-4': {'chunk': 0.5, 'summary': 0.4},
            'gpt-4-32k': {'chunk': 0.6, 'summary': 0.4},
            'gpt-4-turbo': {'chunk': 0.6, 'summary': 0.4},
            'gpt-4-vision': {'chunk': 0.6, 'summary': 0.4},
            'gpt-4-all': {'chunk': 0.6, 'summary': 0.4}
        },
        'anthropic': {
            'claude-3-opus': {'chunk': 0.6, 'summary': 0.4},
            'claude-3-sonnet': {'chunk': 0.6, 'summary': 0.4},
            'claude-3-haiku': {'chunk': 0.5, 'summary': 0.3}
        },
        'gemini': {
            'gemini-pro': {'chunk': 0.6, 'summary': 0.4},
            'gemini-ultra': {'chunk': 0.6, 'summary': 0.4}
        },
        'groq': {
            'mixtral-8x7b-32768': {'chunk': 0.5, 'summary': 0.35},
            'llama2-70b-4096': {'chunk': 0.4, 'summary': 0.3}
        },
        'ollama': {
            'llama3': {'chunk': 0.4, 'summary': 0.3},
            'mistral': {'chunk': 0.4, 'summary': 0.3},
            'mixtral': {'chunk': 0.5, 'summary': 0.35}
        }
    }
    
    def _get_model_ratios(self):
        """현재 모델의 비율을 반환합니다."""
        current_model = self.config.get(f'{self.provider.upper()}_MODEL', '')
        default_ratios = {'chunk': 0.4, 'summary': 0.3}
        return self.MODEL_RATIOS.get(self.provider, {}).get(current_model, default_ratios)

    def summarize(self, content, max_length=None):
        content = self._prepare_content(content)
        
        # 현재 모델의 최대 토큰 수 가져오기
        model_max_tokens = self._get_model_max_tokens()
        current_model = self.config.get(f'{self.provider.upper()}_MODEL', '')
        
        # 현재 모델의 비율 가져오기
        model_ratio = self._get_model_ratios()
        
        # OpenAI의 경우 max_tokens 제한
        if self.provider == 'openai':
            max_chunk_size = min(int(model_max_tokens * model_ratio['chunk']), 4000)  # OpenAI의 max_tokens 제한
            max_summary_tokens = min(int(model_max_tokens * model_ratio['summary']), 4000)
        else:
            max_chunk_size = int(model_max_tokens * model_ratio['chunk'])
            max_summary_tokens = int(model_max_tokens * model_ratio['summary'])
        
        # print(f"[INFO] Current model: {current_model}")
        # print(f"[INFO] Max tokens: {model_max_tokens}")
        # print(f"[INFO] Using max chunk size of {max_chunk_size} tokens")
        # print(f"[INFO] Max summary tokens: {max_summary_tokens}")
        
        # 청크 크기와 요약본 토큰 수 계산
        max_chunk_size = int(model_max_tokens * model_ratio['chunk'])
        max_summary_tokens = int(model_max_tokens * model_ratio['summary'])
        
        # print(f"[INFO] Current model: {current_model}")
        # print(f"[INFO] Max tokens: {model_max_tokens}")
        # print(f"[INFO] Using max chunk size of {max_chunk_size} tokens")
        
        # 긴 텍스트를 의미 단위로 나누는 함수
        def split_into_semantic_chunks(text):
            chunks = []
            current_chunk = []
            current_size = 0
            
            # 단락 단위로 먼저 분리
            paragraphs = text.split('\n\n')
            
            for paragraph in paragraphs:
                # 토큰 수 추정 (한글: 1자당 2-3토큰, 영어: 1자당 0.3-0.5토큰)
                is_korean = any(ord('가') <= ord(c) <= ord('힣') for c in paragraph)
                token_ratio = 2.5 if is_korean else 0.4
                paragraph_size = int(len(paragraph.encode('utf-8')) * token_ratio)
                
                # 단락이 너무 길면 문장 단위로 분리
                if paragraph_size > max_chunk_size:
                    sentences = paragraph.replace('!', '.').replace('?', '.').split('.')
                    for sentence in sentences:
                        if not sentence.strip():
                            continue
                            
                        sentence = sentence.strip() + '.'
                        sentence_size = int(len(sentence.encode('utf-8')) * token_ratio)
                        
                        # 현재 청크가 제한에 근접하면 새 청크 시작
                        if current_size + sentence_size > max_chunk_size:
                            if current_chunk:
                                chunks.append('\n'.join(current_chunk))
                                current_chunk = []
                                current_size = 0
                        
                        current_chunk.append(sentence)
                        current_size += sentence_size
                else:
                    # 현재 청크가 제한에 근접하면 새 청크 시작
                    if current_size + paragraph_size > max_chunk_size:
                        if current_chunk:
                            chunks.append('\n'.join(current_chunk))
                            current_chunk = []
                            current_size = 0
                    
                    if paragraph.strip():
                        current_chunk.append(paragraph)
                        current_size += paragraph_size
            
            # 남은 내용 처리
            if current_chunk:
                chunks.append('\n'.join(current_chunk))
            
            return chunks
        
        try:
            # 텍스트를 의미 단위로 나누기
            chunks = split_into_semantic_chunks(content)
            print(f"[INFO] Dividing text into {len(chunks)} semantic chunks...")
            
            if len(chunks) == 1:
                # 단일 청크인 경우 직접 요약
                system_prompt = self.custom_prompts.get('summarize')
                if not system_prompt:
                    system_prompt = """Please provide a clear and concise summary of the following text, maintaining the key points and overall structure:
1. Maintain the logical flow and relationships between ideas
2. Preserve the hierarchical importance of information
3. Ensure key themes and concepts are properly connected
4. Provide a comprehensive yet concise overview"""
                
                user_lang = self.config.get('LANGUAGE', 'en')
                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "system", "content": f"The user's language is {user_lang}. Please provide the summary in this language."},
                    {"role": "user", "content": chunks[0]}
                ]
                
                return self._call_provider(messages, max_tokens=max_length or int(max_chunk_size * 0.5))
            
            # 여러 청크가 있는 경우 계층적 요약
            summaries = []
            print("[INFO] Starting hierarchical summarization...")
            
            # 1단계: 각 청크의 핵심 내용 추출 (더 짧은 요약)
            for i, chunk in enumerate(chunks, 1):
                print(f"[INFO] Extracting key points from chunk {i}/{len(chunks)}...")
                messages = [
                    {"role": "system", "content": "Extract only the most essential points from this text in a very concise manner:"},
                    {"role": "user", "content": chunk}
                ]
                key_points = self._call_provider(messages, max_tokens=int(max_chunk_size * 0.3))
                summaries.append(key_points)
            
            # 중간 요약들이 너무 길면 다시 나누어 요약
            while len('\n\n'.join(summaries).encode('utf-8')) // 3 > max_chunk_size:
                print("[INFO] Intermediate summaries too long, performing additional summarization...")
                new_summaries = []
                temp_summaries = []
                current_size = 0
                
                for summary in summaries:
                    summary_size = len(summary.encode('utf-8')) // 3
                    if current_size + summary_size > max_chunk_size:
                        if temp_summaries:
                            combined = '\n\n'.join(temp_summaries)
                            messages = [
                                {"role": "system", "content": "Create an extremely concise summary of these points:"},
                                {"role": "user", "content": combined}
                            ]
                            new_summary = self._call_provider(messages, max_tokens=int(max_chunk_size * 0.3))
                            new_summaries.append(new_summary)
                            temp_summaries = [summary]
                            current_size = summary_size
                    else:
                        temp_summaries.append(summary)
                        current_size += summary_size
                
                if temp_summaries:
                    combined = '\n\n'.join(temp_summaries)
                    messages = [
                        {"role": "system", "content": "Create an extremely concise summary of these points:"},
                        {"role": "user", "content": combined}
                    ]
                    new_summary = self._call_provider(messages, max_tokens=int(max_chunk_size * 0.3))
                    new_summaries.append(new_summary)
                
                summaries = new_summaries
            
            # 2단계: 최종 요약 생성
            print("[INFO] Creating final summary...")
            system_prompt = self.custom_prompts.get('summarize')
            if not system_prompt:
                system_prompt = """Please provide a clear and concise summary of the following text, maintaining the key points and overall structure:
1. Maintain the logical flow and relationships between ideas
2. Preserve the hierarchical importance of information
3. Ensure key themes and concepts are properly connected
4. Provide a comprehensive yet concise overview"""

            final_messages = [
                {"role": "system", "content": system_prompt},
                {"role": "system", "content": f"The user's language is {self.language}. Create the summary in this language."},
                {"role": "user", "content": '\n\n'.join(summaries)}
            ]
            
            # max_tokens를 4096 이하로 제한
            max_final_tokens = min(max_length or int(max_chunk_size * 0.5), 4000)  # 여유있게 4000으로 제한
            print(f"[INFO] Using max_tokens of {max_final_tokens} for final summary")
            
            return self._call_provider(final_messages, max_tokens=max_final_tokens)
            
        except Exception as e:
            print(f"[ERROR] Summarization failed: {str(e)}")
            raise

    def translate(self, content, target_language):
        content = self._prepare_content(content)
        
        # 현재 모델의 최대 토큰 수 가져오기
        model_max_tokens = self._get_model_max_tokens()
        current_model = self.config.get(f'{self.provider.upper()}_MODEL', '')
        
        # 현재 모델의 비율 가져오기
        model_ratio = self._get_model_ratios()
        
        # OpenAI의 경우 max_tokens 제한
        if self.provider == 'openai':
            max_chunk_size = min(int(model_max_tokens * model_ratio['chunk']), 4000)
            max_response_tokens = min(int(model_max_tokens * model_ratio['summary']), 4000)
        else:
            max_chunk_size = int(model_max_tokens * model_ratio['chunk'])
            max_response_tokens = int(model_max_tokens * model_ratio['summary'])
        
        print(f"\n[INFO] 번역 작업 시작")
        print(f"[INFO] 현재 모델: {current_model}")
        print(f"[INFO] 최대 토큰 수: {model_max_tokens}")
        print(f"[INFO] 청크 크기: {max_chunk_size} 토큰")
        print(f"[INFO] 응답 토큰 수: {max_response_tokens}")
        
        try:
            import tempfile
            import os
            
            # 입력 텍스트 크기 계산
            content_size = len(content.encode('utf-8'))
            print(f"[INFO] 입력 텍스트 크기: {content_size / 1024:.2f}KB")
            
            # 입력 텍스트가 1MB를 초과하는 경우 임시 파일 사용
            use_temp_files = content_size > 1024 * 1024
            temp_dir = None
            
            if use_temp_files:
                temp_dir = tempfile.mkdtemp(prefix='translation_')
                print(f"[INFO] 임시 디렉토리 생성됨: {temp_dir}")
                print(f"[INFO] 임시 디렉토리 존재 여부: {os.path.exists(temp_dir)}")
            
            # 텍스트를 청크로 나누는 함수
            def split_into_chunks(text):
                chunks = []
                current_chunk = []
                current_size = 0
                chunk_count = 0
                
                print(f"[INFO] 텍스트 분할 시작...")
                
                # 문단 단위로 먼저 분리
                paragraphs = text.split('\n\n')
                print(f"[INFO] 총 {len(paragraphs)}개의 문단 발견")
                
                for i, paragraph in enumerate(paragraphs, 1):
                    # 토큰 수 추정
                    is_korean = any(ord('가') <= ord(c) <= ord('힣') for c in paragraph)
                    token_ratio = 2.5 if is_korean else 0.4
                    paragraph_size = int(len(paragraph.encode('utf-8')) * token_ratio)
                    
                    # 현재 청크가 제한에 근접하면 새 청크 시작
                    if current_size + paragraph_size > max_chunk_size:
                        if current_chunk:
                            chunk_text = '\n'.join(current_chunk)
                            if use_temp_files:
                                chunk_file = os.path.join(temp_dir, f'chunk_{chunk_count}.txt')
                                with open(chunk_file, 'w', encoding='utf-8') as f:
                                    f.write(chunk_text)
                                print(f"[INFO] 청크 {chunk_count} 저장됨: {chunk_file} ({len(chunk_text.encode('utf-8'))/1024:.2f}KB)")
                                chunks.append(chunk_file)
                            else:
                                chunks.append(chunk_text)
                            chunk_count += 1
                            current_chunk = []
                            current_size = 0
                    
                    current_chunk.append(paragraph)
                    current_size += paragraph_size
                    print(f"[INFO] 문단 {i}/{len(paragraphs)} 처리 중... (현재 청크 크기: {current_size} 토큰)")
                
                # 남은 내용 처리
                if current_chunk:
                    chunk_text = '\n'.join(current_chunk)
                    if use_temp_files:
                        chunk_file = os.path.join(temp_dir, f'chunk_{chunk_count}.txt')
                        with open(chunk_file, 'w', encoding='utf-8') as f:
                            f.write(chunk_text)
                        print(f"[INFO] 마지막 청크 {chunk_count} 저장됨: {chunk_file} ({len(chunk_text.encode('utf-8'))/1024:.2f}KB)")
                        chunks.append(chunk_file)
                    else:
                        chunks.append(chunk_text)
                
                print(f"[INFO] 총 {len(chunks)}개의 청크로 분할 완료")
                return chunks
            
            # 텍스트를 청크로 나누기
            chunks = split_into_chunks(content)
            
            # 번역된 청크를 저장할 리스트 또는 임시 파일들
            translated_chunks = []
            
            # 언어 코드 정규화
            normalized_target = normalize_language_code(target_language)
            print(f"[INFO] 대상 언어 코드: {normalized_target}")
            
            # 번역 프롬프트 준비
            system_prompt = self.custom_prompts.get('translate')
            if not system_prompt:
                system_prompt = """Please translate the following text to the specified target language.
                Maintain the original meaning, nuance, and formatting as accurately as possible.
                Preserve all line breaks and paragraph structures."""
            
            # 각 청크 번역
            for i, chunk in enumerate(chunks, 1):
                print(f"\n[INFO] 청크 {i}/{len(chunks)} 번역 시작...")
                
                # 청크 내용 읽기
                if use_temp_files:
                    print(f"[INFO] 청크 파일 읽기: {chunk}")
                    with open(chunk, 'r', encoding='utf-8') as f:
                        chunk_content = f.read()
                    print(f"[INFO] 청크 크기: {len(chunk_content.encode('utf-8'))/1024:.2f}KB")
                else:
                    chunk_content = chunk
                
                # 번역 수행
                print(f"[INFO] AI 모델에 번역 요청 중...")
                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "system", "content": f"Target language: {target_language} (normalized code: {normalized_target})"},
                    {"role": "user", "content": chunk_content}
                ]
                
                translated_text = self._call_provider(messages, max_tokens=max_chunk_size)
                print(f"[INFO] 번역 완료 (결과 크기: {len(translated_text.encode('utf-8'))/1024:.2f}KB)")
                
                # 번역된 텍스트 저장
                if use_temp_files:
                    translated_file = os.path.join(temp_dir, f'translated_{i}.txt')
                    with open(translated_file, 'w', encoding='utf-8') as f:
                        f.write(translated_text)
                    print(f"[INFO] 번역 결과 저장됨: {translated_file}")
                    translated_chunks.append(translated_file)
                else:
                    translated_chunks.append(translated_text)
            
            print("\n[INFO] 번역된 청크 병합 중...")
            # 번역된 청크 합치기
            if use_temp_files:
                final_text = []
                for chunk_file in translated_chunks:
                    print(f"[INFO] 번역된 파일 읽기: {chunk_file}")
                    with open(chunk_file, 'r', encoding='utf-8') as f:
                        final_text.append(f.read())
                result = '\n'.join(final_text)
            else:
                result = '\n'.join(translated_chunks)
            
            print(f"[INFO] 최종 번역 완료 (결과 크기: {len(result.encode('utf-8'))/1024:.2f}KB)")
            return result
            
        except Exception as e:
            print(f"[ERROR] 번역 실패: {str(e)}")
            raise
            
        finally:
            # 임시 파일 및 디렉토리 정리
            if use_temp_files and temp_dir:
                try:
                    import shutil
                    print(f"\n[INFO] 임시 파일 정리 중...")
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                        print(f"[INFO] 임시 디렉토리 삭제됨: {temp_dir}")
                except Exception as e:
                    print(f"[WARNING] 임시 파일 정리 실패: {str(e)}")

    def review_code(self, content):
        content = self._prepare_content(content)
        
        # Use custom prompt if available
        system_prompt = self.custom_prompts.get('review')
        if not system_prompt:
            system_prompt = "Please review the following code for quality, bugs, and performance issues:"
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": content}
        ]
        
        return self._call_provider(messages)

def normalize_language_code(language: str) -> str:
    """
    언어 이름이나 코드를 ISO 639-1 코드로 정규화합니다.
    
    Args:
        language (str): 언어 이름 또는 코드 (예: '한국어', 'korean', 'ko', '영어', 'english', 'en' 등)
        
    Returns:
        str: 정규화된 ISO 639-1 언어 코드
    """
    # 언어 코드 매핑
    LANGUAGE_CODES = {
        # 한국어
        '한국어': 'ko',
        'korean': 'ko',
        'ko': 'ko',
        'kor': 'ko',
        
        # 영어
        '영어': 'en',
        'english': 'en',
        'en': 'en',
        'eng': 'en',
        
        # 일본어
        '일본어': 'ja',
        'japanese': 'ja',
        'ja': 'ja',
        'jpn': 'ja',
        
        # 중국어
        '중국어': 'zh',
        'chinese': 'zh',
        'zh': 'zh',
        'chi': 'zh',
        '중국어(간체)': 'zh-CN',
        '중국어(번체)': 'zh-TW',
        'simplified chinese': 'zh-CN',
        'traditional chinese': 'zh-TW',
        
        # 프랑스어
        '프랑스어': 'fr',
        'french': 'fr',
        'fr': 'fr',
        'fra': 'fr',
        
        # 독일어
        '독일어': 'de',
        'german': 'de',
        'de': 'de',
        'deu': 'de',
        
        # 스페인어
        '스페인어': 'es',
        'spanish': 'es',
        'es': 'es',
        'spa': 'es',
        
        # 이탈리아어
        '이탈리아어': 'it',
        'italian': 'it',
        'it': 'it',
        'ita': 'it',
        
        # 러시아어
        '러시아어': 'ru',
        'russian': 'ru',
        'ru': 'ru',
        'rus': 'ru',
        
        # 베트남어
        '베트남어': 'vi',
        'vietnamese': 'vi',
        'vi': 'vi',
        'vie': 'vi',
        
        # 태국어
        '태국어': 'th',
        'thai': 'th',
        'th': 'th',
        'tha': 'th',
        
        # 인도네시아어
        '인도네시아어': 'id',
        'indonesian': 'id',
        'id': 'id',
        'ind': 'id'
    }
    
    # 입력값 전처리
    normalized = language.lower().strip()
    
    # 매핑된 코드 반환 또는 입력값 그대로 반환
    return LANGUAGE_CODES.get(normalized, normalized)

def summarize_content(content, provider=None, max_length=None):
    processor = TextProcessor()
    return processor.summarize(content, max_length)

def translate_content(content, target_language, provider=None):
    processor = TextProcessor()
    return processor.translate(content, target_language)

def review_code(content, provider=None):
    processor = TextProcessor()
    return processor.review_code(content)

# ============================================================================
# 디버깅 유틸리티 (Debug Utilities)
# ============================================================================

def debug_print(message: str) -> None:
    """
    디버그 메시지를 stderr로 출력합니다.
    
    Args:
        message (str): 출력할 디버그 메시지
    """
    print(message, file=sys.stderr)

# DeepL 번역기 테스트를 위한 코드
async def test_deepl_translator():
    try:
        # DeepL 번역기 초기화
        translator = DeepLTranslator()
        print("[DEBUG] DeepL 번역기가 성공적으로 초기화되었습니다.")
        
        # 테스트할 텍스트
        test_text = "안녕하세요. 이것은 테스트 메시지입니다."
        print(f"[DEBUG] 원본 텍스트: {test_text}")
        
        # 번역 실행
        result = await translator.translate(test_text, target_lang='EN')
        print(f"[DEBUG] 번역 결과: {result}")
        
        # 지원 언어 확인
        languages = await translator.get_supported_languages()
        print(f"[DEBUG] 지원되는 언어 목록: {languages}")
        
    except Exception as e:
        print(f"[ERROR] 테스트 중 오류 발생: {str(e)}")

# 테스트 실행을 위한 코드
if __name__ == "__main__":
    import asyncio
    
    # 테스트 코드 실행
    asyncio.run(test_deepl_translator())

def clean_hwp_text(text):
    """
    HWP에서 추출된 텍스트를 정제합니다.
    """
    import re
    
    if not text:
        return ""
        
    # 1. 불필요한 태그 제거
    text = re.sub(r'<표>|<그림>', '', text)
    
    # 2. 특수 문자 처리
    special_chars = {
        '󰊱': '1.',
        '󰊲': '2.',
        '󰊳': '3.',
        '󰊴': '4.',
        '󰊵': '5.'
    }
    for char, replacement in special_chars.items():
        text = text.replace(char, replacement)
    
    # 3. 연속된 빈 줄 제거
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 4. 줄 시작의 불필요한 공백 제거
    text = '\n'.join(line.strip() for line in text.split('\n'))
    
    # 5. 표 형식 정리 (단위: 등의 표시 처리)
    text = re.sub(r'\(단위\s*:\s*([^)]+)\)', r'(단위: \1)', text)
    
    # 6. 괄호 안의 공백 정리
    text = re.sub(r'\(\s+', '(', text)
    text = re.sub(r'\s+\)', ')', text)
    
    return text.strip()

def extract_structure(text):
    """
    텍스트에서 문서의 구조를 추출합니다.
    """
    import re
    
    if not text:
        return []
        
    sections = []
    current_section = []
    
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            if current_section:
                current_section.append('')
            continue
            
        # 제목 패턴 확인
        is_title = bool(re.match(r'^[0-9０-９]+[\.\s]|^[ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩⅪⅫⅰⅱⅲⅳⅴⅵⅶⅷⅸⅹⅺⅻ][\.\s]', line))
        
        if is_title:
            if current_section:
                sections.append('\n'.join(current_section).strip())
                current_section = []
            current_section.append(line)
        else:
            current_section.append(line)
    
    if current_section:
        sections.append('\n'.join(current_section).strip())
    
    return [section for section in sections if section.strip()]

def convert_hwp_to_pdf(hwp_path: str) -> str:
    """
    HWP 파일을 PDF로 변환합니다.
    
    Args:
        hwp_path (str): HWP 파일 경로
        
    Returns:
        str: 변환된 PDF 파일 경로 (성공 시) 또는 빈 문자열 (실패 시)
    """
    try:
        import os
        import subprocess
        import tempfile
        
        print(f"\nStarting HWP to PDF conversion: {hwp_path}")
        
        # LibreOffice가 설치되어 있는지 확인
        if not is_package_installed('libreoffice'):
            print("LibreOffice is not installed. Please install it first.")
            return ''
        
        # 원본 파일과 같은 위치에 PDF 생성
        output_dir = os.path.dirname(hwp_path)
        original_cwd = os.getcwd()
        
        try:
            # 작업 디렉토리를 출력 디렉토리로 변경
            print(f"Changing directory to: {output_dir}")
            os.chdir(output_dir)
            
            # LibreOffice로 변환 시도
            print("Running LibreOffice conversion command...")
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', os.path.basename(hwp_path)],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            print(f"Conversion command output: {result.stdout}")
            if result.stderr:
                print(f"Conversion command error: {result.stderr}")
            
            if result.returncode == 0:
                # 변환된 PDF 파일 경로
                pdf_path = os.path.splitext(hwp_path)[0] + '.pdf'
                if os.path.exists(pdf_path):
                    print(f"PDF file created successfully: {pdf_path}")
                    return pdf_path
                else:
                    print(f"PDF file not found at expected path: {pdf_path}")
            else:
                print(f"Conversion command failed with return code: {result.returncode}")
                    
        except subprocess.TimeoutExpired:
            print("Conversion timed out after 30 seconds")
        except Exception as e:
            print(f"Conversion error: {str(e)}")
        finally:
            # 원래 작업 디렉토리로 복원
            print(f"Restoring original directory: {original_cwd}")
            os.chdir(original_cwd)
        
        return ''
        
    except Exception as e:
        print(f"Error in convert_hwp_to_pdf: {str(e)}")
        return ''

def get_file_info(path: str) -> Dict[str, Any]:
    """파일의 기본 정보를 조회하는 함수
    
    Args:
        path: 파일 경로
        
    Returns:
        파일 정보를 담은 딕셔너리:
        - filename: 파일명
        - modified_time: 수정일시
        - size: 파일 크기
        - created_time: 생성일시
        - owner: 소유자
        - permissions: 파일 권한
    """
    try:
        # 1. 경로 정규화
        norm_path = normalize_path(path)
        
        # 2. 파일 정보 조회
        stat = os.stat(norm_path)
        
        # 3. 결과 반환
        return {
            'filename': os.path.basename(norm_path),
            'modified_time': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
            'size': stat.st_size,
            'created_time': datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S'),
            'owner': stat.st_uid,
            'permissions': oct(stat.st_mode)[-3:]
        }
    except Exception as e:
        print(f"[ERROR] 파일 정보 조회 실패: {str(e)}")
        raise

def summarize_with_openai(content: str, max_length: int = None) -> str:
    """OpenAI를 사용하여 텍스트를 요약하는 테스트 함수입니다.
    
    Args:
        content (str): 요약할 텍스트 내용
        max_length (int, optional): 최대 요약 길이. Defaults to None.
        
    Returns:
        str: 요약된 텍스트
    """
    try:
        from openai import OpenAI
        import os
        
        # OpenAI 클라이언트 초기화
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            api_key = config.get('OPENAI_API_KEY')
            if not api_key:
                raise ValueError("OpenAI API key is not configured.")
            os.environ['OPENAI_API_KEY'] = api_key
            
        client = OpenAI()  # 환경 변수에서 API 키를 자동으로 가져옴
        
        # 모델 설정
        model = "gpt-4o-mini"
        max_tokens = min(max_length or 4000, 4000)
        
        # 요약 요청
        messages = [
            {"role": "system", "content": "Please provide a clear and concise summary of the following text:"},
            {"role": "system", "content": "Create the summary in Korean."},
            {"role": "user", "content": content}
        ]
        
        response = client.chat.completions.create(
            model=model,
            messages=messages,
            max_tokens=max_tokens
        )
        
        if not response.choices:
            raise ValueError("Empty response from OpenAI API")
            
        return response.choices[0].message.content
        
    except Exception as e:
        print(f"[ERROR] Summarization failed: {str(e)}")
        raise

